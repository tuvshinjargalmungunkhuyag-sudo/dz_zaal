import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

prs = Presentation('Goviyn_Sport_Iltgel_v8.pptx')
# Show slides 8-21 content
for i in range(7, 21):
    slide = prs.slides[i]
    texts = [sh.text_frame.text.strip()[:70] for sh in slide.shapes 
             if sh.has_text_frame and sh.text_frame.text.strip()]
    print(f"\nSLIDE {i+1} ({len(slide.shapes)} shapes): '{texts[0] if texts else ''}' ")
    for t in texts[1:6]:
        print(f"  {t}")
