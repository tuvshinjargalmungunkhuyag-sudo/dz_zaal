import sys, io, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

prs = Presentation('Goviyn_Sport_Iltgel_v8.pptx')
emoji_re = re.compile("[\U00010000-\U0010ffff]|[\U00002600-\U000027BF]|[\U0001F300-\U0001F9FF]", flags=re.UNICODE)

for i, slide in enumerate(prs.slides):
    emojis_found = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text
            found = emoji_re.findall(t)
            if found:
                emojis_found.extend(found)
    if emojis_found:
        unique = list(dict.fromkeys(emojis_found))
        print(f"Slide {i+1}: {''.join(unique)}")
