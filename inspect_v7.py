import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches

prs = Presentation('Goviyn_Sport_Iltgel_v7.pptx')
for i, slide in enumerate(prs.slides):
    print(f'\n=== SLIDE {i+1} ===')
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                print(f'  [{shape.name}] "{t[:120]}"')
        # shape size
        print(f'  shape: {shape.name} l={shape.left/914400:.2f} t={shape.top/914400:.2f} w={shape.width/914400:.2f} h={shape.height/914400:.2f}')
