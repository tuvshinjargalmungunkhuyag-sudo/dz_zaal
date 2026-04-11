import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

prs = Presentation('Goviyn_Sport_Iltgel_v8.pptx')
for i in [2, 3, 4, 5, 6]:  # slides 3-7
    slide = prs.slides[i]
    texts = [sh.text_frame.text.strip()[:60] for sh in slide.shapes 
             if sh.has_text_frame and sh.text_frame.text.strip()]
    print(f"\nSLIDE {i+1} ({len(slide.shapes)} shapes):")
    for t in texts[:10]:
        print(f"  · {t}")
    if len(texts) > 10:
        print(f"  ... {len(texts)} text items total")
