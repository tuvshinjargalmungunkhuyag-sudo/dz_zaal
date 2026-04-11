import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

LOGO_PATH = r"C:\Users\user\Downloads\MUST-logo.png"

NAVY       = RGBColor(0x0A, 0x1A, 0x5C)
BLUE       = RGBColor(0x1A, 0x5C, 0xBB)
BLUE_LIGHT = RGBColor(0x4A, 0x86, 0xD8)
ORANGE     = RGBColor(0xE8, 0x6E, 0x1A)
YELLOW     = RGBColor(0xF5, 0xC5, 0x18)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_D     = RGBColor(0x44, 0x44, 0x44)
GRAY       = RGBColor(0x66, 0x66, 0x66)
GRAY_L     = RGBColor(0xF2, 0xF4, 0xF7)
GREEN_D    = RGBColor(0x14, 0x57, 0x2A)
GREEN      = RGBColor(0x1E, 0x8B, 0x45)
GREEN_L    = RGBColor(0xE8, 0xF5, 0xED)
AMBER      = RGBColor(0xE6, 0x7E, 0x22)
AMBER_L    = RGBColor(0xFF, 0xF3, 0xE0)
PURPLE     = RGBColor(0x4A, 0x14, 0x8C)
RED        = RGBColor(0xC0, 0x39, 0x2B)
DIVIDER    = RGBColor(0xCC, 0xD3, 0xDE)
ARROW_COL  = RGBColor(0x2C, 0x3E, 0x50)

APP_BG   = RGBColor(0xED, 0xE7, 0xDA)
APP_CARD = RGBColor(0xF8, 0xF4, 0xEE)
APP_ACC  = RGBColor(0xD4, 0x70, 0x0F)
APP_TXT  = RGBColor(0x1A, 0x12, 0x08)
APP_TXT2 = RGBColor(0x6B, 0x54, 0x38)
APP_SURF = RGBColor(0xE4, 0xDD, 0xD0)
APP_DIV  = RGBColor(0xD0, 0xC4, 0xB0)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

YEAR       = "2026"
AUTHOR     = "М. Түвшинжаргал (B221940619)"
SUPERVISOR = "Доктор (Ph.D) Д.Эрдэнэтуяа"


def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(0)):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh


def tb(slide, text, x, y, w, h,
       size=13, bold=False, color=BLACK,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def header(slide, title):
    rect(slide, 0, 0, W, Inches(0.76), fill=WHITE)
    rect(slide, 0, 0, Inches(0.12), Inches(0.76), fill=NAVY)
    slide.shapes.add_picture(LOGO_PATH,
                             Inches(0.18), Inches(0.07),
                             Inches(0.60), Inches(0.60))
    tb(slide, title,
       Inches(0.88), Inches(0.13), Inches(12.2), Inches(0.50),
       size=20, bold=True, color=NAVY)
    rect(slide, 0, H - Inches(0.38), W, Inches(0.38), fill=NAVY)
    tb(slide, f"Улаанбаатар хот,  {YEAR} он",
       Inches(0.2), H - Inches(0.35), Inches(5), Inches(0.30),
       size=11, color=YELLOW)


def content_bg(slide):
    rect(slide, 0, Inches(0.76), W,
         H - Inches(0.76) - Inches(0.38), fill=WHITE)


def new_slide():
    return prs.slides.add_slide(blank)


def v_arrow(slide, cx, y0, y1, color=ARROW_COL, lw=Inches(0.045)):
    shaft = y1 - y0 - Inches(0.18)
    if shaft > 0:
        rect(slide, cx - lw // 2, y0, lw, shaft, fill=color)
    aw, ah, steps = Inches(0.22), Inches(0.18), 9
    for i in range(steps):
        frac = i / steps
        wi = int(aw * (1 - frac))
        rect(slide, cx - wi // 2,
             y1 - ah + int((ah / steps) * i),
             wi, int(ah / steps) + Inches(0.002),
             fill=color)


def h_arrow(slide, x0, y, x1, color=ARROW_COL, lw=Inches(0.045)):
    shaft = x1 - x0 - Inches(0.18)
    if shaft > 0:
        rect(slide, x0, y - lw // 2, shaft, lw, fill=color)
    aw, ah, steps = Inches(0.18), Inches(0.20), 8
    for i in range(steps):
        frac = i / steps
        hi = int(ah * (1 - frac))
        rect(slide, x1 - aw + int((aw / steps) * i),
             y - hi // 2, int(aw / steps) + Inches(0.002), hi,
             fill=color)


def bidirectional_v(slide, cx, y0, y1, color=ARROW_COL):
    aw, ah, steps = Inches(0.22), Inches(0.18), 9
    for i in range(steps):
        frac = i / steps
        wi = int(aw * (1 - frac))
        rect(slide, cx - wi // 2,
             y0 + int((ah / steps) * i),
             wi, int(ah / steps) + Inches(0.002),
             fill=color)
    lw = Inches(0.045)
    shaft = y1 - y0 - ah * 2 - Inches(0.02)
    if shaft > 0:
        rect(slide, cx - lw // 2, y0 + ah, lw, shaft, fill=color)
    for i in range(steps):
        frac = i / steps
        wi = int(aw * (1 - frac))
        rect(slide, cx - wi // 2,
             y1 - ah + int((ah / steps) * i),
             wi, int(ah / steps) + Inches(0.002),
             fill=color)


def ph(slide, sx, sy, sw, sh):
    """Draw phone mockup. Returns (cx, cy, cw, ch) - usable content area."""
    rect(slide, sx, sy, sw, sh, fill=RGBColor(0x1A, 0x1A, 0x2A))
    brd   = Inches(0.09)
    notch = Inches(0.28)
    scx = sx + brd
    scy = sy + notch
    scw = sw - 2 * brd
    sch = sh - notch - brd - Inches(0.14)
    rect(slide, scx, scy, scw, sch, fill=APP_BG)
    rect(slide, scx, scy, scw, Inches(0.18), fill=APP_ACC)
    rect(slide, sx + sw // 2 - Inches(0.16), sy + Inches(0.07),
         Inches(0.32), Inches(0.12), fill=RGBColor(0x33, 0x33, 0x44))
    rect(slide, sx + sw // 2 - Inches(0.28),
         sy + sh - brd - Inches(0.08),
         Inches(0.56), Inches(0.06), fill=RGBColor(0x88, 0x88, 0x99))
    cx = scx
    cy = scy + Inches(0.18)
    cw = scw
    ch = sch - Inches(0.18)
    return cx, cy, cw, ch


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1  —  Гарчиг
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, W, H, fill=YELLOW)
rect(s, 0, 0, Inches(0.22), H, fill=NAVY)
rect(s, W - Inches(0.22), 0, Inches(0.22), H, fill=NAVY)
rect(s, 0, H - Inches(0.44), W, Inches(0.44), fill=NAVY)
s.shapes.add_picture(LOGO_PATH, Inches(0.32), Inches(0.14), Inches(0.74), Inches(0.74))
tb(s, "МОНГОЛ УЛСЫН ШИНЖЛЭХ УХААН",
   Inches(1.14), Inches(0.17), Inches(7), Inches(0.27),
   size=10, bold=True, color=NAVY)
tb(s, "ТЕХНОЛОГИЙН ИХ СУРГУУЛЬ",
   Inches(1.14), Inches(0.43), Inches(7), Inches(0.27),
   size=10, bold=True, color=NAVY)
tb(s, "MONGOLIAN UNIVERSITY OF SCIENCE AND TECHNOLOGY",
   Inches(1.14), Inches(0.67), Inches(8), Inches(0.22),
   size=8, color=GRAY)
rect(s, Inches(1.5), Inches(1.52), Inches(10.3), Inches(0.04), fill=NAVY)
tb(s, "ҮҮЛЭН ОРЧИНД СУУРИЛСАН СПОРТ ЗААЛНЫ",
   Inches(1.5), Inches(1.65), Inches(10.3), Inches(0.70),
   size=27, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tb(s, "ЦАГ ЗАХИАЛГЫН 3 ДАВХАРГАТ СИСТЕМИЙН",
   Inches(1.5), Inches(2.33), Inches(10.3), Inches(0.68),
   size=27, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tb(s, "ЗАГВАРЧЛАЛ БА ХЭРЭГЖИЛТ",
   Inches(1.5), Inches(3.00), Inches(10.3), Inches(0.64),
   size=27, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
rect(s, Inches(1.5), Inches(3.76), Inches(10.3), Inches(0.04), fill=NAVY)
tb(s, f"Илтгэгч: {AUTHOR}",
   Inches(4.0), Inches(3.90), Inches(5.6), Inches(0.32),
   size=14, color=NAVY, align=PP_ALIGN.CENTER)
tb(s, f"Удирдагч: {SUPERVISOR}",
   Inches(4.0), Inches(4.24), Inches(5.6), Inches(0.32),
   size=14, color=NAVY, align=PP_ALIGN.CENTER)
tb(s, f"Улаанбаатар хот,  {YEAR} он",
   Inches(0.32), H - Inches(0.42), Inches(5), Inches(0.36),
   size=12, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2  —  Агуулга
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Агуулга")
content_bg(s)
items = [
    ("01", "Судалгааны асуудал ба актуаль байдал"),
    ("02", "Судалгааны зорилго ба зорилт"),
    ("03", "3 Давхаргат системийн архитектур"),
    ("04", "Технологийн стек ба хэрэгслийн сонголт"),
    ("05", "Мобайл аппликейшны UI/UX загвар"),
    ("06", "Захиалгын үйл явц ба давхцал шалгах логик"),
    ("07", "Backend архитектур ба REST API"),
    ("08", "AI туслах: Groq / LLaMA 3.1 нэгтгэл"),
    ("09", "Хэрэгжилтийн явц ба туршилтын үр дүн"),
    ("10", "Дүгнэлт ба цаашдын судалгааны чиглэл"),
]
for i, (num, txt) in enumerate(items):
    col = i % 2
    row = i // 2
    x = Inches(0.65) + col * Inches(6.55)
    y = Inches(1.02) + row * Inches(1.04)
    rect(s, x, y, Inches(6.15), Inches(0.86), fill=GRAY_L, line=DIVIDER, lw=Pt(1))
    rect(s, x, y, Inches(0.06), Inches(0.86), fill=BLUE)
    tb(s, num, x + Inches(0.16), y + Inches(0.14),
       Inches(0.55), Inches(0.55), size=17, bold=True, color=BLUE)
    tb(s, txt, x + Inches(0.74), y + Inches(0.22),
       Inches(5.2), Inches(0.46), size=14, bold=True, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3  —  Судалгааны асуудал
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Судалгааны асуудал ба актуаль байдал")
content_bg(s)
tb(s, "Одоогийн нөхцөл байдал",
   Inches(0.5), Inches(0.90), Inches(8.2), Inches(0.40),
   size=17, bold=True, color=NAVY)
issues = [
    ("01", "Дижитал дэд бүтцийн дутагдал",
     "Спортын заалны захиалгын үйл явц нь утасны яриа болон биечилсэн"
     " бүртгэлд тулгуурладаг тул хэрэглэгчийн цаг хугацаа үрэгддэг."),
    ("02", "Бодит цагийн мэдээллийн хүртээмж",
     "Захиалгын боломж, цагийн слотын нөхцөл байдлыг урьдчилан мэдэх"
     " цахим технологийн боломж байхгүй байна."),
    ("03", "Захиалгын давхардал ба алдааны эрсдэл",
     "Гараар хөтлөгдөх бүртгэлийн системд давхардсан захиалга болон"
     " мэдээллийн зөрөлдөөн тогтмол ажиглагддаг."),
    ("04", "Нэгдсэн мэдээллийн сан ба аналитик дутагдал",
     "Захиалгын статистик, орлогын тайлан, хэрэглэгчийн зан байдлын"
     " дата цуглуулах болон дүн шинжилгээний автоматчлал хийгдээгүй."),
]
for i, (num, title, body) in enumerate(issues):
    col = i % 2
    row = i // 2
    x = Inches(0.45) + col * Inches(6.45)
    y = Inches(1.44) + row * Inches(2.40)
    rect(s, x, y, Inches(6.15), Inches(2.22), fill=WHITE, line=DIVIDER, lw=Pt(1))
    rect(s, x, y, Inches(0.50), Inches(2.22), fill=BLUE)
    tb(s, num, x + Inches(0.60), y + Inches(0.12),
       Inches(0.60), Inches(0.40), size=16, bold=True, color=BLUE)
    tb(s, title, x + Inches(1.22), y + Inches(0.15),
       Inches(4.70), Inches(0.42), size=13, bold=True, color=NAVY)
    tb(s, body, x + Inches(0.62), y + Inches(0.64),
       Inches(5.38), Inches(1.44), size=11, color=GRAY, wrap=True)
rect(s, Inches(9.25), Inches(0.90), Inches(3.75), Inches(5.72), fill=NAVY)
tb(s, "Судалгааны актуаль байдал",
   Inches(9.35), Inches(1.04), Inches(3.55), Inches(0.38),
   size=11, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
rect(s, Inches(9.55), Inches(1.48), Inches(3.14), Inches(0.04), fill=ORANGE)
tb(s, "5+", Inches(9.25), Inches(1.60), Inches(3.75), Inches(0.82),
   size=44, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
tb(s, "спортын байгууллага",
   Inches(9.25), Inches(2.38), Inches(3.75), Inches(0.36),
   size=12, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, Inches(9.55), Inches(2.82), Inches(3.14), Inches(0.04), fill=ORANGE)
tb(s, "Нэгдсэн цахим\nзахиалгын систем\nогт байхгүй",
   Inches(9.25), Inches(2.94), Inches(3.75), Inches(0.88),
   size=13, color=YELLOW, align=PP_ALIGN.CENTER)
rect(s, Inches(9.55), Inches(3.90), Inches(3.14), Inches(0.04), fill=ORANGE)
tb(s, "Шийдлийн зам:\nCloud-д суурилсан\n3 давхаргат систем",
   Inches(9.25), Inches(4.02), Inches(3.75), Inches(1.00),
   size=12, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4  —  Зорилго ба зорилт
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Судалгааны зорилго ба зорилт")
content_bg(s)
rect(s, Inches(0.42), Inches(0.88), Inches(12.48), Inches(1.04), fill=NAVY)
tb(s, "Судалгааны зорилго",
   Inches(0.58), Inches(0.92), Inches(3.8), Inches(0.36),
   size=13, bold=True, color=YELLOW)
tb(s,
   "Өмнөговь аймгийн Даланзадгад хотын спортын байгууллагад зориулан үүлэн орчинд"
   " суурилсан, 3 давхаргат архитектуртай, бодит цагийн цаг захиалгын ухаалаг"
   " системийг загварчлан хэрэгжүүлэх.",
   Inches(0.58), Inches(1.28), Inches(12.10), Inches(0.54),
   size=12, color=WHITE)
goals = [
    ("З-1", "Цаг захиалгын модуль",
     "Бодит цагийн хүртээмж шалгалт, давхардлаас сэргийлэх логик,"
     " хагас / бүтэн талбайн уян хатан захиалгын тогтолцоог хэрэгжүүлэх"),
    ("З-2", "Мобайл интерфэйс",
     "Flutter ашиглан iOS/Android-д зэрэг ажилладаг, Material Design 3-д"
     " нийцсэн хэрэглэгчийн интерфэйс боловсруулах"),
    ("З-3", "AI туслах систем",
     "Groq API ба LLaMA 3.1 8B ашиглан монгол хэлээр харилцах,"
     " заалны мэдээлэл олгох ухаалаг туслах нэвтрүүлэх"),
    ("З-4", "Cloud дэд бүтэц",
     "AWS EC2, Docker container, GitHub Actions CI/CD ашиглан"
     " автомат деплой, масштаблах чадвартай орчин бүрдүүлэх"),
    ("З-5", "Мэдэгдэл ба автоматчлал",
     "Firebase Cloud Messaging болон node-cron ашиглан захиалгын"
     " статусын автомат шинэчлэл, push мэдэгдлийн систем хэрэгжүүлэх"),
    ("З-6", "Аюулгүй байдал",
     "Firebase Auth, OTP баталгаажуулалт, HTTPS шифрлэлт,"
     " орчны хувьсагч (.env)-аар нууц мэдээлэл хамгаалах"),
]
for i, (code, title, body) in enumerate(goals):
    col = i % 3
    row = i // 3
    x = Inches(0.42) + col * Inches(4.20)
    y = Inches(2.10) + row * Inches(2.36)
    rect(s, x, y, Inches(3.98), Inches(2.18), fill=GRAY_L, line=DIVIDER, lw=Pt(1))
    rect(s, x, y, Inches(3.98), Inches(0.48), fill=BLUE)
    tb(s, code + "  —  " + title,
       x + Inches(0.14), y + Inches(0.07),
       Inches(3.70), Inches(0.36), size=12, bold=True, color=WHITE)
    tb(s, body, x + Inches(0.14), y + Inches(0.60),
       Inches(3.70), Inches(1.48), size=10.5, color=GRAY, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5  —  Ач холбогдол
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Судалгааны ач холбогдол")
content_bg(s)
panels = [
    (BLUE,   "Шинжлэх ухааны ач холбогдол",
     "Үүлэн тооцооллын орчинд суурилсан 3 давхаргат архитектурийн загварыг"
     " спортын байгууллагын менежментэд хэрэглэж болохыг нотлон харуулна."
     " REST API ба serverless мэдээллийн санг хосолсон гибрид хандлагыг"
     " судалгааны тайланд тусгана."),
    (GREEN,  "Технологийн ач холбогдол",
     "Flutter cross-platform, Node.js, Firebase, Groq LLaMA зэрэг орчин"
     " үеийн технологиудыг нэгтгэсэн цогц систем бүрдүүлснээр монголын"
     " программ хангамжийн салбарт практик жишиг бий болно."),
    (AMBER,  "Нийгмийн ач холбогдол",
     "Иргэдийн спортын идэвхийг нэмэгдүүлэхэд дэмжлэг болох дижитал"
     " дэд бүтцийг Даланзадгад хотод бүрдүүлж, цахим засаглалын"
     " хөгжилд хувь нэмэр оруулна."),
    (PURPLE, "Эдийн засгийн ач холбогдол",
     "Захиалгын автоматчлал, цаг хуваарийн оновчтой удирдлага, орлогын"
     " дата аналитикаар дамжуулан спортын байгууллагын ашиглалтын"
     " үзүүлэлтийг нэмэгдүүлж, гарах зардлыг бууруулна."),
]
cw = Inches(5.90)
for i, (color, title, body) in enumerate(panels):
    col = i % 2
    row = i // 2
    x = Inches(0.50) + col * Inches(6.55)
    y = Inches(1.02) + row * Inches(2.85)
    rect(s, x, y, cw, Inches(2.62), fill=WHITE, line=DIVIDER, lw=Pt(1.2))
    rect(s, x, y, cw, Inches(0.50), fill=color)
    rect(s, x, y, Inches(0.08), Inches(2.62), fill=color)
    tb(s, title, x + Inches(0.20), y + Inches(0.08),
       cw - Inches(0.28), Inches(0.36), size=14, bold=True, color=WHITE)
    tb(s, body, x + Inches(0.20), y + Inches(0.64),
       cw - Inches(0.30), Inches(1.86), size=11.5, color=GRAY_D, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6  —  3 Давхаргат архитектур
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "3 Давхаргат системийн архитектур")
content_bg(s)

LABEL_X = Inches(0.30);  LABEL_W = Inches(1.00)
MAIN_X  = Inches(1.38);  MAIN_W  = Inches(9.62)
INFRA_X = Inches(11.18); INFRA_W = Inches(1.90)

T1_Y = Inches(0.86);  T1_H = Inches(1.78)
T2_Y = Inches(3.00);  T2_H = Inches(1.78)
T3_Y = Inches(5.14);  T3_H = Inches(1.60)

T1_BG = RGBColor(0xEB, 0xF3, 0xFD)
T2_BG = GREEN_L
T3_BG = AMBER_L

for ty, th, hdr, bg in (
    (T1_Y, T1_H, BLUE,    T1_BG),
    (T2_Y, T2_H, GREEN_D, T2_BG),
    (T3_Y, T3_H, AMBER,   T3_BG),
):
    rect(s, LABEL_X, ty, LABEL_W + MAIN_W, th, fill=bg, line=DIVIDER, lw=Pt(1.2))
    rect(s, LABEL_X, ty, LABEL_W, th, fill=hdr)

for ty, th, num, l1, l2 in (
    (T1_Y, T1_H, "1", "Presentation", "Tier"),
    (T2_Y, T2_H, "2", "Application",  "Tier"),
    (T3_Y, T3_H, "3", "Data",         "Tier"),
):
    tb(s, num,  LABEL_X, ty + Inches(0.12), LABEL_W, Inches(0.46),
       size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, l1,   LABEL_X, ty + Inches(0.58), LABEL_W, Inches(0.36),
       size=9,  bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, l2,   LABEL_X, ty + Inches(0.90), LABEL_W, Inches(0.30),
       size=9,  bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# T1 Flutter box
FLT_X = MAIN_X + Inches(0.14); FLT_W = Inches(2.48)
FLT_Y = T1_Y + Inches(0.18);   FLT_H = T1_H - Inches(0.36)
rect(s, FLT_X, FLT_Y, FLT_W, FLT_H, fill=BLUE, line=NAVY, lw=Pt(1))
tb(s, "Flutter", FLT_X, FLT_Y + Inches(0.10), FLT_W, Inches(0.40),
   size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Mobile Application\niOS  /  Android  (Dart 3.x)",
   FLT_X, FLT_Y + Inches(0.52), FLT_W, Inches(0.72),
   size=9.5, color=RGBColor(0xCF, 0xE2, 0xFF), align=PP_ALIGN.CENTER)

screens_t1 = ["Login", "Home", "Booking\nSheet", "Confirm", "Bookings", "Profile", "Chat", "Notif"]
scr_start_x = MAIN_X + Inches(2.76)
SCR_W = Inches(1.62); SCR_H = Inches(0.64)
for si, scr in enumerate(screens_t1):
    scol = si % 4; srow = si // 4
    sx = scr_start_x + scol * Inches(1.72)
    sy = T1_Y + Inches(0.18) + srow * Inches(0.76)
    rect(s, sx, sy, SCR_W, SCR_H, fill=BLUE_LIGHT, line=NAVY, lw=Pt(0.5))
    tb(s, scr, sx, sy + Inches(0.12), SCR_W, SCR_H - Inches(0.12),
       size=8.5, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

# T2 boxes
app_boxes = [
    ("Express.js\nREST API", "6 endpoint\nJSON / HTTP"),
    ("node-cron\nScheduler",  "Статус\nшинэчлэл"),
    ("Groq API\nLLaMA 3.1",   "AI Chat\n8B model"),
]
APP_W = Inches(2.68); APP_H = T2_H - Inches(0.36); APP_Y = T2_Y + Inches(0.18)
app_clrs = [GREEN_D, RGBColor(0x1A, 0x6B, 0x3A), PURPLE]
for ai, (nm, dt) in enumerate(app_boxes):
    ax = MAIN_X + Inches(0.14) + ai * Inches(3.16)
    rect(s, ax, APP_Y, APP_W, APP_H, fill=app_clrs[ai], line=GREEN_D, lw=Pt(1))
    tb(s, nm, ax, APP_Y + Inches(0.10), APP_W, Inches(0.60),
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, dt, ax, APP_Y + Inches(0.72), APP_W, Inches(0.60),
       size=9, color=RGBColor(0xBB, 0xFF, 0xCC), align=PP_ALIGN.CENTER)

# T3 boxes
data_boxes = [
    ("Firestore",      "NoSQL  •  Compound\nindex  •  Real-time"),
    ("Firebase Auth",  "Phone OTP\nUID  •  ID Token"),
    ("FCM",            "Push мэдэгдэл\nDevice token"),
]
DAT_W = Inches(2.68); DAT_H = T3_H - Inches(0.36); DAT_Y = T3_Y + Inches(0.18)
dat_clrs = [AMBER, RGBColor(0xBF, 0x36, 0x00), RGBColor(0xD4, 0x70, 0x0F)]
for di, (nm, dt) in enumerate(data_boxes):
    dx = MAIN_X + Inches(0.14) + di * Inches(3.16)
    rect(s, dx, DAT_Y, DAT_W, DAT_H, fill=dat_clrs[di], line=AMBER, lw=Pt(1))
    tb(s, nm, dx, DAT_Y + Inches(0.10), DAT_W, Inches(0.40),
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, dt, dx, DAT_Y + Inches(0.52), DAT_W, Inches(0.68),
       size=9, color=RGBColor(0xFF, 0xEE, 0xBB), align=PP_ALIGN.CENTER)

ACX = MAIN_X + MAIN_W // 2
bidirectional_v(s, ACX, T1_Y + T1_H, T2_Y, color=BLUE)
bidirectional_v(s, ACX, T2_Y + T2_H, T3_Y, color=GREEN_D)
tb(s, "HTTP / REST", ACX - Inches(0.56), T1_Y + T1_H + Inches(0.04),
   Inches(1.12), Inches(0.22), size=8, bold=True, color=BLUE, italic=True, align=PP_ALIGN.CENTER)
tb(s, "Firebase Admin SDK", ACX - Inches(0.88), T2_Y + T2_H + Inches(0.04),
   Inches(1.76), Inches(0.22), size=8, bold=True, color=GREEN_D, italic=True, align=PP_ALIGN.CENTER)

rect(s, INFRA_X, T1_Y, INFRA_W, T3_Y + T3_H - T1_Y,
     fill=RGBColor(0x0A, 0x1A, 0x3A), line=BLUE_LIGHT, lw=Pt(1))
tb(s, "Cloud\nInfra", INFRA_X, T1_Y + Inches(0.10), INFRA_W, Inches(0.52),
   size=9, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
infra = [
    ("AWS EC2",        "ap-southeast-1"),
    ("Docker",         "node:20-alpine"),
    ("AWS ECR",        "Image registry"),
    ("GitHub\nActions","CI/CD pipeline"),
]
for ii, (t, d) in enumerate(infra):
    iy = Inches(1.56) + ii * Inches(1.14)
    rect(s, INFRA_X + Inches(0.10), iy, INFRA_W - Inches(0.20), Inches(1.00),
         fill=RGBColor(0x18, 0x2C, 0x6A), line=BLUE_LIGHT, lw=Pt(0.75))
    rect(s, INFRA_X + Inches(0.10), iy, Inches(0.06), Inches(1.00), fill=ORANGE)
    tb(s, t, INFRA_X + Inches(0.22), iy + Inches(0.08),
       INFRA_W - Inches(0.34), Inches(0.34), size=10, bold=True, color=WHITE)
    tb(s, d, INFRA_X + Inches(0.22), iy + Inches(0.44),
       INFRA_W - Inches(0.34), Inches(0.50), size=8.5, color=RGBColor(0xAA, 0xBB, 0xEE))

h_arrow(s, MAIN_X + MAIN_W - Inches(0.10), T2_Y + T2_H // 2,
        INFRA_X - Inches(0.04), color=ORANGE)
tb(s, "deploy", MAIN_X + MAIN_W - Inches(0.04), T2_Y + T2_H // 2 - Inches(0.26),
   Inches(0.80), Inches(0.22), size=8, bold=True, color=ORANGE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7  —  Технологийн стек
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Технологийн стек ба хэрэгслийн сонголт")
content_bg(s)
tech_cols = [
    ("Presentation Tier", [
        ("Flutter 3.x  (Dart)",   "Cross-platform UI framework"),
        ("Firebase Auth SDK",     "Утасны OTP нэвтрэлт"),
        ("HTTP / REST Client",    "API харилцаа"),
        ("Push Notification",     "FCM device token"),
        ("Material Design 3",     "UI / UX дизайн систем"),
    ], BLUE),
    ("Application Tier", [
        ("Node.js 20 LTS",        "JavaScript server runtime"),
        ("Express.js 4.x",        "RESTful API framework"),
        ("Firebase Admin SDK",    "Firestore & Auth удирдлага"),
        ("Groq API  (LLaMA 3.1)", "Хурдан загварын дүүргэлт"),
        ("node-cron",             "Тогтмол ажлын зохицуулагч"),
    ], GREEN_D),
    ("Data & Infra Tier", [
        ("Firebase Firestore",    "NoSQL баримт бичгийн сан"),
        ("AWS EC2",               "Үүлэн сервер хостинг"),
        ("AWS ECR",               "Docker дүрсний бүртгэл"),
        ("Docker  (Alpine)",      "Контейнерчлал"),
        ("GitHub Actions",        "CI/CD автомат деплой"),
    ], AMBER),
]
CW = Inches(3.98)
for ci, (col_title, citems, color) in enumerate(tech_cols):
    x = Inches(0.42) + ci * Inches(4.30)
    rect(s, x, Inches(0.88), CW, Inches(0.50), fill=color)
    tb(s, col_title, x + Inches(0.14), Inches(0.91),
       CW - Inches(0.28), Inches(0.44), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ri, (name, desc) in enumerate(citems):
        by = Inches(1.52) + ri * Inches(1.04)
        rect(s, x, by, CW, Inches(0.92), fill=WHITE, line=DIVIDER, lw=Pt(0.75))
        rect(s, x, by, Inches(0.06), Inches(0.92), fill=color)
        tb(s, name, x + Inches(0.18), by + Inches(0.10),
           CW - Inches(0.26), Inches(0.36), size=13, bold=True, color=color)
        tb(s, desc, x + Inches(0.18), by + Inches(0.48),
           CW - Inches(0.26), Inches(0.38), size=11, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8  —  UI/UX Загвар (Phone Mockups)
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Мобайл аппликейшны UI/UX загвар")
content_bg(s)

PH_W = Inches(2.22)
PH_H = Inches(4.10)
PH_Y = Inches(0.90)
PH_GAP = Inches(0.44)
total_w = 5 * PH_W + 4 * PH_GAP
PH_START = (W - total_w) // 2

# ── Phone 1: Login ────────────────────────────────────────────────────────────
p1x = PH_START + 0 * (PH_W + PH_GAP)
cx, cy, cw, ch = ph(s, p1x, PH_Y, PH_W, PH_H)
rect(s, cx, cy, cw, Inches(0.50), fill=APP_BG)
tb(s, "Говийн Спорт", cx, cy + Inches(0.06), cw, Inches(0.26),
   size=8.5, bold=True, color=APP_TXT, align=PP_ALIGN.CENTER)
tb(s, "Нэвтрэх", cx, cy + Inches(0.30), cw, Inches(0.18),
   size=6.5, color=APP_TXT2, align=PP_ALIGN.CENTER)
rect(s, cx + cw // 2 - Inches(0.28), cy + Inches(0.62),
     Inches(0.56), Inches(0.56), fill=APP_ACC)
tb(s, "G", cx + cw // 2 - Inches(0.28), cy + Inches(0.70),
   Inches(0.56), Inches(0.36), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# progress bars
rect(s, cx + Inches(0.10), cy + Inches(1.32), Inches(0.78), Inches(0.06), fill=APP_ACC)
rect(s, cx + Inches(0.96), cy + Inches(1.32), Inches(0.78), Inches(0.06), fill=APP_DIV)
tb(s, "Утасны дугаар", cx + Inches(0.10), cy + Inches(1.46), cw - Inches(0.20), Inches(0.20),
   size=6.5, bold=True, color=APP_TXT)
rect(s, cx + Inches(0.10), cy + Inches(1.68), cw - Inches(0.20), Inches(0.30),
     fill=WHITE, line=APP_ACC, lw=Pt(1.2))
tb(s, "+976  8800 1234", cx + Inches(0.16), cy + Inches(1.74),
   cw - Inches(0.32), Inches(0.20), size=7, color=APP_TXT2)
rect(s, cx + Inches(0.10), cy + Inches(2.12), cw - Inches(0.20), Inches(0.32), fill=APP_ACC)
tb(s, "Код илгээх", cx + Inches(0.10), cy + Inches(2.18),
   cw - Inches(0.20), Inches(0.22), size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Говийн Спорт  •  Даланзадгад",
   cx, cy + ch - Inches(0.28), cw, Inches(0.24),
   size=5.5, color=APP_TXT2, align=PP_ALIGN.CENTER, italic=True)
tb(s, "Нэвтрэлт", p1x, PH_Y + PH_H + Inches(0.08), PH_W, Inches(0.28),
   size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── Phone 2: Home Screen ──────────────────────────────────────────────────────
p2x = PH_START + 1 * (PH_W + PH_GAP)
cx, cy, cw, ch = ph(s, p2x, PH_Y, PH_W, PH_H)
rect(s, cx, cy, cw, Inches(0.40), fill=APP_SURF)
tb(s, "Говийн Спорт", cx + Inches(0.10), cy + Inches(0.08), cw - Inches(0.44), Inches(0.26),
   size=8.5, bold=True, color=APP_TXT)
rect(s, cx + cw - Inches(0.34), cy + Inches(0.08), Inches(0.24), Inches(0.24), fill=APP_ACC)
rect(s, cx + Inches(0.08), cy + Inches(0.46), cw - Inches(0.16), Inches(0.28),
     fill=WHITE, line=APP_DIV, lw=Pt(0.8))
tb(s, "Хайх...", cx + Inches(0.14), cy + Inches(0.52), cw - Inches(0.28), Inches(0.18),
   size=6, color=APP_TXT2, italic=True)
venues = [
    ("Говийн Арена",     "15,000₮/цаг", "4.8"),
    ("Өмнөговь Спорт",  "12,000₮/цаг", "4.6"),
]
for vi, (vname, vprice, vrate) in enumerate(venues):
    vy = cy + Inches(0.84) + vi * Inches(0.86)
    rect(s, cx + Inches(0.08), vy, cw - Inches(0.16), Inches(0.78),
         fill=APP_CARD, line=APP_DIV, lw=Pt(0.5))
    rect(s, cx + Inches(0.08), vy, Inches(0.48), Inches(0.78), fill=APP_ACC)
    tb(s, vname, cx + Inches(0.62), vy + Inches(0.07),
       cw - Inches(0.76), Inches(0.22), size=7, bold=True, color=APP_TXT)
    tb(s, vprice, cx + Inches(0.62), vy + Inches(0.30),
       cw - Inches(0.76), Inches(0.18), size=6.5, color=APP_ACC, bold=True)
    tb(s, vrate + " | нээлттэй", cx + Inches(0.62), vy + Inches(0.50),
       cw - Inches(0.76), Inches(0.18), size=6, color=APP_TXT2)
rect(s, cx, cy + ch - Inches(0.46), cw, Inches(0.46), fill=APP_SURF)
for ni, nl in enumerate(["Нүүр", "Захиалга", "Профайл", "Chat"]):
    nx = cx + ni * (cw // 4)
    if ni == 0:
        rect(s, nx, cy + ch - Inches(0.46), cw // 4, Inches(0.46), fill=APP_ACC)
    tb(s, nl, nx, cy + ch - Inches(0.26), cw // 4, Inches(0.22),
       size=5, color=WHITE if ni == 0 else APP_TXT2, align=PP_ALIGN.CENTER)
tb(s, "Нүүр дэлгэц", p2x, PH_Y + PH_H + Inches(0.08), PH_W, Inches(0.28),
   size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── Phone 3: Booking Sheet ────────────────────────────────────────────────────
p3x = PH_START + 2 * (PH_W + PH_GAP)
cx, cy, cw, ch = ph(s, p3x, PH_Y, PH_W, PH_H)
rect(s, cx, cy, cw, ch, fill=APP_BG)
sheet_y = cy + Inches(0.70)
rect(s, cx, sheet_y, cw, ch - Inches(0.70), fill=WHITE)
rect(s, cx + cw // 2 - Inches(0.24), sheet_y + Inches(0.07),
     Inches(0.48), Inches(0.06), fill=APP_DIV)
rect(s, cx + Inches(0.10), sheet_y + Inches(0.20), Inches(0.32), Inches(0.32), fill=APP_ACC)
tb(s, "Говийн Арена", cx + Inches(0.48), sheet_y + Inches(0.22),
   cw - Inches(0.58), Inches(0.18), size=7, bold=True, color=APP_TXT)
tb(s, "15,000₮/цаг", cx + Inches(0.48), sheet_y + Inches(0.38),
   cw - Inches(0.58), Inches(0.16), size=6, color=APP_ACC)
half_w = (cw - Inches(0.20)) // 2
rect(s, cx + Inches(0.08), sheet_y + Inches(0.64), half_w, Inches(0.28), fill=APP_ACC)
tb(s, "Бүтэн", cx + Inches(0.08), sheet_y + Inches(0.70), half_w, Inches(0.18),
   size=6, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, cx + Inches(0.08) + half_w + Inches(0.04), sheet_y + Inches(0.64),
     half_w, Inches(0.28), fill=WHITE, line=APP_DIV, lw=Pt(0.5))
tb(s, "Хагас", cx + Inches(0.08) + half_w + Inches(0.04), sheet_y + Inches(0.70),
   half_w, Inches(0.18), size=6, color=APP_TXT2, align=PP_ALIGN.CENTER)
dates = ["Да", "Мя", "Лх", "Пү", "Ба"]
dw = (cw - Inches(0.16)) // 5
for di2, dv in enumerate(dates):
    dx2 = cx + Inches(0.08) + di2 * dw
    rect(s, dx2, sheet_y + Inches(1.04), dw - Inches(0.04), Inches(0.38),
         fill=APP_ACC if di2 == 2 else WHITE, line=APP_DIV, lw=Pt(0.4))
    tb(s, dv, dx2, sheet_y + Inches(1.06), dw - Inches(0.04), Inches(0.16),
       size=5, color=WHITE if di2 == 2 else APP_TXT2, align=PP_ALIGN.CENTER)
    tb(s, str(3 + di2), dx2, sheet_y + Inches(1.20), dw - Inches(0.04), Inches(0.18),
       size=6, bold=True, color=WHITE if di2 == 2 else APP_TXT, align=PP_ALIGN.CENTER)
slot_times = ["08:00", "09:00", "10:00", "11:00", "14:00", "15:00", "16:00", "17:00", "18:00"]
sw2 = (cw - Inches(0.16)) // 3
for si2, st in enumerate(slot_times):
    sc2 = si2 % 3; sr2 = si2 // 3
    slx = cx + Inches(0.08) + sc2 * sw2
    sly = sheet_y + Inches(1.54) + sr2 * Inches(0.30)
    is_sel = si2 in [4, 5]; is_bkd = si2 in [1, 2]
    fc = APP_ACC if is_sel else (APP_SURF if is_bkd else WHITE)
    rect(s, slx, sly, sw2 - Inches(0.03), Inches(0.26),
         fill=fc, line=APP_DIV, lw=Pt(0.4))
    tb(s, st, slx, sly + Inches(0.05), sw2 - Inches(0.03), Inches(0.18),
       size=5.5, color=WHITE if is_sel else (APP_TXT2 if is_bkd else APP_TXT),
       align=PP_ALIGN.CENTER)
rect(s, cx + Inches(0.08), cy + ch - Inches(0.46), cw - Inches(0.16), Inches(0.32), fill=APP_ACC)
tb(s, "Захиалах", cx + Inches(0.08), cy + ch - Inches(0.42),
   cw - Inches(0.16), Inches(0.24), size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Захиалгын хуудас", p3x, PH_Y + PH_H + Inches(0.08), PH_W, Inches(0.28),
   size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── Phone 4: Confirmation ─────────────────────────────────────────────────────
p4x = PH_START + 3 * (PH_W + PH_GAP)
cx, cy, cw, ch = ph(s, p4x, PH_Y, PH_W, PH_H)
rect(s, cx, cy, cw, Inches(0.36), fill=APP_SURF)
tb(s, "Баталгаажуулалт", cx, cy + Inches(0.07), cw, Inches(0.24),
   size=7.5, bold=True, color=APP_TXT, align=PP_ALIGN.CENTER)
rect(s, cx + Inches(0.08), cy + Inches(0.44), cw - Inches(0.16), Inches(1.06), fill=APP_ACC)
tb(s, "Говийн Арена", cx + Inches(0.14), cy + Inches(0.52),
   cw - Inches(0.28), Inches(0.22), size=7.5, bold=True, color=WHITE)
tb(s, "Бүтэн талбай  •  14:00 – 16:00", cx + Inches(0.14), cy + Inches(0.72),
   cw - Inches(0.28), Inches(0.18), size=6, color=RGBColor(0xFF, 0xDD, 0xBB))
rect(s, cx + Inches(0.10), cy + Inches(0.96), cw - Inches(0.20), Inches(0.02), fill=WHITE)
tb(s, "30,000₮", cx + Inches(0.14), cy + Inches(1.02),
   cw - Inches(0.28), Inches(0.24), size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
tb(s, "Таны мэдээлэл", cx + Inches(0.10), cy + Inches(1.60), cw - Inches(0.20), Inches(0.20),
   size=6.5, bold=True, color=APP_TXT)
rect(s, cx + Inches(0.08), cy + Inches(1.82), cw - Inches(0.16), Inches(0.26),
     fill=WHITE, line=APP_DIV, lw=Pt(0.8))
tb(s, "Нэр оруулах...", cx + Inches(0.14), cy + Inches(1.88),
   cw - Inches(0.28), Inches(0.16), size=6, color=APP_TXT2, italic=True)
rect(s, cx + Inches(0.08), cy + Inches(2.12), cw - Inches(0.16), Inches(0.26),
     fill=WHITE, line=APP_DIV, lw=Pt(0.8))
tb(s, "+976 8800 1234", cx + Inches(0.14), cy + Inches(2.18),
   cw - Inches(0.28), Inches(0.16), size=6, color=APP_TXT2)
tb(s, "Төлбөрийн хэлбэр", cx + Inches(0.10), cy + Inches(2.50), cw - Inches(0.20), Inches(0.20),
   size=6.5, bold=True, color=APP_TXT)
pay_w = (cw - Inches(0.24)) // 3
for pi2, pn in enumerate(["QPay", "Карт", "Бэлэн"]):
    px2 = cx + Inches(0.08) + pi2 * (pay_w + Inches(0.04))
    rect(s, px2, cy + Inches(2.74), pay_w, Inches(0.34),
         fill=APP_ACC if pi2 == 0 else WHITE, line=APP_DIV, lw=Pt(0.5))
    tb(s, pn, px2, cy + Inches(2.82), pay_w, Inches(0.18),
       size=5.5, color=WHITE if pi2 == 0 else APP_TXT2, align=PP_ALIGN.CENTER)
rect(s, cx + Inches(0.08), cy + ch - Inches(0.46), cw - Inches(0.16), Inches(0.32), fill=APP_ACC)
tb(s, "Баталгаажуулах", cx + Inches(0.08), cy + ch - Inches(0.42),
   cw - Inches(0.16), Inches(0.24), size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Баталгаажуулалт", p4x, PH_Y + PH_H + Inches(0.08), PH_W, Inches(0.28),
   size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── Phone 5: My Bookings ──────────────────────────────────────────────────────
p5x = PH_START + 4 * (PH_W + PH_GAP)
cx, cy, cw, ch = ph(s, p5x, PH_Y, PH_W, PH_H)
rect(s, cx, cy, cw, Inches(0.36), fill=APP_SURF)
tb(s, "Миний захиалгууд", cx, cy + Inches(0.07), cw, Inches(0.24),
   size=7.5, bold=True, color=APP_TXT, align=PP_ALIGN.CENTER)
rect(s, cx, cy + Inches(0.42), cw, Inches(0.30), fill=APP_SURF)
for ti, tn in enumerate(["Бүгд", "Идэвхтэй", "Хүлээгдэж", "Дууссан"]):
    tx = cx + ti * (cw // 4)
    if ti == 0:
        rect(s, tx, cy + Inches(0.42), cw // 4, Inches(0.30), fill=APP_ACC)
    tb(s, tn, tx, cy + Inches(0.50), cw // 4, Inches(0.18),
       size=5, color=WHITE if ti == 0 else APP_TXT2, align=PP_ALIGN.CENTER)
bkgs = [
    ("Говийн Арена",    "Хүлээгдэж буй", AMBER),
    ("Стадионы Заал",   "Дууссан",        GRAY),
    ("Өмнөговь Спорт",  "Цуцлагдсан",    RED),
]
for bi, (bn, bs, bc) in enumerate(bkgs):
    by = cy + Inches(0.82) + bi * Inches(0.80)
    rect(s, cx + Inches(0.08), by, cw - Inches(0.16), Inches(0.72),
         fill=APP_CARD, line=APP_DIV, lw=Pt(0.5))
    rect(s, cx + Inches(0.08), by, Inches(0.36), Inches(0.72), fill=APP_ACC)
    tb(s, bn, cx + Inches(0.50), by + Inches(0.07),
       cw - Inches(0.62), Inches(0.20), size=6.5, bold=True, color=APP_TXT)
    rect(s, cx + Inches(0.50), by + Inches(0.28), Inches(0.54), Inches(0.18), fill=bc)
    tb(s, bs, cx + Inches(0.50), by + Inches(0.30),
       Inches(0.54), Inches(0.14), size=4.5, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, "2026-04-0" + str(bi + 3) + "  14:00", cx + Inches(0.50), by + Inches(0.50),
       cw - Inches(0.62), Inches(0.16), size=5.5, color=APP_TXT2)
rect(s, cx, cy + ch - Inches(0.44), cw, Inches(0.44), fill=APP_SURF)
for ni, nl in enumerate(["Нүүр", "Захиалга", "Профайл", "Chat"]):
    nx = cx + ni * (cw // 4)
    if ni == 1:
        rect(s, nx, cy + ch - Inches(0.44), cw // 4, Inches(0.44), fill=APP_ACC)
    tb(s, nl, nx, cy + ch - Inches(0.26), cw // 4, Inches(0.22),
       size=5, color=WHITE if ni == 1 else APP_TXT2, align=PP_ALIGN.CENTER)
tb(s, "Миний захиалгууд", p5x, PH_Y + PH_H + Inches(0.08), PH_W, Inches(0.28),
   size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9  —  Захиалгын үйл явц
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Захиалгын үйл явц ба давхцал шалгах логик")
content_bg(s)
steps9 = [
    ("1", "Заал\nсонгох",    BLUE,    "SportVenue\nобъект"),
    ("2", "Өдөр\nсонгох",    NAVY,    "DateTime\nпараметр"),
    ("3", "Цаг\nсонгох",     GREEN_D, "TimeSlot\nжагсаалт"),
    ("4", "Талбай\nсонгох",  PURPLE,  "Full / Half\ncourt type"),
    ("5", "Баталгаажуулах",  AMBER,   "POST\n/api/bookings"),
    ("6", "Код & мэдэгдэл",  RED,     "DBK-XXXXX\nPush notif"),
]
SW = Inches(1.82); SH = Inches(1.88); SY = Inches(1.02)
for i, (num, label, color, sub) in enumerate(steps9):
    sx = Inches(0.42) + i * Inches(2.08)
    rect(s, sx, SY, SW, SH, fill=color)
    rect(s, sx + Inches(0.68), SY + Inches(0.08), Inches(0.46), Inches(0.46), fill=WHITE)
    tb(s, num, sx + Inches(0.68), SY + Inches(0.08), Inches(0.46), Inches(0.46),
       size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(s, label, sx, SY + Inches(0.60), SW, Inches(0.70),
       size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, sub, sx, SY + Inches(1.30), SW, Inches(0.52),
       size=8.5, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER, italic=True)
    if i < 5:
        h_arrow(s, sx + SW, SY + SH // 2, sx + SW + Inches(0.24), color=GRAY)

rect(s, Inches(0.42), Inches(3.18), Inches(5.90), Inches(2.54),
     fill=GRAY_L, line=NAVY, lw=Pt(1.2))
rect(s, Inches(0.42), Inches(3.18), Inches(0.08), Inches(2.54), fill=NAVY)
tb(s, "Давхцал шалгах алгоритм  (Conflict Detection Logic)",
   Inches(0.60), Inches(3.26), Inches(5.62), Inches(0.38),
   size=12, bold=True, color=NAVY)
logic = [
    "1.  Firestore query:  venueId + dateKey + timeSlot + status IN [upcoming, active]",
    "2.  Хэрэв бүтэн талбай (fullCourt): ямар ч захиалга байвал  → 409 Conflict",
    "3.  Хэрэв хагас талбай (halfCourt): бүтэн зах. байвал  → 409 Conflict",
    "4.  Хагас × 2 байвал  → 409 Conflict  (бүтэн болсон тул)",
    "5.  Нөхцөл хангагдвал  doc.add(...)  → 201 Created  +  DBK-XXXXX код",
]
for li, line in enumerate(logic):
    tb(s, line, Inches(0.60), Inches(3.72) + li * Inches(0.38),
       Inches(5.62), Inches(0.34), size=10.5, color=GRAY_D)

rect(s, Inches(6.56), Inches(3.18), Inches(6.40), Inches(2.54),
     fill=GRAY_L, line=NAVY, lw=Pt(1.2))
rect(s, Inches(6.56), Inches(3.18), Inches(0.08), Inches(2.54), fill=GREEN_D)
tb(s, "Захиалгын статусын шилжилт  (State Machine)",
   Inches(6.74), Inches(3.26), Inches(6.12), Inches(0.38),
   size=12, bold=True, color=NAVY)
states = [
    ("upcoming",  BLUE,    Inches(7.00), Inches(3.78)),
    ("active",    GREEN_D, Inches(9.20), Inches(3.78)),
    ("completed", GRAY_D,  Inches(11.2), Inches(3.78)),
    ("cancelled", RED,     Inches(9.20), Inches(4.58)),
]
for st, col, bx, by in states:
    rect(s, bx, by, Inches(1.84), Inches(0.54), fill=col)
    tb(s, st, bx, by + Inches(0.08), Inches(1.84), Inches(0.38),
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
h_arrow(s, Inches(8.84), Inches(4.05), Inches(9.20), color=ARROW_COL)
h_arrow(s, Inches(11.04), Inches(4.05), Inches(11.2), color=ARROW_COL)
tb(s, "cron", Inches(8.84), Inches(3.82), Inches(0.38), Inches(0.22), size=8, color=GRAY, italic=True)
tb(s, "cron", Inches(11.02), Inches(3.82), Inches(0.38), Inches(0.22), size=8, color=GRAY, italic=True)
v_arrow(s, Inches(10.12), Inches(4.32), Inches(4.58), color=RED)
tb(s, "user cancel", Inches(10.16), Inches(4.38), Inches(1.0), Inches(0.22), size=8, color=RED, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10  —  Backend ба REST API
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Backend архитектур ба REST API")
content_bg(s)
tb(s, "REST API Endpoint-уудын тодорхойлолт",
   Inches(0.50), Inches(0.90), Inches(8.10), Inches(0.38),
   size=15, bold=True, color=NAVY)
endpoints = [
    ("POST", "/api/bookings",        "200 / 409", "Захиалга үүсгэх — давхцал шалгаад Firestore-д бичнэ"),
    ("GET",  "/api/bookings?phone=", "200 / 404", "Утасны дугаараар хэрэглэгчийн захиалгын жагсаалт"),
    ("GET",  "/api/schedule",        "200",       "Тухайн заалны өдөрт бодит цагийн слотуудын статус"),
    ("POST", "/api/users",           "200 / 201", "Хэрэглэгч олдвол шинэчлэх, олдохгүй бол шинэ үүсгэх"),
    ("GET",  "/api/users/:phone",    "200 / 404", "Хэрэглэгчийн профайлын мэдээлэл авах"),
    ("POST", "/api/chat",            "200 / 500", "Groq / LLaMA 3.1 8B-д мессэж дамжуулж хариулт авах"),
]
MC = {"POST": BLUE, "GET": GREEN_D}
for i, (method, path, status, desc) in enumerate(endpoints):
    ry = Inches(1.38) + i * Inches(0.88)
    mc = MC.get(method, GRAY)
    rect(s, Inches(0.50), ry, Inches(8.10), Inches(0.76), fill=WHITE, line=DIVIDER, lw=Pt(1))
    rect(s, Inches(0.50), ry, Inches(0.06), Inches(0.76), fill=mc)
    rect(s, Inches(0.58), ry + Inches(0.14), Inches(0.88), Inches(0.46), fill=mc)
    tb(s, method, Inches(0.58), ry + Inches(0.18), Inches(0.88), Inches(0.38),
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, path, Inches(1.58), ry + Inches(0.10), Inches(3.20), Inches(0.36),
       size=12, bold=True, color=NAVY, wrap=False)
    tb(s, status, Inches(4.88), ry + Inches(0.10), Inches(0.80), Inches(0.36),
       size=10, color=GREEN, bold=True)
    tb(s, desc, Inches(5.72), ry + Inches(0.10), Inches(2.78), Inches(0.56),
       size=11, color=GRAY_D, wrap=True)
rect(s, Inches(8.90), Inches(0.90), Inches(4.10), Inches(5.82), fill=NAVY)
tb(s, "Архитектурын шийдэл",
   Inches(9.00), Inches(1.00), Inches(3.88), Inches(0.40),
   size=13, bold=True, color=YELLOW)
notes = [
    ("Middleware Stack",  "cors()  →  express.json()\nenv-based config"),
    ("Error Handling",    "400 validation  /  409 conflict\n500 server error"),
    ("Cron Scheduler",    "node-cron тогтмол ажиллана\nFirestore batch write"),
    ("AI Integration",    "Groq OpenAI-compatible API\nllama-3.1-8b-instant"),
    ("Security",          "HTTPS,  env secrets\nFirebase Admin auth"),
]
for ni, (t, d) in enumerate(notes):
    ny = Inches(1.50) + ni * Inches(0.98)
    rect(s, Inches(9.02), ny, Inches(3.78), Inches(0.84),
         fill=RGBColor(0x14, 0x28, 0x6A), line=BLUE_LIGHT, lw=Pt(0.75))
    rect(s, Inches(9.02), ny, Inches(0.06), Inches(0.84), fill=ORANGE)
    tb(s, t, Inches(9.14), ny + Inches(0.07), Inches(3.54), Inches(0.32),
       size=11, bold=True, color=WHITE)
    tb(s, d, Inches(9.14), ny + Inches(0.42), Inches(3.54), Inches(0.38),
       size=9.5, color=RGBColor(0xAA, 0xBB, 0xEE))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11  —  AI Chat
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "AI туслах: Groq / LLaMA 3.1 нэгтгэл")
content_bg(s)
rect(s, Inches(0.42), Inches(0.90), Inches(5.60), Inches(5.82),
     fill=GRAY_L, line=NAVY, lw=Pt(1.5))
tb(s, "Харилцааны дохионы урсгал  (Message Flow)",
   Inches(0.58), Inches(1.02), Inches(5.28), Inches(0.40),
   size=13, bold=True, color=NAVY)
flow = [
    ("1", "Хэрэглэгч",    "ChatScreen-д мессэж бичнэ",    BLUE),
    ("2", "Flutter HTTP", "POST /api/chat  { messages }",  BLUE_LIGHT),
    ("3", "Express.js",   "System prompt + user content",  GREEN_D),
    ("4", "Groq API",     "llama-3.1-8b-instant model",    PURPLE),
    ("5", "Express.js",   "response.choices[0].message",   GREEN_D),
    ("6", "Flutter UI",   "ChatScreen-д хариулт харуулна", BLUE),
]
for fi, (num, actor, action, color) in enumerate(flow):
    fy = Inches(1.52) + fi * Inches(0.72)
    rect(s, Inches(0.58), fy, Inches(0.40), Inches(0.58), fill=color)
    tb(s, num, Inches(0.58), fy + Inches(0.12), Inches(0.40), Inches(0.36),
       size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, Inches(1.04), fy, Inches(1.40), Inches(0.58), fill=WHITE, line=DIVIDER, lw=Pt(0.75))
    tb(s, actor, Inches(1.06), fy + Inches(0.10), Inches(1.36), Inches(0.38),
       size=10, bold=True, color=color)
    rect(s, Inches(2.50), fy, Inches(3.36), Inches(0.58), fill=WHITE, line=DIVIDER, lw=Pt(0.75))
    tb(s, action, Inches(2.56), fy + Inches(0.10), Inches(3.24), Inches(0.38),
       size=10, color=GRAY_D)
    if fi < 5:
        v_arrow(s, Inches(0.78), fy + Inches(0.58), fy + Inches(1.30), color=color)

rect(s, Inches(6.22), Inches(0.90), Inches(6.78), Inches(2.74), fill=NAVY)
tb(s, "System Prompt агуулга",
   Inches(6.38), Inches(1.00), Inches(6.44), Inches(0.40),
   size=13, bold=True, color=YELLOW)
for pi, pt in enumerate([
    "Байгууллагын нэр, чиг үүрэг: Говийн Спорт",
    "5 заалны нэр, байршил, үнэ, нээлттэй цаг",
    "Цагийн хуваарь: 08:00 — 00:00  (16 слот)",
    "Захиалгын үйл явцын алхамууд",
    "Монгол хэлээр товч, нарийвчлалтай хариулах",
]):
    tb(s, pt, Inches(6.38), Inches(1.48) + pi * Inches(0.38),
       Inches(6.44), Inches(0.34), size=11, color=WHITE)

rect(s, Inches(6.22), Inches(3.78), Inches(6.78), Inches(2.94),
     fill=GRAY_L, line=DIVIDER, lw=Pt(1.2))
tb(s, "Загварын техникийн тодорхойлолт",
   Inches(6.38), Inches(3.88), Inches(6.44), Inches(0.40),
   size=13, bold=True, color=NAVY)
for mi, (k, v) in enumerate([
    ("Загвар",         "meta-llama/Llama-3.1-8B-Instruct"),
    ("Параметр",       "8 тэрбум  (8B)"),
    ("Context window", "128,000 token"),
    ("API endpoint",   "api.groq.com/openai/v1"),
    ("Хариултын хурд", "llama-3.1-8b-instant — маш хурдан"),
    ("Хэл дэмжлэг",   "Монгол хэлийг дэмжинэ (multilingual)"),
]):
    my = Inches(4.34) + mi * Inches(0.40)
    tb(s, k + ":", Inches(6.38), my, Inches(2.00), Inches(0.36),
       size=11, bold=True, color=NAVY)
    tb(s, v, Inches(8.42), my, Inches(4.42), Inches(0.36), size=11, color=GRAY_D)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12  —  Firebase
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Firebase интеграц ба мэдэгдлийн систем")
content_bg(s)
fb_panels = [
    (AMBER, "Firebase Firestore  —  Мэдээллийн сан",
     "Collections болон индексийн бүрдэл:\n\n"
     "  bookings:  venueId, dateKey, timeSlot, courtType, status,\n"
     "             userName, userPhone, code, price, createdAt\n\n"
     "  users:     phone, name, uid, createdAt\n\n"
     "Compound index: venueId + dateKey + timeSlot + status\n"
     "(давхцал шалгах WHERE query-г хурдасгадаг)"),
    (RGBColor(0xBF, 0x36, 0x00), "Firebase Authentication",
     "Баталгаажуулалтын урсгал:\n\n"
     "  1.  Хэрэглэгч утасны дугаар оруулна\n"
     "  2.  Firebase OTP SMS илгээнэ\n"
     "  3.  6 оронтой кодыг баталгаажуулна\n"
     "  4.  UID болон ID token олгогдоно\n"
     "  5.  Token-ийг Firestore users collection-д хадгална\n\n"
     "Нэвтрэлт шаардлагатай таб: Захиалга, Профайл"),
]
for i, (color, title, body) in enumerate(fb_panels):
    x = Inches(0.42) + i * Inches(6.48)
    rect(s, x, Inches(0.90), Inches(6.18), Inches(5.82), fill=WHITE, line=DIVIDER, lw=Pt(1.2))
    rect(s, x, Inches(0.90), Inches(6.18), Inches(0.52), fill=color)
    rect(s, x, Inches(0.90), Inches(0.08), Inches(5.82), fill=color)
    tb(s, title, x + Inches(0.18), Inches(0.96), Inches(5.90), Inches(0.40),
       size=14, bold=True, color=WHITE)
    tb(s, body, x + Inches(0.18), Inches(1.52), Inches(5.90), Inches(5.10),
       size=11, color=GRAY_D, wrap=True)
rect(s, Inches(0.42), Inches(5.76), Inches(12.48), Inches(1.00), fill=NAVY)
tb(s, "Firebase Cloud Messaging (FCM)  —  Push мэдэгдэл",
   Inches(0.58), Inches(5.84), Inches(4.0), Inches(0.36),
   size=12, bold=True, color=YELLOW)
for fi, ev in enumerate(["Захиалга баталгаажих", "Цагаас 30 мин өмнө", "Захиалга дуусмагц", "Промо мэдэгдэл"]):
    fx = Inches(4.90) + fi * Inches(2.00)
    rect(s, fx, Inches(5.84), Inches(1.84), Inches(0.70),
         fill=RGBColor(0x18, 0x2C, 0x6A), line=BLUE_LIGHT, lw=Pt(0.75))
    tb(s, ev, fx + Inches(0.08), Inches(5.92), Inches(1.68), Inches(0.52),
       size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13  —  CI/CD
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "CI/CD Pipeline ба үүлэн дэд бүтэц")
content_bg(s)
steps13 = [
    ("git push\nmain",  NAVY,                      "Developer\nworkstation"),
    ("GitHub\nActions", RGBColor(0x24, 0x29, 0x3E), "YAML workflow\n.github/"),
    ("Docker\nbuild",   BLUE,                       "node:20-alpine\nimage"),
    ("ECR\npush",       RGBColor(0x0D, 0x47, 0xA1), "AWS Elastic\nContainer Reg."),
    ("EC2\ndeploy",     GREEN_D,                    "ap-southeast-1\nSingapore"),
    ("Health\ncheck",   AMBER,                      "Port 3000\nliveness probe"),
]
PW = Inches(1.88); PH_C = Inches(2.10); PY13 = Inches(1.00)
for di, (label, color, sub) in enumerate(steps13):
    px = Inches(0.38) + di * Inches(2.16)
    rect(s, px, PY13, PW, PH_C, fill=color)
    tb(s, str(di + 1), px, PY13 + Inches(0.08), PW, Inches(0.38),
       size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, label, px, PY13 + Inches(0.44), PW, Inches(0.82),
       size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, px + Inches(0.14), PY13 + Inches(1.24), PW - Inches(0.28), Inches(0.03), fill=WHITE)
    tb(s, sub, px, PY13 + Inches(1.32), PW, Inches(0.64),
       size=9, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER, italic=True)
    if di < 5:
        h_arrow(s, px + PW, PY13 + PH_C // 2, px + PW + Inches(0.26), color=GRAY)
rect(s, Inches(0.38), Inches(3.32), Inches(5.88), Inches(2.96),
     fill=GRAY_L, line=NAVY, lw=Pt(1.2))
rect(s, Inches(0.38), Inches(3.32), Inches(0.08), Inches(2.96), fill=NAVY)
tb(s, "Docker container тодорхойлолт",
   Inches(0.56), Inches(3.40), Inches(5.64), Inches(0.38),
   size=12, bold=True, color=NAVY)
for li, line in enumerate([
    "FROM node:20-alpine",
    "WORKDIR /app",
    "COPY package*.json ./  &&  npm ci --only=production",
    "COPY . .",
    "EXPOSE 3000",
    'CMD ["node", "server.js"]',
]):
    rect(s, Inches(0.56), Inches(3.88) + li * Inches(0.36),
         Inches(5.54), Inches(0.32), fill=RGBColor(0x1A, 0x1A, 0x2E))
    tb(s, line, Inches(0.70), Inches(3.88) + li * Inches(0.36),
       Inches(5.26), Inches(0.30), size=9.5, color=RGBColor(0x7E, 0xC8, 0xFF))
rect(s, Inches(6.50), Inches(3.32), Inches(6.48), Inches(2.96),
     fill=GRAY_L, line=NAVY, lw=Pt(1.2))
rect(s, Inches(6.50), Inches(3.32), Inches(0.08), Inches(2.96), fill=GREEN_D)
tb(s, "Орчны хувьсагч  (.env secrets)",
   Inches(6.68), Inches(3.40), Inches(6.20), Inches(0.38),
   size=12, bold=True, color=NAVY)
for ei, (k, v) in enumerate([
    ("PORT",                     "3000"),
    ("GROQ_API_KEY",             "gsk_...  (Groq Console)"),
    ("FIREBASE_SERVICE_ACCOUNT", "JSON string (minified, 1 line)"),
]):
    ey = Inches(3.88) + ei * Inches(0.68)
    tb(s, k, Inches(6.68), ey, Inches(3.10), Inches(0.36), size=11, bold=True, color=GREEN_D)
    tb(s, v, Inches(6.68), ey + Inches(0.36), Inches(6.10), Inches(0.30),
       size=10, color=GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14  —  Хэрэгжилтийн явц
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Хэрэгжилтийн явц ба туршилтын үр дүн")
content_bg(s)
done = [
    ("Flutter мобайл апп",
     "8 дэлгэц бүрэн хөгжүүлсэн — HomeScreen, BookingSheet, ConfirmationScreen,"
     " BookingsScreen, ProfileScreen, ChatScreen, NotificationsScreen, DetailScreen"),
    ("Node.js REST API",
     "6 endpoint бүрэн хэрэгжсэн, production орчинд ажиллаж байна —"
     " bookings, schedule, users, chat"),
    ("Firebase нэгтгэл",
     "Firestore мэдээллийн сан, Phone Auth (OTP), FCM push notification"
     " бүрэн тохируулагдаж туршигдсан"),
    ("AI Chat туслах",
     "Groq / LLaMA 3.1 8B — system prompt-д монгол хэлний заалны мэдээлэл"
     " агуулуулсан, хариултын хурд < 1 секунд"),
    ("Захиалгын conflict логик",
     "Бүтэн / хагас талбайн давхцал шалгалт, Firestore compound index,"
     " 409 алдааны зохицуулалт туршигдсан"),
    ("AWS EC2 + CI/CD",
     "Docker container, AWS ECR, GitHub Actions pipeline бүрэн ажиллаж байна."
     " main branch push → автомат deploy"),
    ("Цагийн хуваарь",
     "08:00 — 00:00, нийт 16 цагийн слот (1 цагийн нэгж),"
     " node-cron автомат статус шинэчлэл ажиллаж байна"),
    ("UI / UX дизайн",
     "Material Design 3, custom AppTheme, warm beige palette,"
     " bottom navigation, modal bottom sheet бүрэн хэрэгжсэн"),
]
for i, (title, desc) in enumerate(done):
    col = i % 2; row = i // 2
    x = Inches(0.38) + col * Inches(6.50)
    y = Inches(0.98) + row * Inches(1.56)
    rect(s, x, y, Inches(6.22), Inches(1.40), fill=GREEN_L, line=GREEN, lw=Pt(1.5))
    rect(s, x + Inches(0.14), y + Inches(0.12), Inches(0.50), Inches(0.26), fill=GREEN)
    tb(s, "DONE", x + Inches(0.14), y + Inches(0.12), Inches(0.50), Inches(0.26),
       size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, x + Inches(0.76), y + Inches(0.10), Inches(5.30), Inches(0.36),
       size=13, bold=True, color=GREEN_D)
    tb(s, desc, x + Inches(0.18), y + Inches(0.54), Inches(5.90), Inches(0.78),
       size=10.5, color=GRAY_D, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15  —  Давуу тал ба Хязгаарлалт
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Системийн давуу тал ба хязгаарлалт")
content_bg(s)
rect(s, Inches(0.38), Inches(0.90), Inches(6.08), Inches(5.72), fill=WHITE, line=GREEN, lw=Pt(1.5))
rect(s, Inches(0.38), Inches(0.90), Inches(6.08), Inches(0.52), fill=GREEN_D)
tb(s, "Системийн давуу тал",
   Inches(0.54), Inches(0.96), Inches(5.84), Inches(0.40),
   size=15, bold=True, color=WHITE)
strengths = [
    ("Cross-platform бүрдэл",
     "Flutter нэг кодоос iOS болон Android хувилбарыг зэрэг бүрдүүлнэ"),
    ("Бодит цагийн давхцал хяналт",
     "Firestore compound index дээр суурилсан хурдан, найдвартай query"),
    ("Уян хатан захиалгын тогтолцоо",
     "Хагас / бүтэн талбайн логик — нэг слотод 2 бие даасан захиалга"),
    ("AI туслах систем",
     "Монгол хэлний system prompt бүхий chatbot хариултыг оновчтой болгоно"),
    ("Serverless + Cloud гибрид",
     "Firebase serverless + AWS EC2 хослолоор зардал ба хурд тэнцвэрждэг"),
    ("Автомат CI/CD pipeline",
     "GitHub Actions — код push хийхэд шууд production-д deploy хийнэ"),
    ("Масштаблах чадвар",
     "Firestore ба EC2 масштаблах чадвартай тул хэрэглэгчийн тоо өсөхөд нийцнэ"),
]
for si, (t, d) in enumerate(strengths):
    sy = Inches(1.52) + si * Inches(0.60)
    rect(s, Inches(0.54), sy, Inches(0.24), Inches(0.24), fill=GREEN)
    tb(s, t, Inches(0.86), sy, Inches(2.30), Inches(0.28), size=11, bold=True, color=GREEN_D)
    tb(s, d, Inches(0.86), sy + Inches(0.28), Inches(5.46), Inches(0.28), size=10, color=GRAY_D)

rect(s, Inches(6.66), Inches(0.90), Inches(6.32), Inches(5.72), fill=WHITE, line=RED, lw=Pt(1.5))
rect(s, Inches(6.66), Inches(0.90), Inches(6.32), Inches(0.52), fill=RED)
tb(s, "Хязгаарлалт ба цаашдын ажил",
   Inches(6.82), Inches(0.96), Inches(6.08), Inches(0.40),
   size=15, bold=True, color=WHITE)
limits = [
    ("Онлайн төлбөр байхгүй",
     "QPay / SocialPay нэгтгэл хийгдэх шаардлагатай"),
    ("Admin panel дутагдалтай",
     "Заалны захиалга удирдлага, орлогын тайлангийн веб самбар хэрэгтэй"),
    ("Load testing хийгдээгүй",
     "Нэгэн зэрэг олон хэрэглэгчийн ачааллын туршилт хийгдэх шаардлагатай"),
    ("Offline sync механизм байхгүй",
     "Интернэтгүй орчинд захиалга боловсруулах логик хэрэгтэй"),
    ("Analytics модуль байхгүй",
     "Хэрэглэгчийн зан байдлын дата, тайлангийн систем нэвтрүүлэх"),
    ("Зурагны хадгалалт",
     "Заалны зургийг Firebase Storage-т шилжүүлэх шаардлагатай"),
]
for li, (t, d) in enumerate(limits):
    ly = Inches(1.52) + li * Inches(0.70)
    rect(s, Inches(6.82), ly, Inches(0.24), Inches(0.24), fill=RED)
    tb(s, t, Inches(7.14), ly, Inches(2.50), Inches(0.28), size=11, bold=True, color=RED)
    tb(s, d, Inches(7.14), ly + Inches(0.28), Inches(5.68), Inches(0.28), size=10, color=GRAY_D)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16  —  Цаашдын чиглэл
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Цаашдын судалгааны чиглэл ба хөгжүүлэлт")
content_bg(s)
future = [
    ("QPay / SocialPay",
     "Монголын тогтвортой төлбөрийн системүүдтэй нэгтгэл хийж захиалгыг"
     " бүрэн автоматчилсан, цаасгүй орчин бүрдүүлэх."),
    ("Admin Web Dashboard",
     "Заалны эздэд зориулсан веб самбар: бодит цагийн хуваарь удирдлага,"
     " орлогын дата аналитик, захиалгын тайлан."),
    ("Үнэлгээ ба сэтгэгдэл",
     "Захиалга дуусмагц 5 одтой үнэлгээ, бичгэн сэтгэгдэл — машин"
     " суралцалтаар заалны чанарын индекс тодорхойлох."),
    ("Байршлын нэгтгэл",
     "Google Maps нэгтгэж заалны газрын зурагтай харуулах,"
     " хэрэглэгчийн байрлалаас ойр заалнуудыг зөвлөмжлөх."),
    ("AI сайжруулалт  (RAG)",
     "Retrieval-Augmented Generation аргачлалаар байнга шинэчлэгдэх"
     " заалны мэдлэгийн сантай холбосон ухаалаг туслах систем болгох."),
    ("Cross-platform өргөтгөл",
     "Flutter Web болон Tablet-д зориулсан responsive layout,"
     " заалны ажилтнуудад зориулсан desktop хувилбар."),
]
for i, (title, desc) in enumerate(future):
    col = i % 3; row = i // 3
    x = Inches(0.38) + col * Inches(4.32)
    y = Inches(1.00) + row * Inches(2.82)
    rect(s, x, y, Inches(4.10), Inches(2.62), fill=WHITE, line=DIVIDER, lw=Pt(1.2))
    rect(s, x, y, Inches(4.10), Inches(0.52), fill=NAVY)
    rect(s, x, y, Inches(0.08), Inches(2.62), fill=BLUE)
    tb(s, str(i + 1) + ".  " + title,
       x + Inches(0.18), y + Inches(0.09), Inches(3.82), Inches(0.38),
       size=13, bold=True, color=WHITE)
    tb(s, desc, x + Inches(0.18), y + Inches(0.66), Inches(3.82), Inches(1.84),
       size=11, color=GRAY_D, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 17  —  Дүгнэлт
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
header(s, "Дүгнэлт")
content_bg(s)
rect(s, Inches(0.38), Inches(0.88), Inches(12.56), Inches(1.14), fill=NAVY)
tb(s,
   "Энэхүү судалгааны ажлын хүрээнд үүлэн орчинд суурилсан, 3 давхаргат архитектур бүхий,"
   " AI туслахтай спортын заалны цаг захиалгын систем амжилттай загварчлагдаж хэрэгжлээ."
   " Flutter (Presentation), Node.js (Application), Firebase (Data) давхаргуудыг"
   " AWS EC2 дэд бүтцэд суурилуулан нэгтгэв.",
   Inches(0.56), Inches(0.96), Inches(12.18), Inches(1.00),
   size=12, color=WHITE)
summary = [
    (BLUE,    "Presentation\nTier",     "Flutter 3.x\n8 дэлгэц\niOS / Android"),
    (GREEN_D, "Application\nTier",     "Node.js 20\nExpress API\nGroq AI"),
    (AMBER,   "Data Tier",             "Firestore\nFirebase Auth\nFCM"),
    (NAVY,    "CI/CD\nInfrastructure", "AWS EC2\nDocker / ECR\nGitHub Actions"),
    (PURPLE,  "AI Integration",        "LLaMA 3.1 8B\nМонгол хэл\nSystem prompt"),
    (RED,     "Conflict Logic",        "Compound index\nHalf / full court\nState machine"),
]
for i, (color, title, sub) in enumerate(summary):
    col = i % 3; row = i // 2
    x = Inches(0.38) + col * Inches(4.32)
    y = Inches(2.22) + row * Inches(2.34)
    rect(s, x, y, Inches(4.10), Inches(2.16), fill=color)
    _h = str(color)
    _r, _g, _b = int(_h[0:2], 16), int(_h[2:4], 16), int(_h[4:6], 16)
    rect(s, x, y, Inches(4.10), Inches(0.50),
         fill=RGBColor(max(_r - 40, 0), max(_g - 40, 0), max(_b - 40, 0)))
    rect(s, x, y, Inches(0.08), Inches(2.16), fill=WHITE)
    tb(s, title, x + Inches(0.18), y + Inches(0.08), Inches(3.82), Inches(0.48),
       size=13, bold=True, color=WHITE)
    tb(s, sub, x + Inches(0.18), y + Inches(0.68), Inches(3.82), Inches(1.36),
       size=12, color=RGBColor(0xCC, 0xDD, 0xFF))
tb(s,
   "Цаашид онлайн төлбөр, admin самбар, RAG-д суурилсан AI,"
   " байршлын нэгтгэлээр системийг өргөтгөн судалгааны үр дүнг баяжуулах боломжтой.",
   Inches(0.38), Inches(6.70), Inches(12.56), Inches(0.40),
   size=11.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 18  —  Баярлалаа
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, W, H, fill=YELLOW)
rect(s, 0, 0, Inches(0.22), H, fill=NAVY)
rect(s, W - Inches(0.22), 0, Inches(0.22), H, fill=NAVY)
rect(s, 0, H - Inches(0.44), W, Inches(0.44), fill=NAVY)
s.shapes.add_picture(LOGO_PATH, Inches(0.32), Inches(0.14), Inches(0.74), Inches(0.74))
rect(s, Inches(1.5), Inches(1.56), Inches(10.3), Inches(0.04), fill=NAVY)
tb(s, "БАЯРЛАЛАА",
   Inches(1.5), Inches(1.70), Inches(10.3), Inches(1.28),
   size=62, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
rect(s, Inches(1.5), Inches(3.08), Inches(10.3), Inches(0.04), fill=NAVY)
tb(s, "Асуулт болон саналд хариулахад бэлэн байна.",
   Inches(1.5), Inches(3.22), Inches(10.3), Inches(0.54),
   size=18, color=NAVY, align=PP_ALIGN.CENTER)
tb(s, f"Илтгэгч:    {AUTHOR}",
   Inches(4.2), Inches(4.18), Inches(5.0), Inches(0.36),
   size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tb(s, f"Удирдагч:  {SUPERVISOR}",
   Inches(4.2), Inches(4.56), Inches(5.0), Inches(0.36),
   size=14, color=NAVY, align=PP_ALIGN.CENTER)
tb(s,
   "Үүлэн орчинд суурилсан спорт заалны цаг захиалгын"
   " 3 давхаргат системийн загварчлал, хэрэгжилт",
   Inches(2.0), Inches(5.20), Inches(9.3), Inches(0.48),
   size=12, color=NAVY, italic=True, align=PP_ALIGN.CENTER)
tb(s, f"Улаанбаатар хот,  {YEAR} он",
   Inches(0.32), H - Inches(0.42), Inches(5), Inches(0.36),
   size=12, color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
#  Save
# ══════════════════════════════════════════════════════════════════════════════
out = r"C:\Users\user\Documents\dz_zaal\Goviyn_Sport_Iltgel_v5.pptx"
prs.save(out)
sys.stdout.buffer.write(b"Saved OK\n")
sys.stdout.buffer.write(f"Slides: {len(prs.slides)}\n".encode())
sys.stdout.buffer.write((out + "\n").encode())
