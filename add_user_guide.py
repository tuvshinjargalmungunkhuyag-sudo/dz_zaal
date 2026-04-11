# -*- coding: utf-8 -*-
"""
User Guide слайдуудыг v5 PPT-д нэмж v6 болгох скрипт
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Өнгөний тодорхойлолт ────────────────────────────────────────────────────
PRIMARY   = RGBColor(0xED, 0xE7, 0xDA)   # #EDE7DA бежийн дэвсгэр
SECONDARY = RGBColor(0xD4, 0x70, 0x0F)   # #D4700F наранжин
ACCENT    = RGBColor(0xB8, 0x5C, 0x08)   # #B85C08 гүн наранжин
CARD      = RGBColor(0xF8, 0xF4, 0xEE)   # #F8F4EE картын дэвсгэр
TEXT      = RGBColor(0x1A, 0x12, 0x08)   # #1A1208 үндсэн текст
SUBTEXT   = RGBColor(0x6B, 0x5B, 0x4B)   # #6B5B4B дэд текст
DIVIDER   = RGBColor(0xD9, 0xCF, 0xC3)   # #D9CFC3 хуваагч
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS   = RGBColor(0x2D, 0x99, 0x63)   # ногоон
PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)   # нил ягаан

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Допоможний функцүүд ────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or TEXT
    return txBox


def slide_bg(slide, color=None):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or PRIMARY


def add_header_bar(slide, title, subtitle=None):
    """Наранжин дээд баар нэмэх"""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill=SECONDARY)
    add_text(slide, title,
             Inches(0.4), Inches(0.1), Inches(10), Inches(0.7),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.75), Inches(10), Inches(0.4),
                 size=14, color=RGBColor(0xFF, 0xE0, 0xC0))


def add_card(slide, x, y, w, h, title, lines, icon="•",
             accent=None, title_size=16, body_size=13):
    """Мэдээллийн карт нэмэх"""
    c = accent or SECONDARY
    add_rect(slide, x, y, w, h, fill=CARD, line=DIVIDER, line_w=Pt(1))
    # Зураас (зүүн талын хөндий)
    add_rect(slide, x, y, Inches(0.06), h, fill=c)
    # Гарчиг
    add_text(slide, f"{icon}  {title}",
             x + Inches(0.15), y + Inches(0.12),
             w - Inches(0.2), Inches(0.35),
             size=title_size, bold=True, color=c)
    # Агуулга
    body = "\n".join(lines) if isinstance(lines, list) else lines
    add_text(slide, body,
             x + Inches(0.15), y + Inches(0.5),
             w - Inches(0.2), h - Inches(0.6),
             size=body_size, color=TEXT)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД ҮҮСГЭГЧ ФУНКЦҮҮД
# ════════════════════════════════════════════════════════════════════════════

def add_ug_cover(prs):
    """UG-1: Хэрэглэгчийн гарын авлагын нэрийн хуудас"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide_bg(slide, PRIMARY)

    # Наранжин тал (зүүн 40%)
    add_rect(slide, 0, 0, Inches(5.3), SLIDE_H, fill=SECONDARY)

    # Логотип хайрцаг
    add_rect(slide, Inches(1.2), Inches(1.5), Inches(2.8), Inches(1.0),
             fill=RGBColor(0xFF, 0xFF, 0xFF), line=None)
    add_text(slide, "🏟️  Говийн Спорт",
             Inches(1.2), Inches(1.55), Inches(2.8), Inches(0.85),
             size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Гарчиг
    add_text(slide, "Хэрэглэгчийн\nГарын Авлага",
             Inches(0.3), Inches(2.8), Inches(4.6), Inches(1.6),
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Спорт заал захиалах алхам алхмаар",
             Inches(0.3), Inches(4.5), Inches(4.6), Inches(0.5),
             size=14, color=RGBColor(0xFF, 0xE0, 0xC0), align=PP_ALIGN.CENTER)

    # Баруун талын агуулга
    add_text(slide, "Энэ гарын авлагад:",
             Inches(5.7), Inches(1.2), Inches(7.0), Inches(0.4),
             size=18, bold=True, color=SECONDARY)

    items = [
        "📱  Апп-ын бүрэн танилцуулга",
        "🔐  Бүртгэл ба нэвтрэлтийн заавар",
        "📅  Захиалга хийх алхам алхмаар",
        "🏟️  4 спорт заалны мэдээлэл",
        "⏰  16 цагийн хуваарь (08:00–00:00)",
        "🤖  AI туслагч чатботын ашиглалт",
        "📋  Захиалга удирдах функцүүд",
    ]
    for i, item in enumerate(items):
        add_text(slide, item,
                 Inches(5.7), Inches(1.85) + i * Inches(0.6),
                 Inches(7.0), Inches(0.55),
                 size=14, color=TEXT)

    # Доод хэсэг
    add_rect(slide, 0, SLIDE_H - Inches(0.5), SLIDE_W, Inches(0.5),
             fill=ACCENT)
    add_text(slide, "ШУТИС • Даланзадгад хотын спорт заалны цаг захиалгын систем • 2025",
             Inches(0.3), SLIDE_H - Inches(0.45), Inches(12.7), Inches(0.4),
             size=11, color=WHITE, align=PP_ALIGN.CENTER)


def add_ug_overview(prs):
    """UG-2: Апп-ын танилцуулга"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, PRIMARY)
    add_header_bar(slide, "Апп-ын танилцуулга",
                   "Говийн Спорт — Даланзадгад хотын спорт заалны нэгдсэн захиалгын систем")

    # 6 товч карт
    items = [
        ("🏟️", "4 Спорт Заал", "Нэг платформд нэгдсэн"),
        ("📅", "16 Цагийн Үе", "08:00–00:00, өдөр бүр"),
        ("⚡", "Бодит Цагийн", "Firebase Firestore"),
        ("🤖", "AI Чатбот", "Claude Haiku + Groq"),
        ("🔔", "Сануулга", "1 цагийн өмнө"),
        ("🔐", "OTP Баталгаа", "6 оронтой email код"),
    ]
    cols, rows = 3, 2
    cw, ch = Inches(3.9), Inches(1.5)
    for i, (icon, title, desc) in enumerate(items):
        col = i % cols
        row = i // cols
        x = Inches(0.3) + col * (cw + Inches(0.25))
        y = Inches(1.5) + row * (ch + Inches(0.2))
        add_rect(slide, x, y, cw, ch, fill=CARD, line=DIVIDER, line_w=Pt(1))
        add_text(slide, icon, x + Inches(0.15), y + Inches(0.1),
                 Inches(0.6), Inches(0.5), size=28)
        add_text(slide, title, x + Inches(0.8), y + Inches(0.1),
                 cw - Inches(0.9), Inches(0.45),
                 size=16, bold=True, color=SECONDARY)
        add_text(slide, desc, x + Inches(0.8), y + Inches(0.6),
                 cw - Inches(0.9), Inches(0.4),
                 size=13, color=SUBTEXT)

    # Дизайны мэдээлэл
    add_rect(slide, Inches(0.3), Inches(5.0), Inches(12.6), Inches(0.9),
             fill=RGBColor(0xFF, 0xF0, 0xE0), line=SECONDARY, line_w=Pt(1))
    add_text(slide,
             "🎨  Дизайн: Material Design 3  •  Montserrat фонт  •  "
             "Цайвар говийн тон (#EDE7DA бежийн дэвсгэр, #D4700F наранжин)  •  Android & iOS",
             Inches(0.5), Inches(5.1), Inches(12.2), Inches(0.75),
             size=13, color=ACCENT, align=PP_ALIGN.CENTER)


def add_ug_venues(prs):
    """UG-3: Спорт заалнуудын жагсаалт"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, PRIMARY)
    add_header_bar(slide, "Спорт Заалнууд",
                   "Аппд бүртгэлтэй 4 спорт заал — нэг дор захиалах боломжтой")

    venues = [
        ("🟠", "Говийн Арена",        "1-р хороо",  "15,000₮/цаг", "⭐ 4.8 (124)", "Гардероб, Душ, Паркинг, WiFi",      RGBColor(0xFF,0x6B,0x35)),
        ("🟢", "Өмнөговь Спорт Заал", "2-р хороо",  "12,000₮/цаг", "⭐ 4.6 (89)",  "Гардероб, Душ, WiFi",               RGBColor(0x00,0xC8,0x96)),
        ("🟣", "Стадионы Спорт Заал", "Стадион",    "8,000₮/цаг",  "⭐ 4.5 (203)", "Гардероб, Душ, Гэрэлтүүлэг, Камер",RGBColor(0x6C,0x63,0xFF)),
        ("🔵", "Цэнтрийн Спорт Заал", "3-р хороо",  "18,000₮/цаг", "⭐ 4.9 (67)",  "Гардероб, Душ, Тренер, WiFi",        RGBColor(0x00,0xAA,0xFF)),
    ]

    for i, (icon, name, loc, price, rating, facilities, color) in enumerate(venues):
        y = Inches(1.45) + i * Inches(1.3)
        add_rect(slide, Inches(0.3), y, Inches(12.6), Inches(1.15),
                 fill=CARD, line=DIVIDER, line_w=Pt(1))
        # Зүүн хөндий өнгө
        add_rect(slide, Inches(0.3), y, Inches(0.15), Inches(1.15), fill=color)
        # Нэр
        add_text(slide, f"{icon}  {name}",
                 Inches(0.6), y + Inches(0.08), Inches(4.0), Inches(0.45),
                 size=17, bold=True, color=TEXT)
        # Байршил
        add_text(slide, f"📍 {loc}",
                 Inches(0.6), y + Inches(0.55), Inches(4.0), Inches(0.35),
                 size=12, color=SUBTEXT)
        # Үнэ
        add_text(slide, price,
                 Inches(4.8), y + Inches(0.08), Inches(2.0), Inches(0.45),
                 size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        # Үнэлгээ
        add_text(slide, rating,
                 Inches(6.8), y + Inches(0.08), Inches(2.2), Inches(0.45),
                 size=14, color=SUBTEXT, align=PP_ALIGN.CENTER)
        # Тоног
        add_text(slide, f"✓ {facilities}",
                 Inches(9.2), y + Inches(0.08), Inches(3.5), Inches(0.55),
                 size=12, color=SUBTEXT)


def add_ug_screens(prs):
    """UG-4: Дэлгэцүүдийн танилцуулга"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, PRIMARY)
    add_header_bar(slide, "Апп-ын Дэлгэцүүд",
                   "4 үндсэн таб: Нүүр | Захиалга | AI Чат | Профайл")

    screens = [
        ("🔐", "Нэвтрэх/\nБүртгэл",
         ["Email + нууц үг", "6 оронтой OTP код", "Rollback механизм"]),
        ("🏠", "Нүүр хуудас",
         ["4 заалны жагсаалт", "Хайлт + шүүлтүүр", "Үнэлгээгээр эрэмбэлэх"]),
        ("📅", "Захиалах\n(Bottom Sheet)",
         ["Бүтэн/Хагас талбай", "14 хоногийн хуваарь", "16 цагийн үе"]),
        ("✅", "Баталгаа",
         ["Нэр, email, төлбөр", "QPay/Карт/Бэлэн", "DBK-XXXXX дугаар"]),
        ("📋", "Захиалгууд",
         ["Бодит цагийн Firestore", "Статус: upcoming/active", "Цуцлах боломж"]),
        ("🤖", "AI Чат",
         ["Claude Haiku 4.5", "Groq Llama нөөц", "Асуулт-хариулт"]),
        ("👤", "Профайл",
         ["Нэр, email, статистик", "Нийт/хүлээгдэж буй", "Гарах функц"]),
        ("🔔", "Мэдэгдэл",
         ["1 цагийн өмнө сануулга", "Локал (FCM-гүй)", "Офлайн ажиллана"]),
    ]

    cw, ch = Inches(2.9), Inches(1.6)
    for i, (icon, name, desc) in enumerate(screens):
        col = i % 4
        row = i // 4
        x = Inches(0.3) + col * (cw + Inches(0.22))
        y = Inches(1.5) + row * (ch + Inches(0.18))
        add_rect(slide, x, y, cw, ch, fill=CARD, line=DIVIDER, line_w=Pt(1))
        add_rect(slide, x, y, cw, Inches(0.05), fill=SECONDARY)
        add_text(slide, icon, x + Inches(0.1), y + Inches(0.1),
                 Inches(0.5), Inches(0.45), size=24)
        add_text(slide, name, x + Inches(0.65), y + Inches(0.08),
                 cw - Inches(0.75), Inches(0.65),
                 size=13, bold=True, color=SECONDARY)
        body = "\n".join(f"• {d}" for d in desc)
        add_text(slide, body, x + Inches(0.1), y + Inches(0.75),
                 cw - Inches(0.2), Inches(0.8),
                 size=11, color=TEXT)


def add_ug_steps(prs):
    """UG-5: Захиалга хийх алхамууд"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, PRIMARY)
    add_header_bar(slide, "Захиалга Хийх Алхамууд",
                   "Хэдхэн товшилтоор спорт заал захиална")

    steps = [
        ("1", "Нэвтрэх / Бүртгүүлэх",
         "Email + нууц үг оруулна. Шинэ хэрэглэгч бол 6 оронтой OTP кодоор баталгаажуулна."),
        ("2", "Заал сонгох",
         "Нүүр хуудасны жагсаалтаас захиалах заалд дарна. Хайлт, шүүлтүүр ашиглаж болно."),
        ("3", "Талбай + Өдөр + Цаг",
         "Бүтэн/Хагас талбай сонгоно. 14 хоногоос өдөр, 16 цагийн үеэс цагаа сонгоно."),
        ("4", "Мэдээлэл + Төлбөр",
         "Нэр (автоматаар), email (автомат) шалгана. QPay/Карт/Бэлэн-ээс сонгоно."),
        ("5", "Баталгаажуулах",
         "'Захиалга баталгаажуулах' товч → DBK-XXXXX дугаар авна. 1 цагийн өмнө сануулга ирнэ."),
        ("6", "Захиалга харах / Цуцлах",
         "'Захиалга' табаас бүх захиалгаа харна. 'Хүлээгдэж буй' захиалгыг цуцлах боломжтой."),
    ]

    # 2 bagana, 3 мөр
    cw, ch = Inches(5.8), Inches(1.45)
    for i, (num, title, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = Inches(0.3) + col * (cw + Inches(0.8))
        y = Inches(1.45) + row * (ch + Inches(0.12))

        # Карт дэвсгэр
        add_rect(slide, x, y, cw, ch, fill=CARD, line=DIVIDER, line_w=Pt(1))

        # Дугаар тойрог
        add_rect(slide, x + Inches(0.1), y + Inches(0.1),
                 Inches(0.6), Inches(0.6), fill=SECONDARY)
        add_text(slide, num,
                 x + Inches(0.1), y + Inches(0.05),
                 Inches(0.6), Inches(0.7),
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Гарчиг
        add_text(slide, title,
                 x + Inches(0.85), y + Inches(0.1),
                 cw - Inches(0.95), Inches(0.45),
                 size=15, bold=True, color=SECONDARY)

        # Тайлбар
        add_text(slide, desc,
                 x + Inches(0.85), y + Inches(0.58),
                 cw - Inches(0.95), Inches(0.8),
                 size=12, color=TEXT)


def add_ug_features(prs):
    """UG-6: Гол функцүүд"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, PRIMARY)
    add_header_bar(slide, "Апп-ын Гол Функцүүд",
                   "Хэрэглэгчид санал болгож буй бүрэн функцийн жагсаалт")

    features_left = [
        ("🔐", "Email OTP баталгаажуулалт",
         "Бүртгэлийн үед Gmail SMTP-ээр 6 оронтой код. Алдаа → Firebase rollback."),
        ("🔍", "Хайлт ба шүүлтүүр",
         "Нэр/төрөл/байршлаар хайх, үнэ эсвэл үнэлгээгээр эрэмбэлэх."),
        ("⏰", "Олон цагийн захиалга",
         "Зэрэгцэх цагуудыг дараалан дарж 2+ цагийн захиалга хийнэ."),
        ("🟣", "Гэрээт цагууд",
         "Байгууллагуудын тогтмол захиалга нил ягаан өнгөөр харагдана."),
        ("❌", "Захиалга цуцлах",
         "'Хүлээгдэж буй' захиалгыг цуцлах боломжтой."),
    ]
    features_right = [
        ("🏟️", "Бүтэн / Хагас талбай",
         "Нэг цагт хоёр хагас талбай зэрэг захиалагдах боломжтой."),
        ("🔔", "Локал сануулга",
         "Захиалгаас 1 цагийн өмнө офлайн сануулга. FCM сервер шаардагдахгүй."),
        ("🤖", "AI туслагч",
         "Claude Haiku 4.5 үндсэн, Groq Llama 3.1-8b нөөц горим."),
        ("📊", "Профайл статистик",
         "Нийт болон хүлээгдэж буй захиалгын тоо профайлд харагдана."),
        ("⚙️", "Автомат статус",
         "Cron ажил 15 минут тутам захиалгын статусыг автоматаар шинэчилнэ."),
    ]

    ch = Inches(1.08)
    for i, (icon, title, desc) in enumerate(features_left):
        y = Inches(1.45) + i * (ch + Inches(0.1))
        add_rect(slide, Inches(0.3), y, Inches(6.0), ch,
                 fill=CARD, line=DIVIDER, line_w=Pt(1))
        add_text(slide, icon, Inches(0.4), y + Inches(0.1),
                 Inches(0.5), Inches(0.5), size=22)
        add_text(slide, title, Inches(0.95), y + Inches(0.08),
                 Inches(5.2), Inches(0.38), size=14, bold=True, color=SECONDARY)
        add_text(slide, desc, Inches(0.95), y + Inches(0.5),
                 Inches(5.2), Inches(0.5), size=12, color=TEXT)

    for i, (icon, title, desc) in enumerate(features_right):
        y = Inches(1.45) + i * (ch + Inches(0.1))
        add_rect(slide, Inches(6.85), y, Inches(6.0), ch,
                 fill=CARD, line=DIVIDER, line_w=Pt(1))
        add_text(slide, icon, Inches(6.95), y + Inches(0.1),
                 Inches(0.5), Inches(0.5), size=22)
        add_text(slide, title, Inches(7.5), y + Inches(0.08),
                 Inches(5.2), Inches(0.38), size=14, bold=True, color=SECONDARY)
        add_text(slide, desc, Inches(7.5), y + Inches(0.5),
                 Inches(5.2), Inches(0.5), size=12, color=TEXT)


def add_ug_tech(prs):
    """UG-7: Техникийн дэлгэрэнгүй"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, PRIMARY)
    add_header_bar(slide, "Техникийн Дэлгэрэнгүй",
                   "Системийн архитектур, технологийн стек, гүйцэтгэлийн хэмжилт")

    # ── Зүүн тал: Архитектур ────────────────────────────────────────────────
    add_text(slide, "3 Давхаргат Архитектур",
             Inches(0.3), Inches(1.35), Inches(6.0), Inches(0.4),
             size=16, bold=True, color=SECONDARY)

    layers = [
        ("📱", "Клиент давхарга",        "Flutter мобайл апп (Android & iOS)"),
        ("⚡", "Логикийн давхарга",      "Node.js/Express REST API (AWS EC2 Токио)"),
        ("🗄️", "Өгөгдлийн давхарга",    "Cloud Firestore (4 collection) + AWS ECR"),
    ]
    for i, (icon, title, desc) in enumerate(layers):
        y = Inches(1.9) + i * Inches(1.05)
        add_rect(slide, Inches(0.3), y, Inches(6.0), Inches(0.9),
                 fill=CARD, line=SECONDARY, line_w=Pt(1.5))
        add_text(slide, icon, Inches(0.45), y + Inches(0.1),
                 Inches(0.5), Inches(0.5), size=22)
        add_text(slide, title, Inches(1.0), y + Inches(0.05),
                 Inches(5.1), Inches(0.38), size=14, bold=True, color=SECONDARY)
        add_text(slide, desc, Inches(1.0), y + Inches(0.48),
                 Inches(5.1), Inches(0.38), size=12, color=TEXT)
        # Сум
        if i < 2:
            add_text(slide, "↓", Inches(2.8), y + Inches(0.9),
                     Inches(0.5), Inches(0.4), size=18, color=SECONDARY,
                     align=PP_ALIGN.CENTER)

    # Firestore collections
    add_text(slide, "Firestore (4 collection):",
             Inches(0.3), Inches(5.1), Inches(6.0), Inches(0.35),
             size=13, bold=True, color=TEXT)
    collections = ["bookings  •  users  •  fixed_bookings  •  email_verifications"]
    add_text(slide, collections[0],
             Inches(0.3), Inches(5.45), Inches(6.0), Inches(0.35),
             size=12, color=SUBTEXT)

    # ── Баруун тал: Хэмжилт ─────────────────────────────────────────────────
    add_text(slide, "Гүйцэтгэлийн Хэмжилт",
             Inches(6.8), Inches(1.35), Inches(6.0), Inches(0.4),
             size=16, bold=True, color=SECONDARY)

    metrics = [
        ("📱", "Апп ачаалах",    "2.3 сек"),
        ("🗄️", "Firestore",      "0.6 сек"),
        ("⚡", "EC2 API",        "0.8 сек"),
        ("🤖", "Claude Haiku",   "1.4 сек"),
        ("🦙", "Groq Llama",     "0.9 сек"),
        ("🔄", "CI/CD deploy",   "3–4 мин"),
    ]
    for i, (icon, label, val) in enumerate(metrics):
        y = Inches(1.88) + i * Inches(0.78)
        add_rect(slide, Inches(6.8), y, Inches(5.9), Inches(0.65),
                 fill=CARD, line=DIVIDER, line_w=Pt(1))
        add_text(slide, icon, Inches(6.9), y + Inches(0.05),
                 Inches(0.5), Inches(0.5), size=20)
        add_text(slide, label, Inches(7.5), y + Inches(0.1),
                 Inches(3.5), Inches(0.4), size=13, color=TEXT)
        add_text(slide, val,
                 Inches(11.0), y + Inches(0.05),
                 Inches(1.6), Inches(0.5),
                 size=16, bold=True, color=SECONDARY, align=PP_ALIGN.RIGHT)

    # CI/CD тайлбар
    add_rect(slide, Inches(6.8), Inches(6.6), Inches(5.9), Inches(0.55),
             fill=RGBColor(0xFF, 0xF0, 0xE0), line=SECONDARY, line_w=Pt(1))
    add_text(slide,
             "🔄  GitHub Actions: Push → Docker Build → ECR Push → EC2 Deploy",
             Inches(7.0), Inches(6.65), Inches(5.6), Inches(0.45),
             size=12, color=ACCENT)


def add_ug_closing(prs):
    """UG-8: Дүгнэлт"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, SECONDARY)

    # Цагаан хайрцаг
    add_rect(slide, Inches(1.5), Inches(1.0), Inches(10.3), Inches(5.5),
             fill=WHITE)

    add_text(slide, "🏟️",
             Inches(5.9), Inches(0.5), Inches(1.5), Inches(1.2),
             size=52, align=PP_ALIGN.CENTER)

    add_text(slide, "Говийн Спорт",
             Inches(1.8), Inches(1.3), Inches(9.8), Inches(0.8),
             size=34, bold=True, color=SECONDARY, align=PP_ALIGN.CENTER)

    add_text(slide,
             "Өмнөговь аймгийн Даланзадгад хотын спорт заалнуудын\n"
             "цаг захиалгыг хялбарчилсан мобайл аппликейшн",
             Inches(1.8), Inches(2.2), Inches(9.8), Inches(0.8),
             size=16, color=SUBTEXT, align=PP_ALIGN.CENTER)

    # Технологиуд
    techs = ["Flutter", "Firebase", "Node.js/Express", "Docker", "AWS EC2",
             "GitHub Actions", "Claude Haiku", "Groq Llama"]
    tw = Inches(1.4)
    start_x = (SLIDE_W - len(techs) * tw) / 2
    for i, tech in enumerate(techs):
        x = start_x + i * tw
        add_rect(slide, x, Inches(3.3), tw - Inches(0.1), Inches(0.5),
                 fill=RGBColor(0xFF, 0xF0, 0xE0), line=SECONDARY, line_w=Pt(1))
        add_text(slide, tech, x, Inches(3.32),
                 tw - Inches(0.1), Inches(0.45),
                 size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(slide, "Хэрэглэгчийн ДҮН:",
             Inches(1.8), Inches(4.1), Inches(9.8), Inches(0.4),
             size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(slide,
             "✅  10 функциональ тест бүгд тэнцсэн  •  "
             "✅  Апп 2.3 сек ачаалдаг  •  "
             "✅  Android & iOS-д ажилладаг",
             Inches(1.8), Inches(4.55), Inches(9.8), Inches(0.45),
             size=13, color=SUCCESS, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.7),
             fill=SECONDARY)
    add_text(slide,
             "ШУТИС  •  Мэдээлэл технологийн сургууль  •  "
             "Тувшинжаргал М.  •  2025",
             Inches(1.6), Inches(5.3), Inches(10.1), Inches(0.5),
             size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# ҮНДСЭН ПРОГРАМ
# ════════════════════════════════════════════════════════════════════════════
def main():
    src = r"C:\Users\user\Documents\dz_zaal\Goviyn_Sport_Iltgel_v5.pptx"
    dst = r"C:\Users\user\Documents\dz_zaal\Goviyn_Sport_Iltgel_v6.pptx"

    print("v5 файл ачааллаж байна...")
    prs = Presentation(src)

    existing = len(prs.slides)
    print(f"  Existing slides: {existing}")

    print("User Guide слайдуудыг нэмж байна...")
    add_ug_cover(prs)
    add_ug_overview(prs)
    add_ug_venues(prs)
    add_ug_screens(prs)
    add_ug_steps(prs)
    add_ug_features(prs)
    add_ug_tech(prs)
    add_ug_closing(prs)

    total = len(prs.slides)
    added = total - existing
    print(f"  Added slides: {added}")
    print(f"  Total slides: {total}")

    prs.save(dst)
    print(f"\n✅  Амжилттай хадгалагдлаа: {dst}")


if __name__ == "__main__":
    main()
