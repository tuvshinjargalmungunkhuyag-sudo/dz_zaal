import sys, io, shutil, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SH

# ── Color palette ──────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x1A, 0x5C)
BLUE   = RGBColor(0x1A, 0x5C, 0xBB)
YELLOW = RGBColor(0xF5, 0xC5, 0x18)
LGRAY  = RGBColor(0xF2, 0xF4, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x1E, 0x8B, 0x4C)
RED    = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xD4, 0x70, 0x0F)
TEAL   = RGBColor(0x12, 0x7A, 0x6B)
PURPLE = RGBColor(0x6C, 0x35, 0x7E)
DKGRAY = RGBColor(0x33, 0x33, 0x44)

# Tier palette
T1_BG = RGBColor(0xD6, 0xEA, 0xFF)   # client – light blue
T1_BD = BLUE
T2_BG = RGBColor(0xFF, 0xF0, 0xCC)   # server – light amber
T2_BD = RGBColor(0xC8, 0x96, 0x00)
T3_BG = RGBColor(0xD5, 0xF5, 0xE3)   # data   – light green
T3_BD = GREEN
CI_BG = RGBColor(0xED, 0xDE, 0xFF)   # ci/cd  – light lavender
CI_BD = PURPLE

def I(x): return Inches(x)
def P(x): return Pt(x)

# ── Shape helpers ───────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(1.2), rounded=False):
    kind = SH.ROUNDED_RECTANGLE if rounded else SH.RECTANGLE
    sh = slide.shapes.add_shape(kind, I(l), I(t), I(w), I(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border; sh.line.width = bw
    else:
        sh.line.fill.background()
    if rounded:
        # set corner radius via XML (smaller = rounder in pptx units: 0=sharp..50000=half)
        from lxml import etree
        prstGeom = sh._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prstGeom is not None:
            avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
            if avLst is None:
                avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
            gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
            gd.set('name', 'adj')
            gd.set('fmla', 'val 16667')
    return sh

def txb(slide, text, l, t, w, h, size=11, bold=False, color=NAVY,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = P(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def ctr(slide, text, l, t, w, h, size=11, bold=False, color=NAVY, italic=False):
    return txb(slide, text, l, t, w, h, size, bold, color, PP_ALIGN.CENTER, italic)

def icon_card(slide, emoji, title, l, t, w, h,
              bg=LGRAY, bd=BLUE, title_size=12, emoji_size=22):
    """Rounded card with big emoji + title"""
    rect(slide, l, t, w, h, fill=bg, border=bd, rounded=True)
    ctr(slide, emoji,  l, t+0.06, w, 0.45, size=emoji_size)
    ctr(slide, title,  l, t+0.52, w, h-0.55, size=title_size, bold=True, color=NAVY)

def arrow_down(slide, cx, y_top, y_bot, color=BLUE, label=''):
    """Vertical arrow + optional label"""
    from pptx.oxml.ns import qn
    from lxml import etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    pns = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    mid = (y_top + y_bot) / 2
    # draw a thin rectangle as shaft
    shaft_w = 0.06
    rect(slide, cx-shaft_w/2, y_top, shaft_w, y_bot-y_top, fill=color, border=None)
    # arrowhead triangle – use a simple narrow rect at bottom
    arw_w = 0.22
    arw_h = 0.14
    rect(slide, cx-arw_w/2, y_bot-arw_h, arw_w, arw_h, fill=color, border=None)
    if label:
        ctr(slide, label, cx-0.7, mid-0.12, 1.4, 0.24, size=8, italic=True, color=DKGRAY)

def tier_box(slide, tier_num, tier_name, tech_title, items,
             l, t, w, h, bg, bd, title_color=NAVY):
    """One tier block in the architecture diagram"""
    rect(slide, l, t, w, h, fill=bg, border=bd, bw=Pt(1.5))
    # tier label pill on left
    rect(slide, l, t+h/2-0.20, 0.36, 0.40, fill=bd, border=None, rounded=True)
    ctr(slide, str(tier_num), l, t+h/2-0.20, 0.36, 0.40, size=14, bold=True, color=WHITE)
    # tech name
    txb(slide, tech_title, l+0.44, t+0.06, w-0.52, 0.34,
        size=13, bold=True, color=title_color)
    txb(slide, tier_name, l+0.44, t+0.38, w-0.52, 0.24,
        size=9, italic=True, color=DKGRAY)
    # items row
    item_w = (w - 0.52 - 0.08*(len(items)-1)) / len(items)
    for i, (icon, lbl) in enumerate(items):
        ix = l + 0.44 + i*(item_w+0.08)
        rect(slide, ix, t+0.65, item_w, h-0.75, fill=WHITE, border=bd, bw=Pt(0.8), rounded=True)
        ctr(slide, icon, ix, t+0.68, item_w, 0.30, size=14)
        ctr(slide, lbl, ix, t+0.98, item_w, h-1.12, size=8, bold=False, color=NAVY)

def clear_content(slide, keep=7):
    """Delete all non-template shapes"""
    shapes = list(slide.shapes)
    for sh in shapes[keep:]:
        sh._element.getparent().remove(sh._element)

# ── Load v7, copy to v8 ─────────────────────────────────────────────
shutil.copy('Goviyn_Sport_Iltgel_v7.pptx', 'Goviyn_Sport_Iltgel_v8.pptx')
prs = Presentation('Goviyn_Sport_Iltgel_v8.pptx')
slides = prs.slides

# Content y-range: 0.80 → 7.10  (header 0.76, footer at 7.12)
CY = 0.85   # content top
CH = 6.20   # content height
CX = 0.30   # left margin
CW = 12.73  # content width

print("Rebuilding slides 3–7 …")

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem: visual "before vs solution" diagram
# ════════════════════════════════════════════════════════════════════
s3 = slides[2]
clear_content(s3)

# ── left column: 4 problems ──
problems = [
    ("📵", "Дижитал дэд\nбүтэц дутмаг",    "Утасны яриа /\nбиечилсэн бүртгэл"),
    ("📊", "Бодит цагийн\nмэдээлэл байхгүй", "Слотын нөхцөлийг\nурьдчилан мэдэхгүй"),
    ("⚠️",  "Захиалгын\nдавхардал",          "Гараар хөтлөх\nбүртгэл → алдаа"),
    ("📉", "Аналитик\nдутагдал",             "Тайлан, статистик\nавтоматжаагүй"),
]

prob_w, prob_h = 2.90, 1.38
prob_x = 0.25
for i, (ico, title, desc) in enumerate(problems):
    py = CY + i*(prob_h + 0.08)
    rect(s3, prob_x, py, prob_w, prob_h, fill=RGBColor(0xFF,0xEB,0xEB), border=RED, bw=Pt(1.2), rounded=True)
    ctr(s3, ico,   prob_x,       py+0.06, 0.55, 0.50, size=20)
    txb(s3, title, prob_x+0.58,  py+0.06, prob_w-0.66, 0.50, size=11, bold=True, color=RED)
    txb(s3, desc,  prob_x+0.58,  py+0.60, prob_w-0.66, 0.70, size=9, color=DKGRAY)

# ── center arrow ──
arr_x = prob_x + prob_w + 0.12
arr_y = CY + 1.6
for dy in range(3):
    ctr(s3, "→", arr_x, arr_y + dy*0.28, 0.55, 0.30, size=22, bold=True, color=BLUE)

# stat box
rect(s3, arr_x-0.02, CY+0.02, 0.60, 0.56, fill=YELLOW, border=NAVY, rounded=True)
ctr(s3, "5+", arr_x-0.02, CY+0.04, 0.60, 0.28, size=16, bold=True, color=NAVY)
ctr(s3, "заал", arr_x-0.02, CY+0.30, 0.60, 0.24, size=7, bold=False, color=NAVY)

# ── right column: solution ──
sol_x = arr_x + 0.68
sol_w = CW - sol_x + CX
sol_y = CY

# big solution box
rect(s3, sol_x, sol_y, sol_w, CH-0.05, fill=RGBColor(0xD6,0xEA,0xFF), border=BLUE, bw=Pt(1.8), rounded=True)
ctr(s3, "☁️ ДЗ-Заал", sol_x, sol_y+0.08, sol_w, 0.40, size=14, bold=True, color=NAVY)
ctr(s3, "Үүлэн орчин дахь цаг захиалгын систем",
    sol_x, sol_y+0.46, sol_w, 0.28, size=9, italic=True, color=BLUE)

solutions = [
    ("📱", "Мобайл апп\n(Android / iOS)"),
    ("⏱️", "Бодит цагийн\nслот хяналт"),
    ("🔒", "Давхардлаас\nсэргийлэх логик"),
    ("📊", "Аналитик\nтайлан"),
    ("🤖", "AI туслагч\n(Claude / Groq)"),
    ("☁️", "AWS EC2\nТокиогийн бүс"),
]
sol_card_w = (sol_w - 0.30) / 3 - 0.06
for i, (ico, lbl) in enumerate(solutions):
    col = i % 3
    row = i // 3
    cx2 = sol_x + 0.15 + col*(sol_card_w+0.06)
    cy2 = sol_y + 0.82 + row*1.52
    rect(s3, cx2, cy2, sol_card_w, 1.40, fill=WHITE, border=BLUE, bw=Pt(0.8), rounded=True)
    ctr(s3, ico, cx2, cy2+0.10, sol_card_w, 0.42, size=18)
    ctr(s3, lbl, cx2, cy2+0.55, sol_card_w, 0.78, size=9, bold=False, color=NAVY)

print("  ✓ Slide 3 rebuilt")

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Goals: compact hexagonal-style grid
# ════════════════════════════════════════════════════════════════════
s4 = slides[3]
clear_content(s4)

# main goal banner
rect(s4, 0.28, CY, CW, 0.68, fill=NAVY, border=None, rounded=True)
ctr(s4, "ЗОРИЛГО — Даланзадгадын спорт байгууллагад зориулсан үүлэн орчинтой 3 давхаргат цаг захиалгын систем бүтээх",
    0.35, CY+0.07, CW-0.14, 0.54, size=10, bold=True, color=WHITE)

goals = [
    ("🗓️", "З-1", "Цаг захиалгын\nмодуль",
     "Бодит цагийн слот шалгалт\nДавхардлаас сэргийлэх логик\nХагас/бүтэн талбай"),
    ("📱", "З-2", "Мобайл\nинтерфэйс",
     "Flutter (Android/iOS)\nMaterial Design 3\nДулаан цөлийн дизайн"),
    ("🤖", "З-3", "AI туслах\nсистем",
     "Claude Haiku 4.5 (үндсэн)\nGroq LLaMA 3.1 (нөөц)\nМонгол хэлний дэмжлэг"),
    ("☁️", "З-4", "Cloud дэд\nбүтэц",
     "AWS EC2 · Docker\nGitHub Actions CI/CD\nАвтомат деплой"),
    ("🔐", "З-5", "Аюулгүй\nнэвтрэлт",
     "Email + OTP баталгаажуулалт\nFirebase Auth\nFirestore аюулгүй дүрмүүд"),
    ("📊", "З-6", "Туршилт &\nбаталгаажуулалт",
     "10 функциональ тест\nGuerilla useability\nГүйцэтгэлийн хэмжилт"),
]

cols, rows = 3, 2
gw = (CW - 0.10) / cols - 0.10
gh = (CH - 0.82) / rows - 0.10

for i, (ico, num, title, bullets) in enumerate(goals):
    col = i % cols
    row = i // cols
    gx = CX + col * (gw + 0.10)
    gy = CY + 0.82 + row * (gh + 0.10)

    rect(s4, gx, gy, gw, gh, fill=LGRAY, border=BLUE, bw=Pt(1.2), rounded=True)
    # colored header bar
    rect(s4, gx, gy, gw, 0.52, fill=BLUE, border=None, rounded=True)
    ctr(s4, ico,   gx,       gy+0.05, 0.55, 0.42, size=18, color=WHITE)
    txb(s4, num,   gx+0.55,  gy+0.06, 0.55, 0.22, size=8, bold=True, color=YELLOW)
    txb(s4, title, gx+0.55,  gy+0.24, gw-0.62, 0.26, size=10, bold=True, color=WHITE)
    # bullets
    for bi, line in enumerate(bullets.split('\n')):
        txb(s4, "• " + line, gx+0.12, gy+0.58+bi*0.28, gw-0.20, 0.28,
            size=9, color=DKGRAY)

print("  ✓ Slide 4 rebuilt")

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Significance: 3-column visual cards
# ════════════════════════════════════════════════════════════════════
s5 = slides[4]
clear_content(s5)

sigs = [
    ("🔬", "Шинжлэх\nухааны", BLUE, T1_BG,
     [
       "3 давхаргат архитектурийн",
       "загварыг баримтжуулсан",
       "NoSQL + REST + Flutter",
       "хослолын кейс судалгаа",
       "Ухаалаг AI нэгтгэлийн",
       "практик жишээ болно",
     ]),
    ("🏙️", "Практик &\nнийгмийн", GREEN, T3_BG,
     [
       "Даланзадгадын 5+ заал",
       "нэг платформ дээр",
       "Хэрэглэгч цаг хэмнэнэ",
       "Давхардлыг арилгана",
       "Бизнесийн аналитик",
       "автоматаар бий болно",
     ]),
    ("⚙️", "Технологийн\nхувь нэмэр", ORANGE, RGBColor(0xFF,0xF0,0xCC),
     [
       "CI/CD pipeline загвар",
       "Docker + AWS EC2",
       "Email OTP баталгаажуулалт",
       "GitHub Actions автоматжуулалт",
       "Firestore аюулгүй дүрмүүд",
       "Цаашид өргөтгөх боломж",
     ]),
]

card_w = (CW - 0.20) / 3 - 0.12
for i, (ico, title, bd_color, bg_color, lines) in enumerate(sigs):
    cx = CX + i * (card_w + 0.16)
    cy = CY

    rect(s5, cx, cy, card_w, CH-0.05, fill=bg_color, border=bd_color, bw=Pt(1.8), rounded=True)
    # icon circle at top
    rect(s5, cx + card_w/2 - 0.40, cy + 0.15, 0.80, 0.80,
         fill=bd_color, border=None, rounded=True)
    ctr(s5, ico, cx + card_w/2 - 0.40, cy + 0.18, 0.80, 0.72, size=24, color=WHITE)
    ctr(s5, title, cx, cy+1.05, card_w, 0.52, size=13, bold=True, color=bd_color)

    for li, line in enumerate(lines):
        txb(s5, "✓ " + line, cx+0.18, cy+1.68+li*0.38, card_w-0.26, 0.36,
            size=9.5, color=DKGRAY)

print("  ✓ Slide 5 rebuilt")

# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — System Architecture: full 3-tier + deployment diagram
# ════════════════════════════════════════════════════════════════════
s6 = slides[5]
clear_content(s6)

# ── Left: 3-tier stack ───────────────────────────────────────────────
tier_x  = CX
tier_w  = 7.80
tier_h  = 1.55
tier_gap = 0.42   # space between tiers (for arrow)
t1y = CY + 0.02
t2y = t1y + tier_h + tier_gap
t3y = t2y + tier_h + tier_gap

# Tier 1 – Flutter Client
tier_box(s6, 1, "Presentation Tier",
         "Flutter  (Dart 3.x)",
         [("📱","Android"), ("🍎","iOS"), ("🎨","Material\nDesign 3"), ("⚡","Provider\nState")],
         tier_x, t1y, tier_w, tier_h, T1_BG, T1_BD)

# Arrow 1→2
arrow_down(s6, tier_x + tier_w/2, t1y+tier_h+0.04, t2y-0.04, BLUE, "HTTPS / REST API")

# Tier 2 – Node.js/Express
tier_box(s6, 2, "Application Tier  —  AWS EC2 ap-northeast-1 (Токио)",
         "Node.js 18  /  Express.js",
         [("🔗","REST\nAPI"), ("📧","Nodemailer\nSMTP"), ("⏰","node-cron\nCron job"), ("🐳","Docker\nContainer")],
         tier_x, t2y, tier_w, tier_h, T2_BG, T2_BD)

# Arrow 2→3
arrow_down(s6, tier_x + tier_w/2, t2y+tier_h+0.04, t3y-0.04, GREEN, "Firebase Admin SDK")

# Tier 3 – Firestore
tier_box(s6, 3, "Data Tier",
         "Cloud Firestore  (NoSQL)",
         [("📂","bookings"), ("👤","users"), ("🔒","fixed_\nbookings"), ("✉️","email_\nverif.")],
         tier_x, t3y, tier_w, tier_h, T3_BG, T3_BD)

# ── Right: supporting services column ───────────────────────────────
rx   = tier_x + tier_w + 0.22
rw   = CW - tier_w - 0.28

# Auth card
rect(s6, rx, t1y, rw, 1.10, fill=RGBColor(0xEE,0xEE,0xFF), border=PURPLE, rounded=True)
ctr(s6, "🔐",        rx,       t1y+0.06, rw, 0.38, size=20)
ctr(s6, "Firebase Auth",       rx, t1y+0.44, rw, 0.26, size=10, bold=True, color=PURPLE)
ctr(s6, "Email OTP · ID Token",rx, t1y+0.70, rw, 0.24, size=8, color=DKGRAY)
ctr(s6, "+ Firestore Rules",   rx, t1y+0.90, rw, 0.20, size=7.5, italic=True, color=PURPLE)

# AI card
ai_y = t1y + 1.18
rect(s6, rx, ai_y, rw, 1.22, fill=RGBColor(0xE8,0xF5,0xE9), border=GREEN, rounded=True)
ctr(s6, "🤖",              rx, ai_y+0.06, rw, 0.38, size=20)
ctr(s6, "AI Туслах",       rx, ai_y+0.44, rw, 0.26, size=10, bold=True, color=GREEN)
ctr(s6, "Claude Haiku 4.5",rx, ai_y+0.70, rw, 0.24, size=8.5, color=DKGRAY)
ctr(s6, "Groq LLaMA 3.1 (нөөц)",rx, ai_y+0.93, rw, 0.22, size=8, color=DKGRAY)

# CI/CD pipeline
ci_y = t2y + 0.08
rect(s6, rx, ci_y, rw, 2.55, fill=RGBColor(0xF3,0xE5,0xFF), border=PURPLE, bw=Pt(1.5), rounded=True)
ctr(s6, "🔄 CI/CD Pipeline", rx, ci_y+0.08, rw, 0.30, size=10, bold=True, color=PURPLE)

steps = [
    ("⬆️", "git push\nmain"),
    ("⚙️", "GitHub\nActions"),
    ("🐳", "Docker\nBuild"),
    ("📦", "AWS\nECR Push"),
    ("🚀", "EC2\nDeploy"),
]
sw = (rw - 0.12) / len(steps)
for si, (sico, slbl) in enumerate(steps):
    sx = rx + 0.06 + si*sw
    rect(s6, sx, ci_y+0.45, sw-0.04, 0.95, fill=WHITE, border=PURPLE, bw=Pt(0.6), rounded=True)
    ctr(s6, sico, sx, ci_y+0.49, sw-0.04, 0.36, size=14)
    ctr(s6, slbl, sx, ci_y+0.84, sw-0.04, 0.52, size=7, color=NAVY)
    if si < len(steps)-1:
        ctr(s6, "→", sx+sw-0.04, ci_y+0.62, 0.10, 0.28, size=9, color=PURPLE)

ctr(s6, "Push → Build → Registry → Deploy  (3–4 мин)",
    rx, ci_y+1.50, rw, 0.24, size=7.5, italic=True, color=PURPLE)
ctr(s6, "AWS ECR  ap-northeast-1 (Токио)",
    rx, ci_y+1.72, rw, 0.24, size=7.5, color=DKGRAY)
ctr(s6, "node:18-alpine  ·  dz-zaal-backend",
    rx, ci_y+1.94, rw, 0.24, size=7.5, color=DKGRAY)

# Performance badges
perf_y = t3y + 0.14
perf_items = [
    ("📱","2.3 с","App boot"),
    ("🗄️","0.6 с","Firestore"),
    ("⚡","0.8 с","API"),
    ("🤖","1.4 с","Claude"),
]
pw = (rw - 0.10) / len(perf_items) - 0.04
for pi, (pico, pval, plbl) in enumerate(perf_items):
    px2 = rx + 0.05 + pi*(pw+0.04)
    rect(s6, px2, perf_y, pw, 0.95, fill=LGRAY, border=NAVY, bw=Pt(0.8), rounded=True)
    ctr(s6, pico, px2, perf_y+0.04, pw, 0.28, size=13)
    ctr(s6, pval, px2, perf_y+0.32, pw, 0.32, size=13, bold=True, color=NAVY)
    ctr(s6, plbl, px2, perf_y+0.65, pw, 0.24, size=7, color=DKGRAY)

print("  ✓ Slide 6 rebuilt")

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Tech Stack: visual 3-layer + decision matrix
# ════════════════════════════════════════════════════════════════════
s7 = slides[6]
clear_content(s7)

# Layer bars (left 60%)
LW = 8.00
ly = CY

layers = [
    ("PRESENTATION\nTIER",  BLUE, T1_BG,
     [("Flutter 3.x","🦋"), ("Dart 3","🎯"), ("Material\nDesign 3","🎨"),
      ("Provider","🔄"), ("Google\nFonts","✏️")]),
    ("APPLICATION\nTIER",   T2_BD, T2_BG,
     [("Node.js 18","🟢"), ("Express.js","⚡"), ("Nodemailer","📧"),
      ("node-cron","⏰"), ("Groq API","🤖")]),
    ("DATA &\nINFRASTRUCTURE", GREEN, T3_BG,
     [("Firestore","🔥"), ("Firebase\nAuth","🔐"), ("AWS EC2\nTokio","☁️"),
      ("Docker","🐳"), ("GitHub\nActions","🔄")]),
]

lh = (CH - 0.10) / 3 - 0.12
for li, (lname, bd, bg, techs) in enumerate(layers):
    lx = CX
    cur_ly = ly + li*(lh+0.12)

    rect(s7, lx, cur_ly, LW, lh, fill=bg, border=bd, bw=Pt(1.8))
    # label on the very left
    rect(s7, lx, cur_ly, 1.10, lh, fill=bd, border=None)
    ctr(s7, lname, lx, cur_ly + lh/2-0.35, 1.10, 0.70, size=9, bold=True, color=WHITE)

    # tech pills
    pill_w = (LW - 1.20 - 0.06*len(techs)) / len(techs)
    for ti, (tname, ticon) in enumerate(techs):
        tx2 = lx + 1.20 + ti*(pill_w+0.06)
        rect(s7, tx2, cur_ly+0.14, pill_w, lh-0.28,
             fill=WHITE, border=bd, bw=Pt(0.8), rounded=True)
        ctr(s7, ticon, tx2, cur_ly+0.20, pill_w, 0.36, size=16)
        ctr(s7, tname, tx2, cur_ly+0.54, pill_w, lh-0.72, size=8, bold=False, color=NAVY)

# Right panel: decision highlights
rx2 = CX + LW + 0.20
rw2 = CW - LW - 0.16
rect(s7, rx2, ly, rw2, CH-0.05, fill=NAVY, border=None, rounded=True)
ctr(s7, "Технологийн\nсонголтын\nшалтгаан",
    rx2, ly+0.14, rw2, 0.70, size=11, bold=True, color=YELLOW)

reasons = [
    ("🦋", "Flutter", "Нэг кодоор iOS/Android — хөгжүүлэлтийн зардал хоёр дахин хэмнэгдэнэ"),
    ("🔥", "Firebase", "Бодит цагийн синк + Auth нэгтгэл — бэкэнд хялбарчлагдана"),
    ("🟢", "Node.js", "JS full-stack — ижил хэл, хурдан prototyping"),
    ("☁️", "AWS EC2", "Монголд хамгийн ойр Токио бүс — latency бага"),
    ("🐳", "Docker", "Орчны тохирголт тогтмол — dev/prod ялгаа арилна"),
    ("🔄", "Actions", "Push-to-deploy — manual орхигдоно"),
]
for ri, (rico, rtech, rdesc) in enumerate(reasons):
    ry2 = ly + 0.94 + ri*0.85
    rect(s7, rx2+0.12, ry2, rw2-0.24, 0.78, fill=RGBColor(0x14,0x2A,0x6E), border=None, rounded=True)
    ctr(s7, rico, rx2+0.12, ry2+0.12, 0.50, 0.50, size=14, color=WHITE)
    txb(s7, rtech, rx2+0.65, ry2+0.08, 1.10, 0.28, size=9, bold=True, color=YELLOW)
    txb(s7, rdesc, rx2+0.65, ry2+0.35, rw2-0.82, 0.38, size=7.5, color=LGRAY)

print("  ✓ Slide 7 rebuilt")

# ════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════
prs.save('Goviyn_Sport_Iltgel_v8.pptx')
print("\n✅  Saved: Goviyn_Sport_Iltgel_v8.pptx  (slides 3–7 fully redesigned)")
print(f"   Total slides: {len(prs.slides)}")
