import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

prs = Presentation('Goviyn_Sport_Iltgel_v7.pptx')
for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t[:80])
    if texts:
        print(f'\nSLIDE {i+1}: {texts[0][:60]}')
        for t in texts[1:4]:
            print(f'  - {t[:70]}')
        if len(texts) > 4:
            print(f'  ... ({len(texts)} shapes total)')
