"""
v3_editable_v5.pptx -> v3_editable_v6.pptx
Email Баталгаажуулалтын Sequence Diagram-ыг шинэ slide болгож,
slide 11-ийн ард байрлуулна. Бүгд editable native PowerPoint shape.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v5.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v6.pptx"

# --- Өнгө ---------------------------------------------------------------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x1B, 0x1F, 0x3A)
INK_SOFT   = RGBColor(0x4A, 0x52, 0x6E)
PAPER      = RGBColor(0xFA, 0xFA, 0xFA)
PAPER2     = RGBColor(0xEE, 0xF1, 0xF8)
ACCENT     = RGBColor(0x1F, 0x3A, 0x68)
ACCENT2    = RGBColor(0xC2, 0x6A, 0x12)
ACCENT3    = RGBColor(0x1E, 0x88, 0x4F)
RETURN     = RGBColor(0x1E, 0x88, 0x4F)   # буцах сум — ногоон

# Lifeline толгойн өнгөнүүд (тус бүрийг ялгаж харуулах)
LIFE_COLORS = [
    RGBColor(0x1F, 0x3A, 0x68),   # Хэрэглэгч — хөх
    RGBColor(0x6A, 0x1B, 0x9A),   # Flutter — нил ягаан
    RGBColor(0xE6, 0x7E, 0x22),   # Firebase — улбар шар
    RGBColor(0x16, 0xA0, 0x85),   # Node.js — ногоон-цэнхэр
    RGBColor(0xC0, 0x39, 0x2B),   # Gmail — улаан
]
LIFE_BG = [
    RGBColor(0xE3, 0xEA, 0xF6),
    RGBColor(0xF1, 0xE5, 0xF7),
    RGBColor(0xFD, 0xEB, 0xD8),
    RGBColor(0xD7, 0xF2, 0xEC),
    RGBColor(0xFD, 0xEC, 0xEA),
]


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0, dash=None):
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)
    if dash is not None:
        ln = line._get_or_add_ln()
        prst = ln.find(qn('a:prstDash'))
        if prst is None:
            prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)


def text_in_shape(shape, text, *, size=11, bold=False, color=INK,
                  align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(40000); tf.margin_right = Emu(40000)
    tf.margin_top = Emu(20000);  tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = align
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            run = p.add_run()
        else:
            p2 = tf.add_paragraph(); p2.alignment = align
            run = p2.add_run()
        run.text = line
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = "Calibri"


def add_textbox(slide, x, y, w, h, text, *, size=11, bold=False,
                italic=False, color=INK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color; run.font.name = "Calibri"
    return tb


def add_arrow(slide, x1, y1, x2, y2, *,
              color=INK, width_pt=1.25, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line = conn.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    ln = line._get_or_add_ln()
    if dash is not None:
        prst = ln.find(qn('a:prstDash'))
        if prst is None:
            prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return conn


# --- Презентац ачаалах --------------------------------------------------
prs = Presentation(SRC)

# Шинэ slide (Blank layout) нэмэх
blank_layout = prs.slide_layouts[6]   # Blank layout
slide = prs.slides.add_slide(blank_layout)

SW = prs.slide_width
SH = prs.slide_height

# --- Title -------------------------------------------------------------
title_tb = slide.shapes.add_textbox(Emu(450_000), Emu(180_000),
                                    SW - Emu(2_500_000), Emu(450_000))
tf = title_tb.text_frame
tf.margin_left = Emu(0); tf.margin_right = Emu(0)
tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r1 = p.add_run(); r1.text = "Бүлэг 2.2:  "
r1.font.size = Pt(14); r1.font.bold = True
r1.font.color.rgb = ACCENT2; r1.font.name = "Calibri"
r2 = p.add_run(); r2.text = "Email Баталгаажуулалт — Sequence Diagram"
r2.font.size = Pt(20); r2.font.bold = True
r2.font.color.rgb = ACCENT; r2.font.name = "Calibri"

# Баруун дээд буланд "UML Sequence" badge
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               SW - Emu(2_100_000), Emu(200_000),
                               Emu(1_700_000), Emu(380_000))
badge.adjustments[0] = 0.45
set_fill(badge, ACCENT)
set_line(badge, ACCENT, 0.5)
badge.shadow.inherit = False
text_in_shape(badge, "UML SEQUENCE", size=10, bold=True, color=WHITE)

# --- Lifeline толгойнууд -----------------------------------------------
LIFELINES = [
    "Хэрэглэгч",
    "Flutter App",
    "Firebase\nAuth",
    "Node.js\nEC2",
    "Gmail\nSMTP",
]

HEAD_W = Emu(2_000_000)
HEAD_H = Emu(750_000)
HEAD_Y = Emu(900_000)

# Тэнцүү зайтай 5 толгой байрлуулах
margin = Emu(450_000)
total_w = SW - 2 * margin
gap = (total_w - 5 * HEAD_W) // 4    # 4 зай

life_x = []   # тус бүрийн center X
for i in range(5):
    x = margin + i * (HEAD_W + gap)
    cx = x + HEAD_W // 2
    life_x.append(cx)

    head = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, HEAD_Y, HEAD_W, HEAD_H)
    head.adjustments[0] = 0.18
    set_fill(head, LIFE_BG[i])
    set_line(head, LIFE_COLORS[i], 1.5)
    head.shadow.inherit = False
    text_in_shape(head, LIFELINES[i],
                  size=14, bold=True, color=LIFE_COLORS[i])

# --- Lifeline босоо dashed шугам ---------------------------------------
LIFE_TOP = HEAD_Y + HEAD_H
LIFE_BOTTOM = Emu(6_400_000)

for cx in life_x:
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      cx, LIFE_TOP, cx, LIFE_BOTTOM)
    ln = conn.line._get_or_add_ln()
    conn.line.color.rgb = INK_SOFT
    conn.line.width = Pt(0.75)
    prst = ln.find(qn('a:prstDash'))
    if prst is None:
        prst = etree.SubElement(ln, qn('a:prstDash'))
    prst.set('val', 'dash')

# --- Message-ууд (8 ширхэг) --------------------------------------------
# (from_idx, to_idx, label, is_return)
MESSAGES = [
    (0, 1, "1. email, нууц үг",   False),
    (1, 2, "2. createUser",        False),
    (2, 1, "3. UID",                True),
    (1, 3, "4. POST /send-code",   False),
    (3, 4, "5. 6-digit code",      False),
    (4, 3, "6. OK",                 True),
    (3, 1, "7. 200 OK",             True),
    (1, 0, "8. Код оруул",          True),
]

# Эхний message Y болон зай
MSG_FIRST_Y = LIFE_TOP + Emu(500_000)
MSG_GAP_Y   = Emu(540_000)

for i, (frm, to, label, is_return) in enumerate(MESSAGES):
    y = MSG_FIRST_Y + i * MSG_GAP_Y
    x1 = life_x[frm]
    x2 = life_x[to]
    color = RETURN if is_return else INK
    dash = "dash" if is_return else None
    width = 1.5 if not is_return else 1.25

    add_arrow(slide, x1, y, x2, y, color=color, width_pt=width, dash=dash)

    # Label — сумны ДЭЭР, төв дээр
    label_w = Emu(2_400_000)
    label_x = (x1 + x2) // 2 - label_w // 2
    add_textbox(slide,
                label_x, y - Emu(330_000),
                label_w, Emu(280_000),
                label,
                size=11, bold=True,
                color=color if is_return else INK,
                italic=is_return,
                align=PP_ALIGN.CENTER)

# --- Доод тайлбарын зурвас ---------------------------------------------
DESC_X = Emu(450_000)
DESC_Y = Emu(6_500_000)
DESC_W = SW - 2 * DESC_X
DESC_H = Emu(280_000)

# Сорилт тайлбар
add_textbox(slide,
            DESC_X, DESC_Y, DESC_W, DESC_H,
            "──── солих хүсэлт      ─ ─ ─ хариу буцах",
            size=10, italic=True, color=INK_SOFT,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# --- Шинэ slide-ыг slide 12-ын байрлалд шилжүүлэх (slide 11-ийн дараа) ----
# python-pptx-д `slides` нь sldIdLst-ийг засаж байрлал өөрчилнө
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
new_slide = slides[-1]                  # хамгийн сүүлд нэмэгдсэн
xml_slides.remove(new_slide)
xml_slides.insert(11, new_slide)        # 12-р slide болгох (0-based 11)

# --- Хадгалах -----------------------------------------------------------
prs.save(DST)
print("Saved:", DST)
print(f"Шинэ Sequence Diagram = Slide 12")
