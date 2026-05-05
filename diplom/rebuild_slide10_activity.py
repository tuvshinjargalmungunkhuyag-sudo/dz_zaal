"""
v3_editable_v3.pptx -> v3_editable_v4.pptx
Activity Diagram-ыг native PowerPoint shape-аар зурна.
LaTeX presentation_v3.tex (line 512) дотор байгаа диаграммыг сэргээв.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v3.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v4.pptx"

# --- Өнгө ----------------------------------------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x1B, 0x1F, 0x3A)
INK_SOFT    = RGBColor(0x4A, 0x52, 0x6E)
ACCENT      = RGBColor(0x1F, 0x3A, 0x68)   # хөх — header
ACCENT2     = RGBColor(0xC2, 0x6A, 0x12)   # улбар шар — error/loop
ACCENT3     = RGBColor(0x1E, 0x88, 0x4F)   # ногоон — success
PAPER       = RGBColor(0xFA, 0xFA, 0xFA)
PAPER2      = RGBColor(0xEE, 0xF1, 0xF8)   # шийдвэр / онцгой
ERROR_FILL  = RGBColor(0xFD, 0xEC, 0xEA)
ERROR_BDR   = RGBColor(0xC0, 0x39, 0x2B)
DECISION_FILL = RGBColor(0xFF, 0xF8, 0xE1)
DECISION_BDR  = RGBColor(0xB7, 0x8A, 0x12)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def no_fill(shape):
    shape.fill.background()


def add_text_in_shape(shape, text, *, size=11, bold=False,
                      color=INK, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(40000)
    tf.margin_right = Emu(40000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_label(slide, x, y, w, h, text, *, size=10,
              color=INK_SOFT, bold=False, italic=True,
              align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_action_box(slide, cx, cy, w, h, text, *,
                   fill=PAPER, border=ACCENT, size=12, bold=False):
    """Үйл ажиллагааны bordered rounded box."""
    x, y = cx - w // 2, cy - h // 2
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.18
    set_fill(box, fill)
    set_line(box, border, 1.25)
    box.shadow.inherit = False
    add_text_in_shape(box, text, size=size, bold=bold, color=INK)
    return box


def add_decision(slide, cx, cy, w, h, text):
    """Шийдвэрийн ромб (diamond)."""
    x, y = cx - w // 2, cy - h // 2
    box = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, w, h)
    set_fill(box, DECISION_FILL)
    set_line(box, DECISION_BDR, 1.25)
    box.shadow.inherit = False
    add_text_in_shape(box, text, size=11, bold=True, color=INK)
    return box


def add_start(slide, cx, cy, d):
    """Start node — дүүрэн хар дугуй."""
    x, y = cx - d // 2, cy - d // 2
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    set_fill(c, INK)
    set_line(c, INK, 0.5)
    return c


def add_end(slide, cx, cy, d):
    """End node — давхар дугуй (outer + inner)."""
    x, y = cx - d // 2, cy - d // 2
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    no_fill(outer)
    set_line(outer, INK, 1.5)
    inner_d = int(d * 0.55)
    ix, iy = cx - inner_d // 2, cy - inner_d // 2
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, ix, iy, inner_d, inner_d)
    set_fill(inner, INK)
    set_line(inner, INK, 0.5)
    return outer


def add_arrow(slide, x1, y1, x2, y2, *, color=INK, width_pt=1.0):
    """Шулуун сум."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line = conn.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    # Сумны үзүүр нэмэх (XML дээр)
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return conn


def add_elbow(slide, points, *, color=INK, width_pt=1.0):
    """Олон цэгийг дамжсан elbow connector — сумны үзүүртэй."""
    # Сегмент тус бүрт connector. Зөвхөн сүүлчийнхэд сумны үзүүр.
    conns = []
    last = len(points) - 1
    for i in range(last):
        x1, y1 = points[i]
        x2, y2 = points[i + 1]
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        conn.line.color.rgb = color
        conn.line.width = Pt(width_pt)
        if i == last - 1:
            from pptx.oxml.ns import qn
            from lxml import etree
            ln = conn.line._get_or_add_ln()
            tail = ln.find(qn('a:tailEnd'))
            if tail is None:
                tail = etree.SubElement(ln, qn('a:tailEnd'))
            tail.set('type', 'triangle')
            tail.set('w', 'med')
            tail.set('h', 'med')
        conns.append(conn)
    return conns


# --- Презентац ачаалах ---------------------------------------------------
prs = Presentation(SRC)
slide = prs.slides[9]   # 10-р slide (empty)

# Хуучин placeholder-уудыг устгах
to_delete = list(slide.shapes)
for sh in to_delete:
    remove_shape(sh)

SW = prs.slide_width      # 12 192 000
SH = prs.slide_height     # 6 858 000

# --- Title ---------------------------------------------------------------
title_tb = slide.shapes.add_textbox(Emu(500_000), Emu(250_000),
                                    SW - Emu(1_000_000), Emu(700_000))
tf = title_tb.text_frame
tf.margin_left = Emu(0); tf.margin_right = Emu(0)
tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r1 = p.add_run(); r1.text = "Бүлэг 2.1: "
r1.font.size = Pt(20); r1.font.bold = True
r1.font.color.rgb = ACCENT2; r1.font.name = "Calibri"
r2 = p.add_run(); r2.text = "Цаг захиалах Activity Diagram"
r2.font.size = Pt(28); r2.font.bold = True
r2.font.color.rgb = ACCENT; r2.font.name = "Calibri"

# Хажуугийн тайлбар badge — баруун дээд буланд
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               SW - Emu(2_300_000), Emu(350_000),
                               Emu(1_900_000), Emu(450_000))
badge.adjustments[0] = 0.30
set_fill(badge, PAPER2)
set_line(badge, ACCENT, 1.0)
badge.shadow.inherit = False
add_text_in_shape(badge, "UML Activity", size=11, bold=True, color=ACCENT)

# --- Нийт layout — гол урсгал төв баганад -------------------------------
CENTER_X = Emu(4_700_000)            # гол баганын төвлөрсөн X
BOX_W    = Emu(3_400_000)
BOX_H    = Emu(550_000)
GAP_Y    = Emu(180_000)              # box хооронд (сум багтаах)
START_Y  = Emu(1_200_000)            # эхний (start node) Y
NODE_D   = Emu(280_000)              # start/end дугуйн диаметр

# --- Y координатууд ------------------------------------------------------
y_start  = START_Y
y_login  = y_start + NODE_D // 2 + GAP_Y + BOX_H // 2
y_date   = y_login + BOX_H + GAP_Y
y_slot   = y_date  + BOX_H + GAP_Y
y_avail  = y_slot  + BOX_H + GAP_Y + Emu(80_000)    # diamond-д орон зай нэмэх
y_create = y_avail + Emu(800_000)
y_code   = y_create + BOX_H + GAP_Y
y_notif  = y_code   + BOX_H + GAP_Y
y_done   = y_notif  + BOX_H + GAP_Y
y_end    = y_done   + BOX_H // 2 + GAP_Y + NODE_D // 2

# --- Үндсэн node-ууд (төв багана) ---------------------------------------
add_start(slide, CENTER_X, y_start, NODE_D)
b_login  = add_action_box(slide, CENTER_X, y_login, BOX_W, BOX_H,
                          "Email/нууц үгээр нэвтрэх")
b_date   = add_action_box(slide, CENTER_X, y_date,  BOX_W, BOX_H,
                          "Огноо сонгох")
b_slot   = add_action_box(slide, CENTER_X, y_slot,  BOX_W, BOX_H,
                          "Цагийн үе сонгох")
DIAM_W   = Emu(1_900_000)
DIAM_H   = Emu(900_000)
b_avail  = add_decision(slide, CENTER_X, y_avail, DIAM_W, DIAM_H,
                        "Цаг боломжтой?")
b_create = add_action_box(slide, CENTER_X, y_create, BOX_W, BOX_H,
                          "Захиалга үүсгэх (Firestore)", border=ACCENT3)
b_code   = add_action_box(slide, CENTER_X, y_code,  BOX_W, BOX_H,
                          "DBK-xxxxx код")
b_notif  = add_action_box(slide, CENTER_X, y_notif, BOX_W, BOX_H,
                          "1 цагийн өмнөх мэдэгдэл")
b_done   = add_action_box(slide, CENTER_X, y_done,  BOX_W, BOX_H,
                          "Баталгаажуулалт", fill=PAPER2,
                          border=ACCENT3, bold=True)
add_end(slide, CENTER_X, y_end, NODE_D)

# --- Алдааны салбар (баруун талд) ---------------------------------------
ERR_X    = CENTER_X + Emu(3_400_000)
ERR_W    = Emu(2_400_000)
ERR_H    = Emu(550_000)
b_err = add_action_box(slide, ERR_X, y_avail, ERR_W, ERR_H,
                       "Алдаа — дахин сонгох",
                       fill=ERROR_FILL, border=ERROR_BDR, bold=True)

# --- Сумнууд: гол урсгал ------------------------------------------------
# helper — top/bottom edge координатыг олох
def top(box, y_center, h=BOX_H):
    return (CENTER_X, y_center - h // 2)
def bot(box, y_center, h=BOX_H):
    return (CENTER_X, y_center + h // 2)

# start -> login
add_arrow(slide, CENTER_X, y_start + NODE_D // 2, CENTER_X, y_login - BOX_H // 2)
# login -> date
add_arrow(slide, CENTER_X, y_login + BOX_H // 2, CENTER_X, y_date - BOX_H // 2)
# date -> slot
add_arrow(slide, CENTER_X, y_date + BOX_H // 2, CENTER_X, y_slot - BOX_H // 2)
# slot -> avail (diamond top)
add_arrow(slide, CENTER_X, y_slot + BOX_H // 2, CENTER_X, y_avail - DIAM_H // 2)
# avail (diamond bottom) -> create
add_arrow(slide, CENTER_X, y_avail + DIAM_H // 2, CENTER_X, y_create - BOX_H // 2)
# Тийм label
add_label(slide,
          CENTER_X + Emu(60_000),
          (y_avail + DIAM_H // 2 + y_create - BOX_H // 2) // 2 - Emu(150_000),
          Emu(700_000), Emu(300_000),
          "Тийм", size=10, color=ACCENT3, bold=True, italic=True,
          align=PP_ALIGN.LEFT)

# create -> code
add_arrow(slide, CENTER_X, y_create + BOX_H // 2, CENTER_X, y_code - BOX_H // 2)
# code -> notif
add_arrow(slide, CENTER_X, y_code + BOX_H // 2, CENTER_X, y_notif - BOX_H // 2)
# notif -> done
add_arrow(slide, CENTER_X, y_notif + BOX_H // 2, CENTER_X, y_done - BOX_H // 2)
# done -> end node
add_arrow(slide, CENTER_X, y_done + BOX_H // 2, CENTER_X, y_end - NODE_D // 2)

# --- avail -> err (баруун тийш) -----------------------------------------
add_arrow(slide,
          CENTER_X + DIAM_W // 2, y_avail,
          ERR_X - ERR_W // 2,    y_avail,
          color=ERROR_BDR, width_pt=1.25)
# Үгүй label
add_label(slide,
          (CENTER_X + DIAM_W // 2 + ERR_X - ERR_W // 2) // 2 - Emu(200_000),
          y_avail - Emu(380_000),
          Emu(700_000), Emu(280_000),
          "Үгүй", size=10, color=ERROR_BDR, bold=True, italic=True)

# --- err буцах урсгал: err.top -> дээш -> зүүн тийш -> slot.right -------
err_top = (ERR_X, y_avail - ERR_H // 2)
turn1   = (ERR_X, y_slot)                      # дээш гарах
turn2   = (CENTER_X + BOX_W // 2 + Emu(80_000), y_slot)
slot_right = (CENTER_X + BOX_W // 2, y_slot)
add_elbow(slide, [err_top, turn1, turn2, slot_right],
          color=ERROR_BDR, width_pt=1.25)
# "loop" tiny label
add_label(slide,
          ERR_X - Emu(450_000), y_slot - Emu(360_000),
          Emu(900_000), Emu(280_000),
          "дахин оролдох", size=9, color=INK_SOFT, italic=True,
          align=PP_ALIGN.CENTER)

# --- Зүүн талд жижиг тайлбар (legend) -----------------------------------
LEG_X = Emu(400_000)
LEG_Y = Emu(1_400_000)
LEG_W = Emu(1_700_000)
ROW_H = Emu(380_000)

legend_items = [
    ("●",  "Эхлэл",            INK),
    ("◯",  "Төгсгөл",          INK),
    ("◇",  "Шийдвэр",          DECISION_BDR),
    ("▭",  "Үйл ажиллагаа",    ACCENT),
    ("▭",  "Алдаа",            ERROR_BDR),
]

# Legend толгой
add_label(slide, LEG_X, LEG_Y - Emu(380_000),
          LEG_W, Emu(320_000),
          "ТАЙЛБАР", size=10, color=ACCENT, bold=True, italic=False,
          align=PP_ALIGN.LEFT)

for i, (sym, label, color) in enumerate(legend_items):
    add_label(slide,
              LEG_X, LEG_Y + i * ROW_H,
              Emu(280_000), ROW_H,
              sym, size=14, color=color, bold=True, italic=False,
              align=PP_ALIGN.LEFT)
    add_label(slide,
              LEG_X + Emu(320_000), LEG_Y + i * ROW_H,
              LEG_W - Emu(320_000), ROW_H,
              label, size=10, color=INK, italic=False,
              align=PP_ALIGN.LEFT)

# --- Хадгалах ------------------------------------------------------------
prs.save(DST)
print("Saved:", DST)
