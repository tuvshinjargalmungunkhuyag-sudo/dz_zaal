"""
v3_editable_v12.pptx -> v3_editable_v13.pptx
Slide 21 "Хэрэгжилтийн Үр Дүн" — 8 хэрэгжүүлсэн зүйлсийн checklist.
Дээр төгссөн статистик banner, доор 2x4 grid карт.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v12.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v13.pptx"

# --- Өнгө ---------------------------------------------------------------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x1B, 0x1F, 0x3A)
INK_SOFT   = RGBColor(0x4A, 0x52, 0x6E)
PAPER      = RGBColor(0xFA, 0xFA, 0xFA)
PAPER2     = RGBColor(0xEE, 0xF1, 0xF8)
ACCENT     = RGBColor(0x1F, 0x3A, 0x68)
ACCENT2    = RGBColor(0xC2, 0x6A, 0x12)
SUCCESS    = RGBColor(0x1E, 0x88, 0x4F)
SUCCESS_FILL = RGBColor(0xE7, 0xF5, 0xEC)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def text_in_shape(shape, text, *, size=11, bold=False, color=INK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  italic=False, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(40_000); tf.margin_right = Emu(40_000)
    tf.margin_top = Emu(20_000);  tf.margin_bottom = Emu(20_000)
    tf.vertical_anchor = anchor
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font


def add_textbox(slide, x, y, w, h, text, *, size=11, bold=False,
                italic=False, color=INK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    return tb


def add_result_card(slide, x, y, w, h, title, subtitle, *, color=ACCENT):
    """Хэрэгжилтийн нэг item-ийн карт. Зүүн талд checkmark, баруун талд DONE."""
    # Үндсэн карт
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.10
    set_fill(card, WHITE)
    set_line(card, RGBColor(0xCF, 0xD8, 0xEC), 1.0)
    card.shadow.inherit = False

    # Зүүн талд checkmark circle
    chk_d = Emu(560_000)
    chk_x = x + Emu(150_000)
    chk_y = y + (h - chk_d) // 2
    chk = slide.shapes.add_shape(MSO_SHAPE.OVAL, chk_x, chk_y, chk_d, chk_d)
    set_fill(chk, SUCCESS)
    set_line(chk, SUCCESS, 0.25)
    text_in_shape(chk, "✓", size=22, bold=True, color=WHITE)

    # DONE badge — баруун талд
    DONE_W = Emu(750_000)
    DONE_H = Emu(360_000)
    done_x = x + w - DONE_W - Emu(150_000)
    done_y = y + (h - DONE_H) // 2
    done = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  done_x, done_y, DONE_W, DONE_H)
    done.adjustments[0] = 0.45
    set_fill(done, SUCCESS_FILL)
    set_line(done, SUCCESS, 1.0)
    text_in_shape(done, "DONE", size=10, bold=True, color=SUCCESS)

    # Title (дээр)
    text_x = chk_x + chk_d + Emu(220_000)
    text_w = done_x - text_x - Emu(150_000)
    add_textbox(slide,
                text_x, y + Emu(180_000),
                text_w, Emu(380_000),
                title, size=12, bold=True, color=color,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    # Subtitle (доор)
    add_textbox(slide,
                text_x, y + Emu(580_000),
                text_w, Emu(380_000),
                subtitle, size=10, color=INK_SOFT,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    return card


def add_stat(slide, cx, cy, w, h, big_value, label):
    """Дугаар + label-тай KPI tile."""
    x, y = cx - w // 2, cy - h // 2
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.10
    set_fill(box, SUCCESS)
    set_line(box, SUCCESS, 0.5)
    box.shadow.inherit = False
    # Том value
    add_textbox(slide,
                x, y + Emu(80_000), w, h // 2,
                big_value, size=28, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
    # Label
    add_textbox(slide,
                x, y + h // 2 + Emu(40_000), w, h // 2 - Emu(80_000),
                label, size=10, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    return box


# --- Презентац ачаалах --------------------------------------------------
prs = Presentation(SRC)
slide = prs.slides[20]   # 21-р slide

# Picture болон хуучин title-ыг устгах (frame үлдээх)
KEEP = {"Rectangle 1", "Rectangle 2", "TextBox 3",
        "Rectangle 4", "SlideNum"}
to_delete = []
for sh in slide.shapes:
    if sh.name not in KEEP:
        to_delete.append(sh)
for sh in to_delete:
    remove_shape(sh)

SW = prs.slide_width
SH = prs.slide_height

# --- Title ---------------------------------------------------------------
title_tb = slide.shapes.add_textbox(Emu(420_000), Emu(120_000),
                                    SW - Emu(2_500_000), Emu(450_000))
tf = title_tb.text_frame
tf.margin_left = Emu(0); tf.margin_right = Emu(0)
tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r1 = p.add_run(); r1.text = "Бүлэг 3.2:  "
r1.font.size = Pt(14); r1.font.bold = True
r1.font.color.rgb = ACCENT2; r1.font.name = "Calibri"
r2 = p.add_run(); r2.text = "Хэрэгжилтийн үр дүн"
r2.font.size = Pt(20); r2.font.bold = True
r2.font.color.rgb = ACCENT; r2.font.name = "Calibri"

# Баруун дээд "100% COMPLETE" badge
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               SW - Emu(2_300_000), Emu(150_000),
                               Emu(1_900_000), Emu(380_000))
badge.adjustments[0] = 0.45
set_fill(badge, SUCCESS)
set_line(badge, SUCCESS, 0.5)
badge.shadow.inherit = False
text_in_shape(badge, "● 100% COMPLETE", size=10, bold=True, color=WHITE)

# --- Дээд KPI banner — 3 stat tile -------------------------------------
KPI_Y = Emu(800_000)
KPI_H = Emu(950_000)
KPI_GAP = Emu(180_000)
kpi_total_w = Emu(8_400_000)
kpi_w = (kpi_total_w - 2 * KPI_GAP) // 3
kpi_left = (SW - kpi_total_w) // 2

stats = [
    ("8 / 8",    "Хэрэгжсэн зорилт"),
    ("4.3 / 5", "Хэрэглэхэд хялбар (UX)"),
    ("3 — 4 мин", "Автомат деплой хугацаа"),
]
for i, (val, lbl) in enumerate(stats):
    cx = kpi_left + i * (kpi_w + KPI_GAP) + kpi_w // 2
    add_stat(slide, cx, KPI_Y + KPI_H // 2, kpi_w, KPI_H, val, lbl)

# --- 8 ширхэг checklist карт (2 баган × 4 эгнээ) ----------------------
ITEMS = [
    ("Flutter мобайл апп",     "8 дэлгэц · iOS / Android",                 ACCENT),
    ("Node.js REST API",       "8 endpoint · production EC2",              ACCENT),
    ("Firebase нэгтгэл",       "Firestore · Auth OTP · локал мэдэгдэл",    ACCENT),
    ("AI Chat туслах",         "Claude Haiku + Groq Llama · < 1 сек",      ACCENT),
    ("Давхцал шалгалт",        "Compound index · atomic transaction",      ACCENT),
    ("CI/CD Pipeline",         "GitHub Actions · автомат деплой",          ACCENT),
    ("Туршилт (нэгтгэсэн)",   "9 функциональ + 5 гүйцэтгэлийн алхам",    ACCENT),
    ("Хэрэглэхэд хялбар",      "4.3 / 5.0 оноо (Usability scale)",         ACCENT),
]

CARD_TOP = Emu(1_950_000)
CARD_W   = Emu(5_550_000)
CARD_H   = Emu(1_000_000)
ROW_GAP  = Emu(120_000)
COL_GAP  = Emu(180_000)
LEFT_X   = (SW - 2 * CARD_W - COL_GAP) // 2

for i, (title, sub, color) in enumerate(ITEMS):
    col = i % 2
    row = i // 2
    cx = LEFT_X + col * (CARD_W + COL_GAP)
    cy = CARD_TOP + row * (CARD_H + ROW_GAP)
    add_result_card(slide, cx, cy, CARD_W, CARD_H, title, sub, color=color)

# --- Хадгалах -----------------------------------------------------------
prs.save(DST)
print("Done")
