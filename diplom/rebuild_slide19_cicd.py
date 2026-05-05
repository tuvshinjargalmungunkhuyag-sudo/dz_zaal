"""
v3_editable_v10.pptx -> v3_editable_v11.pptx
Slide 19 "CI/CD Pipeline" — ЦЭВЭР ХЭВТЭЭ pipeline хэлбэрээр шинээр зурна.
Бүх элемент жигд тохирсон, slide-аас гарахгүй.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v10.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v11.pptx"

# --- Өнгө ---------------------------------------------------------------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x1B, 0x1F, 0x3A)
INK_SOFT   = RGBColor(0x4A, 0x52, 0x6E)
PAPER      = RGBColor(0xFA, 0xFA, 0xFA)
PAPER2     = RGBColor(0xEE, 0xF1, 0xF8)
ACCENT     = RGBColor(0x1F, 0x3A, 0x68)
ACCENT2    = RGBColor(0xC2, 0x6A, 0x12)

BUILD_DARK   = RGBColor(0x1F, 0x3A, 0x68)   # хөх — build phase
BUILD_FILL   = RGBColor(0xE3, 0xEA, 0xF6)
DEPLOY_DARK  = RGBColor(0x1E, 0x88, 0x4F)   # ногоон — deploy phase
DEPLOY_FILL  = RGBColor(0xE7, 0xF5, 0xEC)
DECISION_DARK= RGBColor(0xB7, 0x8A, 0x12)
DECISION_FILL= RGBColor(0xFF, 0xF8, 0xE1)
ERROR_DARK   = RGBColor(0xC0, 0x39, 0x2B)
ERROR_FILL   = RGBColor(0xFD, 0xEC, 0xEA)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0, dash=None):
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)
    if dash is not None:
        ln = line._get_or_add_ln()
        prst = ln.find(qn('a:prstDash'))
        if prst is None:
            prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)


def text_in_shape(shape, text, *, size=11, bold=False, color=INK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  italic=False, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(40_000); tf.margin_right = Emu(40_000)
    tf.margin_top = Emu(20_000);  tf.margin_bottom = Emu(20_000)
    tf.vertical_anchor = anchor
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font


def add_textbox(slide, x, y, w, h, text, *, size=11, bold=False,
                italic=False, color=INK, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    return tb


def add_arrow(slide, x1, y1, x2, y2, *,
              color=INK_SOFT, width_pt=1.5, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line = conn.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    ln = line._get_or_add_ln()
    if dash is not None:
        prst = ln.find(qn('a:prstDash'))
        if prst is None:
            prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return conn


def add_step_card(slide, cx, cy, w, h, num, label, *, color_fill, color_dark):
    """Дугаартай (01..06) хэвтээ алхмын карт. Доороо нэртэй."""
    x, y = cx - w // 2, cy - h // 2

    # Үндсэн карт
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.10
    set_fill(card, color_fill)
    set_line(card, color_dark, 1.5)
    card.shadow.inherit = False

    # Дугаар бөмбөлөг (карт дотор дээр төв дээр)
    num_d = Emu(520_000)
    num_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - num_d // 2, y + Emu(180_000),
        num_d, num_d)
    set_fill(num_circle, color_dark)
    set_line(num_circle, color_dark, 0.25)
    text_in_shape(num_circle, num, size=14, bold=True, color=WHITE)

    # Алхмын нэр (доор)
    add_textbox(slide,
                x + Emu(80_000),
                y + Emu(180_000) + num_d + Emu(80_000),
                w - Emu(160_000), Emu(420_000),
                label, size=11, bold=True, color=color_dark,
                align=PP_ALIGN.CENTER, font="Consolas")
    return card


def add_decision(slide, cx, cy, w, h, text, *, fill, dark):
    x, y = cx - w // 2, cy - h // 2
    box = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, w, h)
    set_fill(box, fill)
    set_line(box, dark, 1.5)
    box.shadow.inherit = False
    text_in_shape(box, text, size=11, bold=True, color=dark)
    return box


def add_chevron(slide, cx, cy, w, h, color):
    """Жижигхэн өнгөт chevron (pentagon) сум."""
    x, y = cx - w // 2, cy - h // 2
    arrow = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, x, y, w, h)
    set_fill(arrow, color)
    set_line(arrow, color, 0.25)
    return arrow


# --- Презентац ачаалах --------------------------------------------------
prs = Presentation(SRC)
slide = prs.slides[18]   # 19-р slide

# Хуучин CI/CD shape-уудыг устгах (гадна frame Rectangle-уудыг үлдээх)
KEEP = {"Rectangle 1", "Rectangle 2", "TextBox 3",
        "Rectangle 4", "SlideNum"}
to_delete = []
for sh in slide.shapes:
    if sh.name not in KEEP:
        to_delete.append(sh)
for sh in to_delete:
    remove_shape(sh)

SW = prs.slide_width
SH = prs.slide_height

# --- Title ---------------------------------------------------------------
title_tb = slide.shapes.add_textbox(Emu(420_000), Emu(120_000),
                                    SW - Emu(2_500_000), Emu(450_000))
tf = title_tb.text_frame
tf.margin_left = Emu(0); tf.margin_right = Emu(0)
tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r1 = p.add_run(); r1.text = "Бүлэг 3.1:  "
r1.font.size = Pt(14); r1.font.bold = True
r1.font.color.rgb = ACCENT2; r1.font.name = "Calibri"
r2 = p.add_run(); r2.text = "CI/CD Pipeline — Docker + AWS EC2"
r2.font.size = Pt(20); r2.font.bold = True
r2.font.color.rgb = ACCENT; r2.font.name = "Calibri"

# Баруун дээд "AUTOMATED" badge
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               SW - Emu(2_300_000), Emu(150_000),
                               Emu(1_900_000), Emu(380_000))
badge.adjustments[0] = 0.45
set_fill(badge, DEPLOY_DARK)
set_line(badge, DEPLOY_DARK, 0.5)
badge.shadow.inherit = False
text_in_shape(badge, "● AUTOMATED  ~3 МИН", size=10, bold=True, color=WHITE)

# --- Phase headers (BUILD / DEPLOY) -------------------------------------
PHASE_Y = Emu(900_000)
PHASE_H = Emu(380_000)
LEFT_PAD = Emu(450_000)

# Хагас slide width-ийг хэмжих
half_w = (SW - 2 * LEFT_PAD) // 2

# BUILD phase header
add_textbox(slide,
            LEFT_PAD, PHASE_Y, half_w, PHASE_H,
            "▶ BUILD PHASE",
            size=12, bold=True, color=BUILD_DARK,
            align=PP_ALIGN.LEFT)
# DEPLOY phase header
add_textbox(slide,
            LEFT_PAD + half_w, PHASE_Y, half_w, PHASE_H,
            "▶ DEPLOY PHASE",
            size=12, bold=True, color=DEPLOY_DARK,
            align=PP_ALIGN.LEFT)

# Phase header-ийн доор тус бүрд жижиг өнгөт шугам
sep_y = PHASE_Y + PHASE_H + Emu(40_000)
sep1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              LEFT_PAD, sep_y,
                              half_w - Emu(200_000), Emu(40_000))
set_fill(sep1, BUILD_DARK)
sep1.line.fill.background()

sep2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              LEFT_PAD + half_w, sep_y,
                              half_w - Emu(200_000), Emu(40_000))
set_fill(sep2, DEPLOY_DARK)
sep2.line.fill.background()

# --- Pipeline (нэг хэвтээ мөр) -------------------------------------------
PIPE_CY = Emu(2_900_000)             # pipeline-ийн төв y
CARD_W  = Emu(1_400_000)
CARD_H  = Emu(1_500_000)
DIAM_W  = Emu(1_400_000)
DIAM_H  = Emu(1_400_000)
CHEV_W  = Emu(280_000)
CHEV_H  = Emu(380_000)

# 6 step + 1 diamond = 7 unit; chevron-ыг тэдгээрийн дунд
# Эхлээд pipeline-ыг төвлөрүүлж байрлуулъя
unit_w = CARD_W            # бүх unit-ыг ижил W болгон тооцоолж zay тооцох
total_units_w = 6 * CARD_W + DIAM_W
gaps = 7                   # 7 chevron (4 step хооронд + diamond-ийн өмнөх + дараах)
total_chev_w = gaps * (CHEV_W + Emu(80_000) * 2)   # chevron + 2 талын margin
pipe_w = total_units_w + total_chev_w
pipe_left = (SW - pipe_w) // 2

# Алхмууд + diamond байрлал
items = [
    ("step", "01", "git push",       BUILD_FILL,  BUILD_DARK),
    ("step", "02", "GH Actions",     BUILD_FILL,  BUILD_DARK),
    ("step", "03", "Docker build",   BUILD_FILL,  BUILD_DARK),
    ("diamond", "", "Build OK?",     DECISION_FILL, DECISION_DARK),
    ("step", "04", "ECR push",       DEPLOY_FILL, DEPLOY_DARK),
    ("step", "05", "SSH to EC2",     DEPLOY_FILL, DEPLOY_DARK),
    ("step", "06", "docker run",     DEPLOY_FILL, DEPLOY_DARK),
]

cursor = pipe_left
diamond_cx = None
last_card_right_edge = None
positions = []

for i, (kind, num, label, fill, dark) in enumerate(items):
    if kind == "step":
        cx = cursor + CARD_W // 2
        add_step_card(slide, cx, PIPE_CY, CARD_W, CARD_H, num, label,
                      color_fill=fill, color_dark=dark)
        positions.append((kind, cx, CARD_W))
        cursor += CARD_W
    else:  # diamond
        cx = cursor + DIAM_W // 2
        add_decision(slide, cx, PIPE_CY, DIAM_W, DIAM_H, label,
                     fill=fill, dark=dark)
        diamond_cx = cx
        positions.append((kind, cx, DIAM_W))
        cursor += DIAM_W

    # Chevron (хамгийн сүүлчийн item-аас бусдад)
    if i < len(items) - 1:
        cursor += Emu(80_000)
        chev_color = items[i + 1][4]   # дараагийн item-ийн dark өнгө
        add_chevron(slide,
                    cursor + CHEV_W // 2, PIPE_CY,
                    CHEV_W, CHEV_H, chev_color)
        cursor += CHEV_W + Emu(80_000)

# Diamond-ийн "Тийм" label (диамантийн доор)
add_textbox(slide,
            diamond_cx - Emu(400_000), PIPE_CY + DIAM_H // 2 - Emu(40_000),
            Emu(800_000), Emu(280_000),
            "Тийм →", size=10, bold=True, italic=True,
            color=DEPLOY_DARK, align=PP_ALIGN.CENTER)

# --- Fail branch (diamond-аас доош) -------------------------------------
FAIL_CY = PIPE_CY + DIAM_H // 2 + Emu(800_000)
FAIL_W = Emu(2_400_000)
FAIL_H = Emu(580_000)

# diamond → fail сум (босоо)
add_arrow(slide,
          diamond_cx, PIPE_CY + DIAM_H // 2,
          diamond_cx, FAIL_CY - FAIL_H // 2,
          color=ERROR_DARK, width_pt=1.5)

# "Үгүй" label сумны хажууд
add_textbox(slide,
            diamond_cx + Emu(60_000),
            PIPE_CY + DIAM_H // 2 + Emu(80_000),
            Emu(700_000), Emu(280_000),
            "Үгүй", size=10, bold=True, italic=True,
            color=ERROR_DARK, align=PP_ALIGN.LEFT)

# Fail card
fx = diamond_cx - FAIL_W // 2
fy = FAIL_CY - FAIL_H // 2
fail_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  fx, fy, FAIL_W, FAIL_H)
fail_box.adjustments[0] = 0.20
set_fill(fail_box, ERROR_FILL)
set_line(fail_box, ERROR_DARK, 1.5)
fail_box.shadow.inherit = False
add_textbox(slide,
            fx + Emu(140_000), fy, Emu(420_000), FAIL_H,
            "✕", size=22, bold=True, color=ERROR_DARK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_textbox(slide,
            fx + Emu(560_000), fy, FAIL_W - Emu(640_000), FAIL_H,
            "Workflow fail — лог GitHub UI",
            size=11, bold=True, color=ERROR_DARK,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# --- Доод тайлбар самбар (4 баган) -------------------------------------
INFO_X = Emu(450_000)
INFO_Y = Emu(5_500_000)
INFO_W = SW - 2 * INFO_X
INFO_H = Emu(1_100_000)

info_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  INFO_X, INFO_Y, INFO_W, INFO_H)
info_box.adjustments[0] = 0.08
set_fill(info_box, PAPER2)
set_line(info_box, ACCENT, 1.0)
info_box.shadow.inherit = False

# 4 баганатай мэдээллийн tile
infos = [
    ("⏱", "ХУГАЦАА",  "3–4 минут",          ACCENT),
    ("🔄", "ҮЕ ШАТ",   "6 алхамтай",         BUILD_DARK),
    ("⚙", "ЗЭРЭГЦЭХ", "бүтэн автомат",      ACCENT2),
    ("📊", "ХЯНАЛТ",   "GitHub UI · log",    DEPLOY_DARK),
]

col_w = (INFO_W - Emu(400_000)) // 4
for i, (icon, label, value, color) in enumerate(infos):
    cx = INFO_X + Emu(200_000) + i * col_w + col_w // 2

    # Icon (том)
    add_textbox(slide,
                cx - Emu(400_000), INFO_Y + Emu(150_000),
                Emu(800_000), Emu(380_000),
                icon, size=20, bold=True, color=color,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Label (жижиг)
    add_textbox(slide,
                cx - Emu(900_000), INFO_Y + Emu(560_000),
                Emu(1_800_000), Emu(240_000),
                label, size=9, bold=True, color=color,
                align=PP_ALIGN.CENTER)
    # Value (том)
    add_textbox(slide,
                cx - Emu(900_000), INFO_Y + Emu(800_000),
                Emu(1_800_000), Emu(280_000),
                value, size=11, bold=True, color=INK,
                align=PP_ALIGN.CENTER)

    # Багана хооронд босоо тусгаарлагч
    if i < 3:
        sep_x = INFO_X + Emu(200_000) + (i + 1) * col_w
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     sep_x, INFO_Y + Emu(280_000),
                                     Emu(15_000), INFO_H - Emu(560_000))
        set_fill(sep, RGBColor(0xCF, 0xD8, 0xEC))
        sep.line.fill.background()

# --- Хадгалах -----------------------------------------------------------
prs.save(DST)
print("Done")
