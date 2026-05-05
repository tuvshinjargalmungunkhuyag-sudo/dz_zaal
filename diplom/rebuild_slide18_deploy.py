"""
v3_editable_v8.pptx -> v3_editable_v9.pptx
Slide 18 дээрх 4 хоосон rectangle (github/Local host/aws/playstore)-г
устгаад, Local Host → GitHub → AWS → Play Store чиглэлийн deployment
architecture-ыг native PowerPoint shape-аар зурна.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v8.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v9.pptx"

# --- Өнгө ----------------------------------------------------------------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x1B, 0x1F, 0x3A)
INK_SOFT   = RGBColor(0x4A, 0x52, 0x6E)
PAPER      = RGBColor(0xFA, 0xFA, 0xFA)
PAPER2     = RGBColor(0xEE, 0xF1, 0xF8)
ACCENT     = RGBColor(0x1F, 0x3A, 0x68)
ARROW_GRAY = RGBColor(0x6B, 0x72, 0x8A)

# Stage өнгөнүүд (тус бүр өөр өнгө — brand-like)
STAGE_COLORS = [
    {"dark": RGBColor(0x16, 0x6B, 0xAB), "fill": RGBColor(0xE3, 0xEF, 0xF8)},  # Local
    {"dark": RGBColor(0x24, 0x29, 0x3E), "fill": RGBColor(0xE5, 0xE7, 0xEC)},  # GitHub
    {"dark": RGBColor(0xE6, 0x7E, 0x22), "fill": RGBColor(0xFD, 0xEC, 0xD8)},  # AWS
    {"dark": RGBColor(0x1E, 0x88, 0x4F), "fill": RGBColor(0xE7, 0xF5, 0xEC)},  # Play Store
]


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0, dash=None):
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)
    if dash is not None:
        ln = line._get_or_add_ln()
        prst = ln.find(qn('a:prstDash'))
        if prst is None:
            prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)


def no_line(shape):
    shape.line.fill.background()


def text_in_shape(shape, text, *, size=12, bold=False, color=INK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  italic=False, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(40_000); tf.margin_right = Emu(40_000)
    tf.margin_top = Emu(20_000);  tf.margin_bottom = Emu(20_000)
    tf.vertical_anchor = anchor
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font


def add_textbox(slide, x, y, w, h, text, *, size=12, bold=False,
                italic=False, color=INK, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    return tb


def add_arrow(slide, x1, y1, x2, y2, *, color=ARROW_GRAY, width_pt=1.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line = conn.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    ln = line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return conn


def add_stage_card(slide, cx, cy, w, h, *,
                   icon, stage_label, name, sub_items, color):
    """Үе шатны том карт: icon, stage label, нэр, доор tool-ууд."""
    x, y = cx - w // 2, cy - h // 2
    fill = color["fill"]; dark = color["dark"]

    # Үндсэн карт
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.07
    set_fill(card, WHITE)
    set_line(card, dark, 2.0)
    card.shadow.inherit = False

    # Дээд талын өнгөт зурвас (header)
    HEAD_H = Emu(800_000)
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x + Emu(40_000),
                                  y + Emu(40_000),
                                  w - Emu(80_000),
                                  HEAD_H)
    set_fill(head, dark)
    set_line(head, dark, 0.25)

    # Stage label (TIER 01 etc) — толгойн дээд буланд
    add_textbox(slide,
                x + Emu(60_000), y + Emu(60_000),
                w - Emu(120_000), Emu(280_000),
                stage_label, size=9, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)

    # Icon (том emoji/symbol)
    add_textbox(slide,
                x + Emu(60_000), y + Emu(280_000),
                w - Emu(120_000), Emu(560_000),
                icon, size=32, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Stage name (header-ын доор, том)
    add_textbox(slide,
                x + Emu(40_000), y + Emu(40_000) + HEAD_H + Emu(150_000),
                w - Emu(80_000), Emu(450_000),
                name, size=16, bold=True, color=dark,
                align=PP_ALIGN.CENTER)

    # Tool жагсаалт
    base_y = y + Emu(40_000) + HEAD_H + Emu(620_000)
    row_h = Emu(320_000)
    for i, item in enumerate(sub_items):
        # Жижиг bullet
        add_textbox(slide,
                    x + Emu(180_000), base_y + i * row_h,
                    Emu(180_000), row_h,
                    "•", size=14, bold=True, color=dark,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide,
                    x + Emu(380_000), base_y + i * row_h,
                    w - Emu(420_000), row_h,
                    item, size=10, color=INK_SOFT,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    return card


def add_chevron_arrow(slide, x, y, w, h, label, color):
    """Stage-уудын дунд chevron-маягийн анхаарал татсан label-тай сум."""
    # Pentagon shape-аар chevron хийх
    arrow = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, x, y, w, h)
    set_fill(arrow, color["dark"])
    set_line(arrow, color["dark"], 0.25)
    text_in_shape(arrow, label, size=10, bold=True, color=WHITE)
    return arrow


# --- Презентац ачаалах --------------------------------------------------
prs = Presentation(SRC)
slide = prs.slides[17]   # 18-р slide

# Бүх хуучин shape-ыг устгах
to_delete = list(slide.shapes)
for sh in to_delete:
    remove_shape(sh)

SW = prs.slide_width
SH = prs.slide_height

# --- Title ---------------------------------------------------------------
title_tb = slide.shapes.add_textbox(Emu(420_000), Emu(180_000),
                                    SW - Emu(2_500_000), Emu(450_000))
tf = title_tb.text_frame
tf.margin_left = Emu(0); tf.margin_right = Emu(0)
tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r1 = p.add_run(); r1.text = "Бүлэг 3.1:  "
r1.font.size = Pt(14); r1.font.bold = True
r1.font.color.rgb = STAGE_COLORS[2]["dark"]
r1.font.name = "Calibri"
r2 = p.add_run(); r2.text = "Үүлэн орчны Деплой Архитектур"
r2.font.size = Pt(20); r2.font.bold = True
r2.font.color.rgb = ACCENT; r2.font.name = "Calibri"

# Баруун дээд буланд "AUTO DEPLOY" badge
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               SW - Emu(2_300_000), Emu(180_000),
                               Emu(1_900_000), Emu(380_000))
badge.adjustments[0] = 0.45
set_fill(badge, ACCENT)
set_line(badge, ACCENT, 0.5)
badge.shadow.inherit = False
text_in_shape(badge, "3–4 МИН АВТОМАТ", size=10, bold=True, color=WHITE)

# --- 4 stage card-ыг хэвтээ дараалуулах ---------------------------------
stages = [
    {
        "icon": "</>",
        "stage_label": "STAGE 01",
        "name": "Local Host",
        "items": ["VS Code, Cursor", "Flutter SDK", "Docker desktop", "git commit"],
        "color": STAGE_COLORS[0],
    },
    {
        "icon": "⎇",
        "stage_label": "STAGE 02",
        "name": "GitHub",
        "items": ["Repository (main)", "GitHub Actions", "Workflow YAML", "secrets / IAM"],
        "color": STAGE_COLORS[1],
    },
    {
        "icon": "☁",
        "stage_label": "STAGE 03",
        "name": "AWS Cloud",
        "items": ["ECR — Docker registry", "EC2 (ap-southeast-1)", "Node.js контейнер", "REST API endpoint"],
        "color": STAGE_COLORS[2],
    },
    {
        "icon": "▶",
        "stage_label": "STAGE 04",
        "name": "Play Store",
        "items": ["Android .aab", "Internal testing", "Production track", "Хэрэглэгчид"],
        "color": STAGE_COLORS[3],
    },
]

# Layout тохиргоо
CARD_W = Emu(2_400_000)
CARD_H = Emu(3_400_000)
ARROW_W = Emu(450_000)
ARROW_H = Emu(450_000)

# Card-ын төв y
CARD_CY = Emu(3_500_000)

# Эхний card-ын зүүн margin
total_w = 4 * CARD_W + 3 * ARROW_W + Emu(360_000)   # 4 card + 3 arrow + жижиг зай
left_margin = (SW - total_w) // 2

# Stage card болон chevron arrow-уудыг ээлжлэн зурах
arrow_labels = ["push", "deploy", "publish"]

for i, stage in enumerate(stages):
    # Card-ын төв X
    cx = left_margin + CARD_W // 2 + i * (CARD_W + ARROW_W + Emu(120_000))
    add_stage_card(slide, cx, CARD_CY, CARD_W, CARD_H,
                   icon=stage["icon"],
                   stage_label=stage["stage_label"],
                   name=stage["name"],
                   sub_items=stage["items"],
                   color=stage["color"])

    # Chevron сум — энэ card-ийн дараа (хамгийн сүүлчийн card-аас бусдад)
    if i < 3:
        ax = cx + CARD_W // 2 + Emu(60_000)
        ay = CARD_CY - ARROW_H // 2
        # Цайвар chevron — шилжих stage-ийн өнгөөр
        add_chevron_arrow(slide, ax, ay, ARROW_W, ARROW_H,
                          arrow_labels[i],
                          stages[i + 1]["color"])

# --- Доод тайлбар: workflow дарааллын товч тайлбар ---------------------
DESC_X = Emu(450_000)
DESC_Y = Emu(5_700_000)
DESC_W = SW - 2 * DESC_X
DESC_H = Emu(900_000)

desc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  DESC_X, DESC_Y, DESC_W, DESC_H)
desc_box.adjustments[0] = 0.10
set_fill(desc_box, PAPER2)
set_line(desc_box, ACCENT, 1.0)
desc_box.shadow.inherit = False

# Доод хэсгийн тайлбар (4 баганаар)
sub_w = (DESC_W - Emu(400_000)) // 4
for i, (label, info) in enumerate([
        ("➀ git push",         "локалаас main руу"),
        ("➁ Actions trigger",  "Docker build + ECR push"),
        ("➂ EC2 redeploy",     "SSH + docker pull/run"),
        ("➃ App release",      "AAB → Production")]):
    sub_x = DESC_X + Emu(200_000) + i * sub_w
    add_textbox(slide,
                sub_x, DESC_Y + Emu(120_000),
                sub_w - Emu(80_000), Emu(280_000),
                label, size=11, bold=True,
                color=stages[i]["color"]["dark"],
                align=PP_ALIGN.LEFT)
    add_textbox(slide,
                sub_x, DESC_Y + Emu(420_000),
                sub_w - Emu(80_000), Emu(380_000),
                info, size=10, color=INK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

# --- Хадгалах ----------------------------------------------------------
prs.save(DST)
print("Done")
