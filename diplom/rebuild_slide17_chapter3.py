"""
v3_editable_v7.pptx -> v3_editable_v8.pptx
17-р slide дээр "Бүлэг 3 — Хэрэгжүүлэлт" section cover-ыг native
PowerPoint shape-аар зурна (slide 8-той ижил дизайн).
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v7.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v8.pptx"

# --- Өнгөний палет (slide 8-тай ижил) -----------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x1B, 0x1F, 0x3A)
INK_SOFT    = RGBColor(0x4A, 0x52, 0x6E)
ACCENT      = RGBColor(0x1F, 0x3A, 0x68)   # хөх
ACCENT2     = RGBColor(0xC2, 0x6A, 0x12)   # улбар шар
ACCENT3     = RGBColor(0x1E, 0x88, 0x4F)   # ногоон
ACCENT4     = RGBColor(0xC0, 0x39, 0x2B)   # улаан
PAPER3      = RGBColor(0xF5, 0xF7, 0xFB)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    return tb


# --- Презентац ачаалах --------------------------------------------------
prs = Presentation(SRC)
slide = prs.slides[16]   # 17-р slide

# Хуучин placeholder-уудыг бүгдийг устгах
to_delete = list(slide.shapes)
for sh in to_delete:
    remove_shape(sh)

SW = prs.slide_width
SH = prs.slide_height

# --- 1) Дэвсгэр (зүүн талд цайвар хөх блок) -----------------------------
LEFT_BAND_W = Emu(3_500_000)
left_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   0, 0, LEFT_BAND_W, SH)
set_fill(left_band, ACCENT)
no_line(left_band)

# CHAPTER label дугаарын дээр
add_text(slide, Emu(300000), Emu(1_050_000),
         LEFT_BAND_W - Emu(600000), Emu(350000),
         "CHAPTER", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)

# Том "03" дугаар (chapter number)
add_text(slide, Emu(300000), Emu(1_400_000),
         LEFT_BAND_W - Emu(600000), Emu(2_500_000),
         "03", size=200, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Доод талын тэмдэглэгээ
add_text(slide, Emu(300000), Emu(4_300_000),
         LEFT_BAND_W - Emu(600000), Emu(350000),
         "M.Tүвшинжаргал · 2026", size=10, italic=True,
         color=RGBColor(0xCF, 0xD8, 0xEC),
         align=PP_ALIGN.CENTER)

# --- 2) Баруун талд гарчиг ----------------------------------------------
TXT_X = LEFT_BAND_W + Emu(700_000)
TXT_W = SW - TXT_X - Emu(500_000)

# БҮЛЭГ 3 жижиг гарчиг
add_text(slide, TXT_X, Emu(1_200_000), TXT_W, Emu(700_000),
         "БҮЛЭГ 3", size=22, bold=True, color=ACCENT2,
         align=PP_ALIGN.LEFT)

# Үндсэн нэр
add_text(slide, TXT_X, Emu(1_750_000), TXT_W, Emu(1_200_000),
         "Үүлэн орчны\nхэрэгжүүлэлт",
         size=44, bold=True, color=ACCENT,
         align=PP_ALIGN.LEFT)

# Шугам
line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    TXT_X, Emu(3_400_000),
                                    Emu(900_000), Emu(70_000))
set_fill(line_shape, ACCENT2)
no_line(line_shape)

# --- 3) 4 сэдэв (badge) ------------------------------------------------
topics = [
    ("Docker",        "контейнержуулалт",      ACCENT),
    ("AWS EC2",       "ap-southeast-1",        ACCENT2),
    ("GitHub Actions","CI/CD pipeline",        ACCENT3),
    ("Туршилт",       "функциональ · нэгтгэл", ACCENT4),
]

BADGE_TOP = Emu(3_700_000)
BADGE_W   = Emu(1_750_000)
BADGE_H   = Emu(1_500_000)
BADGE_GAP = Emu(150_000)

for i, (title, sub, color) in enumerate(topics):
    bx = TXT_X + i * (BADGE_W + BADGE_GAP)

    # Үндсэн карт
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  bx, BADGE_TOP, BADGE_W, BADGE_H)
    card.adjustments[0] = 0.10
    set_fill(card, PAPER3)
    set_line(card, color, 1.5)
    card.shadow.inherit = False

    # Дээр өнгөт зурвас
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     bx + Emu(40000), BADGE_TOP + Emu(40000),
                                     BADGE_W - Emu(80000), Emu(80000))
    set_fill(top_bar, color)
    no_line(top_bar)

    # Дугаар (01..04) дугуй badge
    num_d = Emu(360000)
    num_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        bx + (BADGE_W - num_d) // 2, BADGE_TOP + Emu(220000),
        num_d, num_d)
    set_fill(num_circle, color)
    set_line(num_circle, color, 0.25)
    add_text(slide,
             bx + (BADGE_W - num_d) // 2, BADGE_TOP + Emu(220000),
             num_d, num_d,
             f"0{i+1}", size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Сэдвийн нэр
    add_text(slide, bx + Emu(100000), BADGE_TOP + Emu(700000),
             BADGE_W - Emu(200000), Emu(380000),
             title, size=14, bold=True, color=color,
             align=PP_ALIGN.CENTER)

    # Тайлбар
    add_text(slide, bx + Emu(100000), BADGE_TOP + Emu(1_080_000),
             BADGE_W - Emu(200000), Emu(380000),
             sub, size=10, color=INK_SOFT,
             align=PP_ALIGN.CENTER)

# --- Хадгалах -----------------------------------------------------------
prs.save(DST)
print("Done")
