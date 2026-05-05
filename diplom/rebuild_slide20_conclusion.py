"""
v3_editable_v11.pptx -> v3_editable_v12.pptx
Slide 20 "Дүгнэлт" — top summary + 3 баганатай зөв дизайн.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v11.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v12.pptx"

# --- Өнгө ----------------------------------------------------------------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x1B, 0x1F, 0x3A)
INK_SOFT   = RGBColor(0x4A, 0x52, 0x6E)
PAPER      = RGBColor(0xFA, 0xFA, 0xFA)
PAPER2     = RGBColor(0xEE, 0xF1, 0xF8)
ACCENT     = RGBColor(0x1F, 0x3A, 0x68)
NEW_DARK   = RGBColor(0x1E, 0x88, 0x4F)   # ногоон — Шинэлэг тал
NEW_FILL   = RGBColor(0xE7, 0xF5, 0xEC)
IMPACT_DARK= RGBColor(0xC2, 0x6A, 0x12)   # улбар шар — Ач холбогдол
IMPACT_FILL= RGBColor(0xFF, 0xF1, 0xE0)
FUTURE_DARK= RGBColor(0x1F, 0x3A, 0x68)   # хөх — Цаашдын хөгжил
FUTURE_FILL= RGBColor(0xE3, 0xEA, 0xF6)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def text_in_shape(shape, text, *, size=11, bold=False, color=INK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  italic=False, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(60_000); tf.margin_right = Emu(60_000)
    tf.margin_top = Emu(40_000);  tf.margin_bottom = Emu(40_000)
    tf.vertical_anchor = anchor
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font


def add_textbox(slide, x, y, w, h, text, *, size=11, bold=False,
                italic=False, color=INK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    return tb


def add_column(slide, x, y, w, h, *, icon, header, items, fill, dark):
    """Гарчиг + dot bullet жагсаалттай багана карт."""
    # Үндсэн карт
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    set_fill(card, WHITE)
    set_line(card, dark, 1.5)
    card.shadow.inherit = False

    # Дээд талын өнгөт header bar
    HEAD_H = Emu(750_000)
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x + Emu(40_000), y + Emu(40_000),
                                  w - Emu(80_000), HEAD_H)
    set_fill(head, dark)
    no_line(head)

    # Том icon (header дотор зүүн талд)
    add_textbox(slide,
                x + Emu(60_000), y + Emu(40_000),
                Emu(700_000), HEAD_H,
                icon, size=26, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Header текст (баруун талд)
    add_textbox(slide,
                x + Emu(800_000), y + Emu(40_000),
                w - Emu(840_000), HEAD_H,
                header, size=15, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # Bullet жагсаалт
    LIST_TOP = y + Emu(40_000) + HEAD_H + Emu(220_000)
    row_h = Emu(580_000)
    for i, item in enumerate(items):
        ry = LIST_TOP + i * row_h

        # Зүүн талын өнгөт цэг
        dot_d = Emu(180_000)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Emu(220_000), ry + Emu(120_000),
            dot_d, dot_d)
        set_fill(dot, dark)
        set_line(dot, dark, 0.25)

        # Текст
        add_textbox(slide,
                    x + Emu(480_000), ry,
                    w - Emu(560_000), row_h,
                    item, size=12, color=INK,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    return card


# --- Презентац ачаалах --------------------------------------------------
prs = Presentation(SRC)
slide = prs.slides[19]   # 20-р slide

# Бүх placeholder-ыг устгана
to_delete = list(slide.shapes)
for sh in to_delete:
    remove_shape(sh)

SW = prs.slide_width
SH = prs.slide_height

# --- Title ---------------------------------------------------------------
add_textbox(slide,
            Emu(450_000), Emu(120_000),
            SW - Emu(900_000), Emu(450_000),
            "Дүгнэлт",
            size=24, bold=True, color=ACCENT,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# --- Дээд summary banner ------------------------------------------------
SUM_X = Emu(450_000)
SUM_Y = Emu(700_000)
SUM_W = SW - 2 * SUM_X
SUM_H = Emu(950_000)

sum_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 SUM_X, SUM_Y, SUM_W, SUM_H)
sum_box.adjustments[0] = 0.06
set_fill(sum_box, PAPER2)
set_line(sum_box, ACCENT, 1.5)
sum_box.shadow.inherit = False

# Зүүн талын өнгөт зурвас
left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  SUM_X + Emu(60_000),
                                  SUM_Y + Emu(80_000),
                                  Emu(120_000),
                                  SUM_H - Emu(160_000))
set_fill(left_bar, ACCENT)
no_line(left_bar)

# Зүүн дээд буланд жижиг "✓ ҮР ДҮН" tag
tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             SUM_X + Emu(280_000),
                             SUM_Y + Emu(120_000),
                             Emu(1_400_000),
                             Emu(280_000))
tag.adjustments[0] = 0.45
set_fill(tag, ACCENT)
set_line(tag, ACCENT, 0.25)
text_in_shape(tag, "✓ ҮР ДҮН", size=10, bold=True, color=WHITE)

# Summary текст
add_textbox(slide,
            SUM_X + Emu(280_000),
            SUM_Y + Emu(450_000),
            SUM_W - Emu(360_000),
            SUM_H - Emu(490_000),
            "Flutter, Node.js, Firebase давхаргуудыг AWS EC2 үүлэн дэд бүтэцэд "
            "суурилуулан нэгтгэсэн 3 давхаргат ухаалаг захиалгын систем "
            "production орчинд амжилттай хэрэгжлээ.",
            size=13, color=INK,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

# --- 3 баганатай дэлгэрэнгүй -------------------------------------------
COL_TOP = SUM_Y + SUM_H + Emu(280_000)
COL_H   = SH - COL_TOP - Emu(450_000)        # доод тал хүртэл
COL_GAP = Emu(280_000)
COL_W   = (SW - 2 * Emu(450_000) - 2 * COL_GAP) // 3

# Багана 1 — Шинэлэг тал
add_column(slide,
           Emu(450_000), COL_TOP, COL_W, COL_H,
           icon="✦",
           header="Шинэлэг тал",
           items=[
               "3 давхаргат үүлэн архитектур",
               "Давхар горимын AI чатбот",
               "fixed_bookings функц",
           ],
           fill=NEW_FILL, dark=NEW_DARK)

# Багана 2 — Ач холбогдол
add_column(slide,
           Emu(450_000) + COL_W + COL_GAP, COL_TOP, COL_W, COL_H,
           icon="◉",
           header="Ач холбогдол",
           items=[
               "Иргэдэд хялбар, ил тод",
               "Захиалгын ачаалал бууруулна",
               "Орон нутгийн жишиг",
           ],
           fill=IMPACT_FILL, dark=IMPACT_DARK)

# Багана 3 — Цаашдын хөгжил
add_column(slide,
           Emu(450_000) + 2 * (COL_W + COL_GAP), COL_TOP, COL_W, COL_H,
           icon="➤",
           header="Цаашдын хөгжил",
           items=[
               "QPay / SocialPay төлбөр",
               "HTTPS / SSL аюулгүй холболт",
               "Веб админ дэлгэц",
           ],
           fill=FUTURE_FILL, dark=FUTURE_DARK)

# --- Хадгалах -----------------------------------------------------------
prs.save(DST)
print("Done")
