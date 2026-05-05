"""
v3.pptx -> v3_editable.pptx
4-р slide (index 3) дээрх зургийг устгаж, native PowerPoint shape-аар
"Судалгааны асуудал ба шийдэл" диаграммыг зурна.
Бүх text box, rectangle, arrow тус тусдаа засагдана.
"""

from copy import deepcopy
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable.pptx"

# --- Өнгөний палет (presentation_v3.tex-тэй ойролцоо) ----------------------
COLOR_BG          = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_INK         = RGBColor(0x1B, 0x1F, 0x3A)
COLOR_PAPER       = RGBColor(0xFA, 0xFA, 0xFA)
COLOR_PAPER2      = RGBColor(0xEE, 0xF1, 0xF8)
COLOR_PROBLEM     = RGBColor(0xC0, 0x39, 0x2B)   # улаан — асуудал
COLOR_PROBLEM_BG  = RGBColor(0xFD, 0xEC, 0xEA)
COLOR_SOLUTION    = RGBColor(0x1E, 0x88, 0x4F)   # ногоон — шийдэл
COLOR_SOLUTION_BG = RGBColor(0xE7, 0xF5, 0xEC)
COLOR_ACCENT      = RGBColor(0x1F, 0x3A, 0x68)   # хөх — accent
COLOR_ACCENT_BG   = RGBColor(0xE3, 0xEA, 0xF6)


def remove_shape(shape):
    """Slide-аас shape устгана."""
    sp = shape._element
    sp.getparent().remove(sp)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0):
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)


def add_rounded_rect(slide, x, y, w, h, fill, line_color, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.10
    set_fill(shp, fill)
    set_line(shp, line_color, line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color=COLOR_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_bullet_box(slide, x, y, w, h, items, *,
                   fill, border, text_color=COLOR_INK, size=14):
    """Bullet point-уудтай rounded rectangle нэмнэ."""
    box = add_rounded_rect(slide, x, y, w, h, fill, border, line_w=1.5)
    # Bullet текстийг тус тусын text box-аар хийж, хэрэглэгч бүгдийг засаж чадна
    pad_x = Emu(180000)
    pad_y = Emu(150000)
    line_h = Emu(420000)
    for i, item in enumerate(items):
        ty = y + pad_y + i * line_h
        # Bullet цэг
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + pad_x, ty + Emu(120000),
            Emu(110000), Emu(110000))
        set_fill(dot, border)
        set_line(dot, border, 0.25)
        # Текст
        add_text(slide, x + pad_x + Emu(220000), ty,
                 w - 2 * pad_x - Emu(220000), line_h,
                 item, size=size, color=text_color)
    return box


# --- Презентац ачаалах ----------------------------------------------------
prs = Presentation(SRC)
slide = prs.slides[3]

# Хуучин shape-уудыг устгана (title placeholder-ыг үлдээнэ)
to_delete = []
for sh in slide.shapes:
    if sh.shape_type == 13:                          # PICTURE
        to_delete.append(sh)
    elif sh.name == "Content Placeholder 3":
        to_delete.append(sh)
for sh in to_delete:
    remove_shape(sh)

# --- Title-ыг шинэчлэх ----------------------------------------------------
for sh in slide.shapes:
    if sh.name == "Title 1" and sh.has_text_frame:
        sh.text_frame.text = "Судалгааны асуудал ба шийдэл"
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(32)
                r.font.bold = True
                r.font.color.rgb = COLOR_ACCENT
        break

SW = prs.slide_width        # 12 192 000
SH = prs.slide_height       # 6 858 000

# --- Layout тогтмолууд ----------------------------------------------------
TOP = Emu(1_700_000)               # title-аас доош
COL_W = Emu(5_200_000)
COL_H = Emu(3_900_000)
LEFT_X = Emu(300_000)
RIGHT_X = SW - LEFT_X - COL_W
HEAD_H = Emu(550_000)
GAP = Emu(120_000)

# --- Зүүн багана: АСУУДАЛ ---------------------------------------------------
# Header
head_l = add_rounded_rect(slide, LEFT_X, TOP, COL_W, HEAD_H,
                          COLOR_PROBLEM, COLOR_PROBLEM)
add_text(slide, LEFT_X, TOP, COL_W, HEAD_H,
         "Тулгамдаж буй асуудлууд",
         size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bullets
problems = [
    "Утсаар захиалахад шугам завгүй",
    "Биечлэн очих → цаг, зардал",
    "Бодит цагийн хуваарь байхгүй",
    "Цуцлах, сануулга байхгүй",
    "Цаасан бүртгэл → алдаа, өгөгдөл алдагдах",
    "No-show, тогтмол захиалга ялгахад хүндрэл",
]
add_bullet_box(slide, LEFT_X, TOP + HEAD_H + GAP, COL_W, COL_H,
               problems, fill=COLOR_PROBLEM_BG, border=COLOR_PROBLEM,
               size=14)

# --- Баруун багана: ШИЙДЭЛ ------------------------------------------------
head_r = add_rounded_rect(slide, RIGHT_X, TOP, COL_W, HEAD_H,
                          COLOR_SOLUTION, COLOR_SOLUTION)
add_text(slide, RIGHT_X, TOP, COL_W, HEAD_H,
         "Системийн шийдэл",
         size=20, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

solutions = [
    "Flutter мобайл аппликейшн (iOS / Android)",
    "Firestore real-time цагийн хуваарь",
    "DBK-xxxxx захиалгын баталгаажуулалт",
    "Локал мэдэгдэл (1 цагийн өмнө)",
    "Cloud Firestore — найдвартай хадгалалт",
    "fixed_bookings + AI чатбот туслах",
]
add_bullet_box(slide, RIGHT_X, TOP + HEAD_H + GAP, COL_W, COL_H,
               solutions, fill=COLOR_SOLUTION_BG, border=COLOR_SOLUTION,
               size=14)

# --- Дунд талын сум (асуудал → шийдэл) ------------------------------------
arrow_w = Emu(900_000)
arrow_h = Emu(500_000)
arrow_x = LEFT_X + COL_W + (RIGHT_X - LEFT_X - COL_W - arrow_w) // 2
arrow_y = TOP + HEAD_H + GAP + (COL_H - arrow_h) // 2
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               arrow_x, arrow_y, arrow_w, arrow_h)
set_fill(arrow, COLOR_ACCENT)
set_line(arrow, COLOR_ACCENT, 0.5)

# Сумны доор/дээр label
add_text(slide, arrow_x - Emu(200000), arrow_y - Emu(450000),
         arrow_w + Emu(400000), Emu(400000),
         "ШИЙДЭХ", size=11, bold=True, color=COLOR_ACCENT,
         align=PP_ALIGN.CENTER)

# --- Доод хэсэг: 3 түлхүүр компонент --------------------------------------
KEY_TOP = TOP + HEAD_H + GAP + COL_H + Emu(250_000)
KEY_H = Emu(600_000)
KEY_GAP = Emu(200_000)
KEY_W = (SW - 2 * LEFT_X - 2 * KEY_GAP) // 3

keys = [
    ("Мобайл аппликейшн", COLOR_PROBLEM),
    ("3 давхаргат архитектур", COLOR_ACCENT),
    ("Үүлэн дэд бүтэц (AWS)", COLOR_SOLUTION),
]
for i, (label, color) in enumerate(keys):
    kx = LEFT_X + i * (KEY_W + KEY_GAP)
    box = add_rounded_rect(slide, kx, KEY_TOP, KEY_W, KEY_H,
                           COLOR_PAPER2, color, line_w=1.5)
    add_text(slide, kx, KEY_TOP, KEY_W, KEY_H,
             label, size=15, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# --- Хадгалах -------------------------------------------------------------
prs.save(DST)
print("Saved:", DST)
