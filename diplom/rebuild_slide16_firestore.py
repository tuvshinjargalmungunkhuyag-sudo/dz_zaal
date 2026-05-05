"""
v3_editable_v6.pptx -> v3_editable_v7.pptx
16-р slide "Firebase Firestore — Өгөгдлийн Схем" зургийг native PowerPoint
shape-аар сэргээж зурна. ER-style 4 collection card.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v6.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v7.pptx"

# --- Өнгө ----------------------------------------------------------------
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x1B, 0x1F, 0x3A)
INK_SOFT  = RGBColor(0x4A, 0x52, 0x6E)
PAPER     = RGBColor(0xFA, 0xFA, 0xFA)
PAPER2    = RGBColor(0xF5, 0xF7, 0xFB)
ACCENT    = RGBColor(0x1F, 0x3A, 0x68)
RELATION  = RGBColor(0x6B, 0x72, 0x8A)

# Collection-уудын онцолсон өнгө
USERS_COL    = RGBColor(0x1E, 0x88, 0x4F)   # ногоон
BOOKINGS_COL = RGBColor(0x6A, 0x1B, 0x9A)   # ягаан
FIXED_COL    = RGBColor(0xC2, 0x6A, 0x12)   # улбар шар
EMAIL_COL    = RGBColor(0x1F, 0x3A, 0x68)   # хөх

DOCID_BG = RGBColor(0xFF, 0xF8, 0xE1)       # цайвар шар (key талбарыг тэмдэглэх)
DOCID_BD = RGBColor(0xB7, 0x8A, 0x12)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0, dash=None):
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)
    if dash is not None:
        ln = line._get_or_add_ln()
        prst = ln.find(qn('a:prstDash'))
        if prst is None:
            prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)


def no_fill(shape):
    shape.fill.background()


def text_in_shape(shape, text, *, size=11, bold=False, color=INK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(60_000); tf.margin_right = Emu(60_000)
    tf.margin_top = Emu(20_000);  tf.margin_bottom = Emu(20_000)
    tf.vertical_anchor = anchor
    tf.clear()
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color; run.font.name = "Calibri"


def add_textbox(slide, x, y, w, h, text, *, size=11, bold=False,
                italic=False, color=INK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    return tb


def add_collection_card(slide, x, y, w, h, name, fields, color):
    """Collection нэр + талбаруудын жагсаалт.
       fields: [(name, type, is_key)] жагсаалт.
    """
    HEADER_H = Emu(450_000)

    # 1) Үндсэн карт background (paper)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    set_fill(card, PAPER)
    set_line(card, color, 1.5)
    card.shadow.inherit = False

    # 2) Толгойн зурвас (өнгөт, дээд хэсэг)
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x + Emu(40_000),
                                  y + Emu(40_000),
                                  w - Emu(80_000),
                                  HEADER_H)
    set_fill(head, color)
    set_line(head, color, 0.25)
    text_in_shape(head, name, size=14, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Толгойн өмнөх жижиг "📁" badge маягийн icon — text-ээр
    add_textbox(slide,
                x + Emu(80_000), y + Emu(40_000),
                Emu(400_000), HEADER_H,
                "▤", size=14, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # 3) Талбаруудын мөр тус бүрийг textbox-аар нэмэх
    field_top = y + Emu(40_000) + HEADER_H + Emu(150_000)
    row_h     = Emu(330_000)
    pad_x     = Emu(180_000)

    for i, field in enumerate(fields):
        fy = field_top + i * row_h
        if isinstance(field, tuple):
            fname, ftype, is_key = field
        else:
            fname, ftype, is_key = field, "", False

        # PK indicator
        if is_key:
            # Жижиг шар badge талбарын зүүн талд
            key_w = Emu(280_000); key_h = Emu(220_000)
            key_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x + pad_x, fy + Emu(40_000),
                key_w, key_h)
            key_box.adjustments[0] = 0.45
            set_fill(key_box, DOCID_BG)
            set_line(key_box, DOCID_BD, 0.5)
            text_in_shape(key_box, "PK", size=8, bold=True, color=DOCID_BD)
            name_x = x + pad_x + key_w + Emu(120_000)
        else:
            # Жижиг bullet ●
            add_textbox(slide,
                        x + pad_x, fy,
                        Emu(280_000), row_h,
                        "•", size=14, bold=True, color=color,
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            name_x = x + pad_x + Emu(280_000)

        # Талбарын нэр
        add_textbox(slide,
                    name_x, fy,
                    Emu(1_400_000), row_h,
                    fname, size=11,
                    bold=is_key,
                    color=INK, align=PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.MIDDLE,
                    font="Consolas")

        # Талбарын төрөл (баруун талд)
        if ftype:
            add_textbox(slide,
                        x + pad_x + Emu(1_700_000), fy,
                        w - pad_x - Emu(1_900_000), row_h,
                        ftype, size=10, italic=True,
                        color=INK_SOFT, align=PP_ALIGN.LEFT,
                        anchor=MSO_ANCHOR.MIDDLE,
                        font="Consolas")

    return card


def add_relation(slide, x1, y1, x2, y2, label):
    """Хоёр collection-ийн хооронд dashed харилцааны шугам."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line = conn.line
    line.color.rgb = RELATION
    line.width = Pt(1.25)
    ln = line._get_or_add_ln()
    prst = ln.find(qn('a:prstDash'))
    if prst is None:
        prst = etree.SubElement(ln, qn('a:prstDash'))
    prst.set('val', 'dash')
    # Сумны үзүүр
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'sm')
    tail.set('h', 'sm')

    # Label жижиг хайрцагт
    cx = (x1 + x2) // 2
    cy = (y1 + y2) // 2
    lw = Emu(900_000); lh = Emu(280_000)
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                cx - lw // 2, cy - lh // 2, lw, lh)
    bg.adjustments[0] = 0.45
    set_fill(bg, WHITE)
    set_line(bg, RELATION, 0.5)
    text_in_shape(bg, label, size=9, bold=True, italic=True,
                  color=RELATION)
    return conn


# --- Презентац ачаалах --------------------------------------------------
prs = Presentation(SRC)
slide = prs.slides[15]   # 16-р slide

# Зураг ба title-ыг устгах (хүрээний rectangle-уудыг үлдээх)
to_delete = []
for sh in slide.shapes:
    if sh.shape_type == 13:                  # Picture
        to_delete.append(sh)
    elif sh.name == "TextBox 5":              # Хуучин title
        to_delete.append(sh)
for sh in to_delete:
    remove_shape(sh)

SW = prs.slide_width
SH = prs.slide_height

# --- Title ------------------------------------------------------------
title_tb = slide.shapes.add_textbox(Emu(420_000), Emu(120_000),
                                    SW - Emu(2_500_000), Emu(450_000))
tf = title_tb.text_frame
tf.margin_left = Emu(0); tf.margin_right = Emu(0)
tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r1 = p.add_run(); r1.text = "Бүлэг 2.4:  "
r1.font.size = Pt(14); r1.font.bold = True
r1.font.color.rgb = USERS_COL; r1.font.name = "Calibri"
r2 = p.add_run(); r2.text = "Firebase Firestore — Өгөгдлийн схем"
r2.font.size = Pt(20); r2.font.bold = True
r2.font.color.rgb = ACCENT; r2.font.name = "Calibri"

# Баруун дээд буланд "NoSQL · ER" badge
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               SW - Emu(2_100_000), Emu(150_000),
                               Emu(1_700_000), Emu(380_000))
badge.adjustments[0] = 0.45
set_fill(badge, ACCENT)
set_line(badge, ACCENT, 0.5)
badge.shadow.inherit = False
text_in_shape(badge, "4 COLLECTIONS", size=10, bold=True, color=WHITE)

# --- Collection card-уудын байрлал (2×2 grid) -------------------------
CARD_W = Emu(2_900_000)
CARD_H = Emu(2_400_000)
GRID_GAP_X = Emu(2_400_000)
GRID_GAP_Y = Emu(280_000)

# Сүлжээний эхлэл — төвийг харгалзан
total_w = 2 * CARD_W + GRID_GAP_X
LEFT_X = (SW - total_w) // 2
TOP_Y  = Emu(900_000)

# 4 cell-ийн x,y координат
cells = {
    "users":     (LEFT_X,                       TOP_Y),
    "bookings":  (LEFT_X + CARD_W + GRID_GAP_X, TOP_Y),
    "fixed":     (LEFT_X,                       TOP_Y + CARD_H + GRID_GAP_Y),
    "email":     (LEFT_X + CARD_W + GRID_GAP_X, TOP_Y + CARD_H + GRID_GAP_Y),
}

# --- 1) users -----------------------------------------------------------
ux, uy = cells["users"]
add_collection_card(slide, ux, uy, CARD_W, CARD_H, "users", [
    ("uid",            "doc ID",     True),
    ("name",           "string",     False),
    ("email",          "string",     False),
    ("emailVerified",  "bool",       False),
    ("createdAt",      "timestamp",  False),
], USERS_COL)

# --- 2) bookings --------------------------------------------------------
bx, by = cells["bookings"]
add_collection_card(slide, bx, by, CARD_W, CARD_H, "bookings", [
    ("code",     "DBK-xxxxx",      True),
    ("userId",   "string  (FK)",   False),
    ("venueId",  "string",         False),
    ("date",     "string  (Y-M-D)",False),
    ("timeSlot", "string",         False),
    ("status",   "upcoming | cancelled", False),
], BOOKINGS_COL)

# --- 3) fixed_bookings --------------------------------------------------
fx, fy = cells["fixed"]
add_collection_card(slide, fx, fy, CARD_W, CARD_H, "fixed_bookings", [
    ("venueId",          "string",  False),
    ("organizationName", "string",  False),
    ("dayOfWeek",        "0..6",     False),
    ("timeSlot",         "string",  False),
    ("isActive",         "bool",    False),
], FIXED_COL)

# --- 4) email_verifications --------------------------------------------
ex, ey = cells["email"]
add_collection_card(slide, ex, ey, CARD_W, CARD_H, "email_verifications", [
    ("uid",       "doc ID",     True),
    ("code",      "6 оронтой",  False),
    ("email",     "string",     False),
    ("expiresAt", "timestamp",  False),
], EMAIL_COL)

# --- Харилцаа: users.uid → bookings.userId (хэвтээ) -------------------
# users-ийн right edge → bookings-ийн left edge (дунд)
y_mid = uy + CARD_H // 2
add_relation(slide,
             ux + CARD_W, y_mid,
             bx,          y_mid,
             "userId")

# --- Харилцаа: users.uid → email_verifications.uid -------------------
# users.bottom-right corner → email.top-left corner (диагональ)
add_relation(slide,
             ux + CARD_W - Emu(400_000), uy + CARD_H,
             ex + Emu(400_000),          ey,
             "uid")

# --- Доод тайлбар: PK тэмдэглэгээ ------------------------------------
LEG_X = Emu(450_000)
LEG_Y = Emu(6_500_000)
LEG_W = SW - 2 * LEG_X
LEG_H = Emu(280_000)

add_textbox(slide,
            LEG_X, LEG_Y, LEG_W, LEG_H,
            "▤  Collection нэр       PK  doc ID түлхүүр       ─ ─ ─  Гадаад түлхүүр харилцаа",
            size=10, italic=True, color=INK_SOFT,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# --- Хадгалах ---------------------------------------------------------
prs.save(DST)
print("Done")
