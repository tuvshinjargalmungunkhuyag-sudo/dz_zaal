"""
v3_editable.pptx -> v3_editable.pptx (in-place дээр slide 6 шинэчлэх)
6-р slide дээр "3 давхаргат архитектурын онол" диаграммыг native PowerPoint
shape-аар зурна.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v2.pptx"

# --- Өнгөний палет --------------------------------------------------------
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
INK          = RGBColor(0x1B, 0x1F, 0x3A)
SUBTLE       = RGBColor(0x6B, 0x72, 0x8A)
ACCENT_DARK  = RGBColor(0x1F, 0x3A, 0x68)

# 3 давхрын өнгө
PRES_FILL    = RGBColor(0xE3, 0xEA, 0xF6)   # цайвар хөх
PRES_DARK    = RGBColor(0x1F, 0x3A, 0x68)   # хар хөх
APP_FILL     = RGBColor(0xFF, 0xF1, 0xE0)   # цайвар улбар шар
APP_DARK     = RGBColor(0xC2, 0x6A, 0x12)   # хар улбар шар
DATA_FILL    = RGBColor(0xE7, 0xF5, 0xEC)   # цайвар ногоон
DATA_DARK    = RGBColor(0x1E, 0x88, 0x4F)   # хар ногоон


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_tier(slide, x, y, w, h, *,
             tier_label, title, subtitle, fill, dark):
    """Нэг давхарга — өнгөт зурвас + гарчиг + тайлбар."""
    # Үндсэн карт
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.08
    set_fill(card, fill)
    set_line(card, dark, 1.5)
    card.shadow.inherit = False

    # Зүүн талын өнгөт зурвас (highlight bar)
    bar_w = Emu(160000)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x, y + Emu(80000), bar_w, h - Emu(160000))
    set_fill(bar, dark)
    set_line(bar, dark, 0.25)

    # Tier нэр (TIER 1/2/3)
    add_text(slide, x + bar_w + Emu(250000), y + Emu(150000),
             Emu(2_400_000), Emu(380000),
             tier_label, size=11, bold=True, color=dark,
             align=PP_ALIGN.LEFT)

    # Гол гарчиг
    add_text(slide, x + bar_w + Emu(250000), y + Emu(550000),
             w - bar_w - Emu(450000), Emu(550000),
             title, size=20, bold=True, color=dark,
             align=PP_ALIGN.LEFT)

    # Дэд тайлбар
    add_text(slide, x + bar_w + Emu(250000), y + Emu(1_120_000),
             w - bar_w - Emu(450000), Emu(450000),
             subtitle, size=12, color=INK,
             align=PP_ALIGN.LEFT)

    return card


def add_arrow_with_label(slide, x_center, y_top, height, label):
    """Босоо чиглэлийн сум + хажууд label."""
    arrow_w = Emu(420000)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        x_center - arrow_w // 2, y_top, arrow_w, height)
    set_fill(arrow, ACCENT_DARK)
    set_line(arrow, ACCENT_DARK, 0.5)

    # Сумны баруун талд label
    add_text(slide, x_center + arrow_w, y_top + (height // 2) - Emu(150000),
             Emu(3_500_000), Emu(350000),
             label, size=11, italic=True, color=SUBTLE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# --- Презентац ачаалах ---------------------------------------------------
prs = Presentation(SRC)
slide = prs.slides[5]   # 6-р slide

# Хуучин content placeholder-ыг устгаад title-ыг л үлдээх
to_delete = []
for sh in slide.shapes:
    if sh.name == "Content Placeholder 2" or sh.shape_type == 13:
        to_delete.append(sh)
for sh in to_delete:
    remove_shape(sh)

# Title-ыг сайжруулах
for sh in slide.shapes:
    if sh.name == "Title 1" and sh.has_text_frame:
        sh.text_frame.clear()
        p = sh.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "3 давхаргат архитектурын онол"
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = ACCENT_DARK
        break

SW = prs.slide_width      # 12 192 000
SH = prs.slide_height     # 6 858 000

# --- Layout тогтмол ------------------------------------------------------
TIER_W = Emu(8_400_000)
TIER_H = Emu(1_300_000)
TIER_X = (SW - TIER_W) // 2
GAP = Emu(450_000)            # давхрын хоорондох зай (сум багтах)

ARROW_H = GAP - Emu(60_000)

# Эхний tier (Presentation) дээд талд
TOP_Y = Emu(1_700_000)

PRES_Y = TOP_Y
APP_Y  = PRES_Y + TIER_H + GAP
DATA_Y = APP_Y + TIER_H + GAP

# --- 3 давхарга ----------------------------------------------------------
add_tier(slide, TIER_X, PRES_Y, TIER_W, TIER_H,
         tier_label="TIER 1 — PRESENTATION",
         title="Flutter Мобайл клиент",
         subtitle="Android · iOS  ·  UI/UX дэлгэц, навигаци, локал мэдэгдэл",
         fill=PRES_FILL, dark=PRES_DARK)

add_tier(slide, TIER_X, APP_Y, TIER_W, TIER_H,
         tier_label="TIER 2 — APPLICATION",
         title="Node.js / Express REST API",
         subtitle="AWS EC2 · Docker контейнер  ·  Бизнесийн логик, email баталгаажуулалт",
         fill=APP_FILL, dark=APP_DARK)

add_tier(slide, TIER_X, DATA_Y, TIER_W, TIER_H,
         tier_label="TIER 3 — DATA",
         title="Firebase",
         subtitle="Cloud Firestore · Authentication · Cloud Functions",
         fill=DATA_FILL, dark=DATA_DARK)

# --- Сум 1: Presentation -> Application ---------------------------------
arrow_x = TIER_X + TIER_W // 2 - Emu(2_400_000)   # төв тал руу шилжүүлсэн
add_arrow_with_label(slide,
                     x_center=arrow_x,
                     y_top=PRES_Y + TIER_H + Emu(30000),
                     height=ARROW_H,
                     label="REST API + Firebase SDK")

# --- Сум 2: Application -> Data -----------------------------------------
add_arrow_with_label(slide,
                     x_center=arrow_x,
                     y_top=APP_Y + TIER_H + Emu(30000),
                     height=ARROW_H,
                     label="Firestore Admin SDK")

# --- Зүүн талын чиглэлийн label-ууд (клиент / логик / өгөгдөл) ---------
LABEL_X = TIER_X - Emu(1_300_000)
LABEL_W = Emu(1_200_000)

for y, txt, color in [
        (PRES_Y, "клиент",  PRES_DARK),
        (APP_Y,  "логик",   APP_DARK),
        (DATA_Y, "өгөгдөл", DATA_DARK)]:
    add_text(slide, LABEL_X, y, LABEL_W, TIER_H,
             txt, size=14, bold=True, italic=True, color=color,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# --- Хадгалах ------------------------------------------------------------
prs.save(DST)
print("Saved:", DST)
