"""
v3_editable_v4.pptx -> v3_editable_v5.pptx
14-р slide "AI Туслах — Хоёр горим" диаграммыг native PowerPoint
shape-аар зурна. Бүх элемент editable.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

SRC = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v4.pptx"
DST = r"C:\Users\user\Documents\dz_zaal\diplom\v3_editable_v5.pptx"

# --- Өнгө ---------------------------------------------------------------
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x1B, 0x1F, 0x3A)
INK_SOFT   = RGBColor(0x4A, 0x52, 0x6E)
PAPER      = RGBColor(0xFA, 0xFA, 0xFA)
PAPER2     = RGBColor(0xEE, 0xF1, 0xF8)

ACCENT     = RGBColor(0x1F, 0x3A, 0x68)   # хөх — header
APP_FILL   = RGBColor(0xE3, 0xEA, 0xF6)
APP_DARK   = RGBColor(0x1F, 0x3A, 0x68)

PRIMARY_FILL = RGBColor(0xE7, 0xF5, 0xEC)   # цайвар ногоон
PRIMARY_DARK = RGBColor(0x1E, 0x88, 0x4F)   # ногоон
FALLBACK_FILL= RGBColor(0xFF, 0xF1, 0xE0)   # цайвар улбар шар
FALLBACK_DARK= RGBColor(0xC2, 0x6A, 0x12)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb, width_pt=1.0, dash=None):
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)
    if dash is not None:
        ln = line._get_or_add_ln()
        prst = ln.find(qn('a:prstDash'))
        if prst is None:
            prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)


def no_fill(shape):
    shape.fill.background()


def text_in_shape(shape, text, *, size=12, bold=False, color=INK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(60_000)
    tf.margin_right = Emu(60_000)
    tf.margin_top = Emu(40_000)
    tf.margin_bottom = Emu(40_000)
    tf.vertical_anchor = anchor
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            run = p.add_run()
        else:
            p2 = tf.add_paragraph()
            p2.alignment = align
            run = p2.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"


def add_textbox(slide, x, y, w, h, text, *, size=12, bold=False,
                italic=False, color=INK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_card(slide, cx, cy, w, h, title, subtitle, *,
             fill, border, dark, mode_label=None):
    """Гарчиг + дэд тайлбар + (заавал биш) дээр mode badge-тай rounded карт."""
    x, y = cx - w // 2, cy - h // 2
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.10
    set_fill(box, fill)
    set_line(box, dark, 1.5)
    box.shadow.inherit = False

    # Зүүн талын highlight bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x, y + Emu(60_000),
                                 Emu(120_000), h - Emu(120_000))
    set_fill(bar, dark)
    set_line(bar, dark, 0.25)

    # Mode badge (дээд буланд)
    if mode_label:
        bw = Emu(900_000); bh = Emu(280_000)
        bx = x + w - bw - Emu(140_000)
        by = y + Emu(120_000)
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       bx, by, bw, bh)
        badge.adjustments[0] = 0.45
        set_fill(badge, dark)
        set_line(badge, dark, 0.25)
        text_in_shape(badge, mode_label, size=9, bold=True, color=WHITE)

    # Гол гарчиг
    title_y = y + Emu(150_000) + (Emu(360_000) if mode_label else Emu(120_000))
    add_textbox(slide,
                x + Emu(220_000), title_y,
                w - Emu(280_000), Emu(420_000),
                title, size=15, bold=True, color=dark,
                align=PP_ALIGN.LEFT)

    # Дэд тайлбар
    add_textbox(slide,
                x + Emu(220_000), title_y + Emu(420_000),
                w - Emu(280_000), Emu(380_000),
                subtitle, size=11, color=INK_SOFT,
                align=PP_ALIGN.LEFT)
    return box


def add_arrow(slide, x1, y1, x2, y2, *,
              color=INK, width_pt=1.25, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line = conn.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    ln = line._get_or_add_ln()
    if dash is not None:
        prst = ln.find(qn('a:prstDash'))
        if prst is None:
            prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return conn


# --- Презентац ачаалах --------------------------------------------------
prs = Presentation(SRC)
slide = prs.slides[13]   # 14-р slide

# Хуучин зураг ба placeholder-уудыг устгана (чимэглэлийн frame үлдээх)
to_delete = []
for sh in slide.shapes:
    if sh.shape_type == 13:                     # Picture
        to_delete.append(sh)
    elif sh.name == "TextBox 5":                 # Хуучин title
        to_delete.append(sh)
for sh in to_delete:
    remove_shape(sh)

SW = prs.slide_width
SH = prs.slide_height

# --- Title ---------------------------------------------------------------
title_x = Emu(420_000)
title_y = Emu(120_000)
title_w = SW - Emu(2_500_000)
title_tb = slide.shapes.add_textbox(title_x, title_y, title_w, Emu(450_000))
tf = title_tb.text_frame
tf.margin_left = Emu(0); tf.margin_right = Emu(0)
tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r1 = p.add_run(); r1.text = "Бүлэг 2.2:  "
r1.font.size = Pt(14); r1.font.bold = True
r1.font.color.rgb = FALLBACK_DARK; r1.font.name = "Calibri"
r2 = p.add_run(); r2.text = "AI Туслах — Хоёр горим"
r2.font.size = Pt(20); r2.font.bold = True
r2.font.color.rgb = ACCENT; r2.font.name = "Calibri"

# Баруун дээд буланд "Dual Mode" badge
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               SW - Emu(2_100_000), Emu(150_000),
                               Emu(1_700_000), Emu(380_000))
badge.adjustments[0] = 0.45
set_fill(badge, ACCENT)
set_line(badge, ACCENT, 0.5)
badge.shadow.inherit = False
text_in_shape(badge, "DUAL MODE", size=10, bold=True, color=WHITE)

# --- Layout: 3 багана ---------------------------------------------------
# Col 1 (Flutter App)        Col 2 (CF / EC2 stack)        Col 3 (LLM stack)
COL1_X = Emu(1_500_000)
COL2_X = Emu(5_400_000)
COL3_X = Emu(9_700_000)

# Дээд эгнээ (primary), доод эгнээ (fallback)
ROW_TOP_Y    = Emu(2_200_000)
ROW_BOTTOM_Y = Emu(4_500_000)

CARD_W = Emu(2_900_000)
CARD_H = Emu(1_350_000)

APP_W  = Emu(2_500_000)
APP_H  = Emu(1_400_000)
# Flutter App-ыг хоёр эгнээний дунд байрлуулах
APP_Y = (ROW_TOP_Y + ROW_BOTTOM_Y) // 2

# --- 1) Flutter App -----------------------------------------------------
add_card(slide, COL1_X, APP_Y, APP_W, APP_H,
         "Flutter App", "Хэрэглэгчийн чат UI",
         fill=APP_FILL, border=APP_DARK, dark=APP_DARK)

# --- 2) Cloud Functions (primary middle) -------------------------------
add_card(slide, COL2_X, ROW_TOP_Y, CARD_W, CARD_H,
         "Firebase\nCloud Functions",
         "API key Flutter-т харагдахгүй",
         fill=PRIMARY_FILL, border=PRIMARY_DARK, dark=PRIMARY_DARK,
         mode_label="ҮНДСЭН")

# --- 3) Node.js / EC2 (fallback middle) --------------------------------
add_card(slide, COL2_X, ROW_BOTTOM_Y, CARD_W, CARD_H,
         "Node.js / EC2",
         "Cloud Function алдаа → backup",
         fill=FALLBACK_FILL, border=FALLBACK_DARK, dark=FALLBACK_DARK,
         mode_label="НӨӨЦ")

# --- 4) Claude Haiku 4.5 (primary right) -------------------------------
add_card(slide, COL3_X, ROW_TOP_Y, CARD_W, CARD_H,
         "Claude Haiku 4.5",
         "Anthropic · < 1 сек хариу",
         fill=PRIMARY_FILL, border=PRIMARY_DARK, dark=PRIMARY_DARK)

# --- 5) Groq Llama 3.1-8b (fallback right) -----------------------------
add_card(slide, COL3_X, ROW_BOTTOM_Y, CARD_W, CARD_H,
         "Groq Llama 3.1-8b",
         "LPU дэд бүтэц · хурдан",
         fill=FALLBACK_FILL, border=FALLBACK_DARK, dark=FALLBACK_DARK)

# --- Сумнууд ------------------------------------------------------------
# (1) App.right -> CF.left  (primary, ногоон)
add_arrow(slide,
          COL1_X + APP_W // 2, APP_Y - Emu(250_000),
          COL2_X - CARD_W // 2, ROW_TOP_Y,
          color=PRIMARY_DARK, width_pt=1.75)
# label "Cloud Function"
mid_x = (COL1_X + APP_W // 2 + COL2_X - CARD_W // 2) // 2
mid_y_top = (APP_Y - Emu(250_000) + ROW_TOP_Y) // 2
add_textbox(slide,
            mid_x - Emu(900_000), mid_y_top - Emu(380_000),
            Emu(1_800_000), Emu(280_000),
            "Cloud Function call",
            size=10, italic=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)

# (2) App.right -> EC2.left (fallback, dashed улбар)
add_arrow(slide,
          COL1_X + APP_W // 2, APP_Y + Emu(250_000),
          COL2_X - CARD_W // 2, ROW_BOTTOM_Y,
          color=FALLBACK_DARK, width_pt=1.75, dash="dash")
mid_y_bot = (APP_Y + Emu(250_000) + ROW_BOTTOM_Y) // 2
add_textbox(slide,
            mid_x - Emu(900_000), mid_y_bot + Emu(120_000),
            Emu(1_800_000), Emu(280_000),
            "fallback if timeout",
            size=10, italic=True, color=FALLBACK_DARK, align=PP_ALIGN.CENTER)

# (3) CF.right -> Claude.left
add_arrow(slide,
          COL2_X + CARD_W // 2, ROW_TOP_Y,
          COL3_X - CARD_W // 2, ROW_TOP_Y,
          color=PRIMARY_DARK, width_pt=1.75)

# (4) EC2.right -> Groq.left
add_arrow(slide,
          COL2_X + CARD_W // 2, ROW_BOTTOM_Y,
          COL3_X - CARD_W // 2, ROW_BOTTOM_Y,
          color=FALLBACK_DARK, width_pt=1.75)

# --- Доод тайлбарын зурвас ---------------------------------------------
DESC_Y = Emu(5_700_000)
DESC_H = Emu(900_000)
DESC_X = Emu(450_000)
DESC_W = SW - 2 * DESC_X

desc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  DESC_X, DESC_Y, DESC_W, DESC_H)
desc_box.adjustments[0] = 0.12
set_fill(desc_box, PAPER2)
set_line(desc_box, ACCENT, 1.0)
desc_box.shadow.inherit = False

# 2 баганын тайлбар desc_box дотор
half_w = (DESC_W - Emu(200_000)) // 2

# Зүүн (Үндсэн)
add_textbox(slide,
            DESC_X + Emu(200_000), DESC_Y + Emu(120_000),
            half_w - Emu(100_000), Emu(280_000),
            "● ҮНДСЭН ГОРИМ",
            size=11, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.LEFT)
add_textbox(slide,
            DESC_X + Emu(200_000), DESC_Y + Emu(420_000),
            half_w - Emu(100_000), Emu(420_000),
            "Claude Haiku 4.5 нь Cloud Function-аар дуудагдаж < 1 секундэд хариулдаг",
            size=11, color=INK, align=PP_ALIGN.LEFT)

# Баруун (Нөөц)
add_textbox(slide,
            DESC_X + half_w + Emu(300_000), DESC_Y + Emu(120_000),
            half_w - Emu(100_000), Emu(280_000),
            "● НӨӨЦ ГОРИМ",
            size=11, bold=True, color=FALLBACK_DARK, align=PP_ALIGN.LEFT)
add_textbox(slide,
            DESC_X + half_w + Emu(300_000), DESC_Y + Emu(420_000),
            half_w - Emu(100_000), Emu(420_000),
            "Cloud Function алдаа гарвал Groq Llama автоматаар идэвхждэг",
            size=11, color=INK, align=PP_ALIGN.LEFT)

# --- Хадгалах -----------------------------------------------------------
prs.save(DST)
print("Saved:", DST)
