"""make_v9.py  —  Professional rebuild of Goviyn_Sport_Iltgel.pptx
No emojis; academic Mongolian language; clean geometric diagrams.
"""
import sys, io, shutil, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SH

# ── Palette ────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0A, 0x1A, 0x5C)
BLUE    = RGBColor(0x1A, 0x5C, 0xBB)
BLUE_L  = RGBColor(0xD6, 0xEA, 0xFF)
AMBER   = RGBColor(0xA8, 0x6A, 0x00)
AMBER_L = RGBColor(0xFF, 0xF0, 0xCC)
GREEN   = RGBColor(0x15, 0x6B, 0x36)
GREEN_L = RGBColor(0xD5, 0xF5, 0xE3)
RED     = RGBColor(0xA0, 0x28, 0x18)
RED_L   = RGBColor(0xFD, 0xED, 0xEB)
PURPLE  = RGBColor(0x5A, 0x28, 0x7A)
PURP_L  = RGBColor(0xF0, 0xE6, 0xFF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF2, 0xF4, 0xF7)
MGRAY   = RGBColor(0xCC, 0xCC, 0xD8)
DGRAY   = RGBColor(0x33, 0x33, 0x44)
YELLOW  = RGBColor(0xF5, 0xC5, 0x18)

def I(x): return Inches(x)
def P(x): return Pt(x)

# ── Low-level helpers ───────────────────────────────────────────────
def box(slide, l, t, w, h, fill=None, border=None, bw=P(1.2), rounded=False):
    kind = SH.ROUNDED_RECTANGLE if rounded else SH.RECTANGLE
    sh = slide.shapes.add_shape(kind, I(l), I(t), I(w), I(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border; sh.line.width = bw
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size=11, bold=False, color=NAVY,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = P(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb

def ctxt(slide, text, l, t, w, h, **kw):
    return txt(slide, text, l, t, w, h, align=PP_ALIGN.CENTER, **kw)

def label_pill(slide, text, l, t, w=0.60, h=0.26, bg=BLUE, fg=WHITE, size=8):
    """Small colored pill label — replaces emoji icon."""
    box(slide, l, t, w, h, fill=bg, border=None, rounded=True)
    ctxt(slide, text, l, t, w, h, size=size, bold=True, color=fg)

def num_circle(slide, n, l, t, d=0.38, bg=NAVY, fg=WHITE, size=13):
    """Filled circle with number."""
    box(slide, l, t, d, d, fill=bg, border=None, rounded=True)
    ctxt(slide, str(n), l, t, d, d, size=size, bold=True, color=fg)

def horiz_arrow(slide, x1, y, x2, color=BLUE, label=''):
    """Horizontal arrow from x1 to x2 at vertical position y."""
    shaft_h = 0.04
    shaft_y = y + 0.10
    box(slide, x1, shaft_y, x2-x1-0.10, shaft_h, fill=color, border=None)
    # triangle head
    box(slide, x2-0.12, y, 0.12, 0.24, fill=color, border=None)
    if label:
        ctxt(slide, label, x1, y-0.20, x2-x1, 0.20, size=7.5,
             italic=True, color=color)

def vert_arrow(slide, x, y1, y2, color=BLUE, label=''):
    """Vertical arrow with optional label."""
    shaft_w = 0.04
    mx = x - shaft_w/2
    box(slide, mx, y1, shaft_w, y2-y1-0.10, fill=color, border=None)
    box(slide, x-0.10, y2-0.12, 0.20, 0.12, fill=color, border=None)
    if label:
        ctxt(slide, label, x-0.75, (y1+y2)/2-0.12, 1.50, 0.24,
             size=7.5, italic=True, color=color)

def clear_content(slide, keep=7):
    for sh in list(slide.shapes)[keep:]:
        sh._element.getparent().remove(sh._element)

# ── Emoji stripper (for slides that aren't rebuilt) ────────────────
EMOJI_RE = re.compile(
    "[\U00010000-\U0010ffff]"
    "|[\U0001F300-\U0001FFFF]"
    "|[\u2600-\u27BF]"
    "|[\u2B00-\u2BFF]"
    "|[\u2300-\u23FF]"
    "|[\u25A0-\u25FF]"
    "|[\u2700-\u27BF]",
    flags=re.UNICODE
)

def strip_emojis_slide(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                cleaned = EMOJI_RE.sub('', run.text)
                # replace ✓ checkmark with dash
                cleaned = cleaned.replace('\u2713', '-').replace('\u2605', '')
                run.text = cleaned

# ── Load v7 → copy as v9 ───────────────────────────────────────────
shutil.copy('Goviyn_Sport_Iltgel_v7.pptx', 'Goviyn_Sport_Iltgel_v9.pptx')
prs = Presentation('Goviyn_Sport_Iltgel_v9.pptx')
slides = prs.slides

CY = 0.85   # content top
CH = 6.20   # content height
CX = 0.28   # left margin
CW = 12.77  # content width

print("=== Building Goviyn_Sport_Iltgel_v9.pptx ===\n")

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Research Problem: gap-analysis diagram
# ════════════════════════════════════════════════════════════════════
s = slides[2]; clear_content(s)

# Section header strip
box(s, CX, CY, CW, 0.44, fill=NAVY, border=None)
ctxt(s, "Даланзадгадын спортын байгууллагуудын одоогийн нөхцөл байдал ба судалгааны актуаль байдал",
     CX+0.10, CY+0.04, CW-0.20, 0.36, size=10, bold=True, color=WHITE)

# ── 4 problem cards on the left (2-col × 2-row) ──
problems = [
    ("П-1", "Дижитал дэд бүтцийн дутагдал",
     "Спортын заалны захиалга нь утасны яриа болон биечилсэн бүртгэлд тулгуурладаг тул хэрэглэгчийн цаг зарцуулалт болон алдааны эрсдэл өндөр байна."),
    ("П-2", "Бодит цагийн мэдээллийн хүртээмжгүй байдал",
     "Захиалгын боломжит слотын нөхцөл байдлыг урьдчилан мэдэх цахим технологийн боломж одоогоор байхгүй байна."),
    ("П-3", "Захиалгын давхардал ба мэдээллийн зөрчил",
     "Гараар хөтлөгдөх бүртгэлийн системд давхардсан захиалга болон мэдээллийн зөрчил тогтмол ажиглагдаж байна."),
    ("П-4", "Нэгдсэн аналитик дэд бүтцийн дутагдал",
     "Захиалгын статистик, орлогын тайлан, хэрэглэгчийн зан байдлын дата цуглуулах болон дүн шинжилгээний автоматчлал хийгдээгүй."),
]

pw = 5.80; ph = 1.38
for i, (num, title, desc) in enumerate(problems):
    col = i % 2; row = i // 2
    px = CX + col*(pw+0.18)
    py = CY + 0.52 + row*(ph+0.12)
    box(s, px, py, pw, ph, fill=RED_L, border=RED, bw=P(1.0))
    # number pill
    box(s, px, py, 0.48, ph, fill=RED, border=None)
    ctxt(s, num, px, py+ph/2-0.18, 0.48, 0.36, size=8.5, bold=True, color=WHITE)
    txt(s, title, px+0.56, py+0.08, pw-0.64, 0.30, size=10, bold=True, color=RED)
    txt(s, desc,  px+0.56, py+0.42, pw-0.64, 0.88, size=8.5, color=DGRAY)

# ── Bottom: research context ──
ctx_y = CY + 0.52 + 2*(ph+0.12) + 0.06
box(s, CX, ctx_y, CW, CH - (ctx_y-CY), fill=LGRAY, border=BLUE, bw=P(1.0))

# Stat box
box(s, CX+0.12, ctx_y+0.10, 2.00, 0.90, fill=BLUE, border=None, rounded=True)
ctxt(s, "5+",          CX+0.12, ctx_y+0.12, 2.00, 0.40, size=22, bold=True, color=WHITE)
ctxt(s, "спортын заал", CX+0.12, ctx_y+0.56, 2.00, 0.30, size=8.5, color=YELLOW)

# Gap text
txt(s, "Нэгдсэн цахим захиалгын систем огт байхгүй  —  "
       "Бүх байгууллага тусдаа, гараар ажилладаг",
    CX+2.24, ctx_y+0.22, 5.00, 0.56, size=9.5, color=DGRAY)

horiz_arrow(s, CX+7.30, ctx_y+0.34, CX+8.80, color=BLUE, label='')

# Solution
box(s, CX+8.84, ctx_y+0.10, 3.80, 0.90, fill=BLUE, border=None, rounded=True)
ctxt(s, "Шийдэл: Үүлэн орчинд суурилсан\n3 давхаргат захиалгын систем",
     CX+8.84, ctx_y+0.12, 3.80, 0.78, size=9, bold=True, color=WHITE)

print("  Slide 3 done")

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Research Goals: structured goal grid
# ════════════════════════════════════════════════════════════════════
s = slides[3]; clear_content(s)

# Main goal banner
box(s, CX, CY, CW, 0.70, fill=NAVY, border=None, rounded=True)
txt(s, "Зорилго:  Өмнөговь аймгийн Даланзадгад хотын спортын байгууллагуудад зориулан үүлэн орчинд суурилсан, "
       "3 давхаргат архитектуртай, бодит цагийн давхцал хяналт бүхий мобайл захиалгын системийг загварчлан хэрэгжүүлэх.",
    CX+0.16, CY+0.08, CW-0.32, 0.55, size=9.5, bold=False, color=YELLOW)

goals = [
    ("З-1", "Цаг захиалгын модуль",    BLUE, BLUE_L,
     ["Бодит цагийн хүртээмж шалгалт",
      "Давхардлаас сэргийлэх логик",
      "Хагас / бүтэн талбайн уян захиалга"]),
    ("З-2", "Мобайл интерфэйс",        BLUE, BLUE_L,
     ["Flutter — iOS / Android",
      "Material Design 3 стандарт",
      "Дулаан цөлийн өнгөний схем"]),
    ("З-3", "AI туслах систем",         GREEN, GREEN_L,
     ["Anthropic Claude Haiku 4.5 (үндсэн)",
      "Groq LLaMA 3.1-8B (нөөц горим)",
      "Монгол хэлний асуулт-хариулт"]),
    ("З-4", "Үүлэн дэд бүтэц",          AMBER, AMBER_L,
     ["AWS EC2  ap-northeast-1 (Токио)",
      "Docker container ажиллуулалт",
      "GitHub Actions CI/CD автоматжуулалт"]),
    ("З-5", "Баталгаажуулалт & аюулгүй байдал", PURPLE, PURP_L,
     ["Email + 6 оронтой OTP код",
      "Firebase Auth — ID token",
      "Firestore аюулгүй байдлын дүрмүүд"]),
    ("З-6", "Туршилт & баталгаажуулалт",         RED, RED_L,
     ["10 функциональ тест кейс",
      "Ашиглалтын хялбар байдлын туршилт",
      "Гүйцэтгэлийн хэмжилт & шинжилгээ"]),
]

cols = 3; gw = (CW - 0.20) / cols - 0.12
gh = (CH - 0.85) / 2 - 0.10

for i, (num, title, bd, bg, buls) in enumerate(goals):
    col = i % cols; row = i // cols
    gx = CX + col*(gw+0.12)
    gy = CY + 0.80 + row*(gh+0.10)
    box(s, gx, gy, gw, gh, fill=bg, border=bd, bw=P(1.2))
    # header bar
    box(s, gx, gy, gw, 0.46, fill=bd, border=None)
    ctxt(s, num,   gx+0.04, gy+0.06, 0.50, 0.34, size=11, bold=True, color=WHITE)
    txt(s,  title, gx+0.58, gy+0.08, gw-0.66, 0.32, size=10, bold=True, color=WHITE)
    for bi, b in enumerate(buls):
        txt(s, "  " + b, gx+0.10, gy+0.54+bi*0.32, gw-0.18, 0.30,
            size=9, color=DGRAY)

print("  Slide 4 done")

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Research Significance: 3-column academic cards
# ════════════════════════════════════════════════════════════════════
s = slides[4]; clear_content(s)

sigs = [
    ("Шинжлэх ухааны\nач холбогдол", BLUE, BLUE_L,
     "Үүлэн тооцоолол дахь 3 давхаргат архитектурийн загварыг баримтжуулна.",
     ["NoSQL + REST + Flutter хослолын практик кейс судалгаа",
      "AI нэгтгэлийн (fallback chain) загварчлалын жишээ",
      "Жижиг хот дахь дижитал шилжилтийн эмпирик баримт",
      "Ашиглалтын хялбар байдлын туршилтын арга зүйн хэрэглээ"]),
    ("Практик &\nнийгмийн ач холбогдол", GREEN, GREEN_L,
     "Даланзадгадын 5+ спортын байгууллагын үйл ажиллагааны үр ашгийг дээшлүүлнэ.",
     ["Хэрэглэгчийн цаг хугацааны зарцуулалтыг эрс бууруулна",
      "Захиалгын давхардлыг бүрэн арилгана",
      "Орлогын статистик ба аналитикийн автоматжуулалт",
      "Монголын жижиг хотод тохирсон replicate загвар болно"]),
    ("Технологийн\nхувь нэмэр", AMBER, AMBER_L,
     "Монголын хэмжээнд ховор хэрэгждэг орчин үеийн технологийн стекийн туршлага.",
     ["CI/CD pipeline (GitHub Actions → Docker → ECR → EC2)",
      "Email OTP баталгаажуулалтын хэрэгжилтийн загвар",
      "Firestore compound index ба security rules жишээ",
      "Flutter Provider state management хэрэгжилт"]),
]

cw = (CW - 0.24) / 3 - 0.08
for i, (title, bd, bg, intro, buls) in enumerate(sigs):
    cx = CX + i*(cw+0.12)
    box(s, cx, CY, cw, CH-0.05, fill=bg, border=bd, bw=P(1.8))
    # colored header
    box(s, cx, CY, cw, 0.90, fill=bd, border=None)
    ctxt(s, title, cx+0.06, CY+0.06, cw-0.12, 0.80,
         size=11.5, bold=True, color=WHITE)
    # intro sentence
    txt(s, intro, cx+0.14, CY+1.00, cw-0.28, 0.62, size=9, italic=True, color=bd)
    # separator line
    box(s, cx+0.12, CY+1.68, cw-0.24, 0.025, fill=bd, border=None)
    for bi, b in enumerate(buls):
        txt(s, b, cx+0.24, CY+1.74+bi*0.68, cw-0.36, 0.58, size=8.5, color=DGRAY)
        box(s, cx+0.12, CY+1.74+bi*0.68+0.16, 0.08, 0.08, fill=bd, border=None, rounded=True)

print("  Slide 5 done")

# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — System Architecture: full 3-tier + deployment + services
# ════════════════════════════════════════════════════════════════════
s = slides[5]; clear_content(s)

TW = 7.60   # tier block width
TX = CX
TH = 1.50   # tier height
TG = 0.44   # gap between tiers (for arrow)
T1Y = CY + 0.04
T2Y = T1Y + TH + TG
T3Y = T2Y + TH + TG

RX  = TX + TW + 0.22
RW  = CW - TW - 0.18

# helper: draw one tier block
def tier(slide, n, tier_lbl, tech_lbl, region, items,
         lx, ty, w, h, bd, bg, lbl_bg):
    box(slide, lx, ty, w, h, fill=bg, border=bd, bw=P(1.5))
    # left numbered tab
    box(slide, lx, ty, 0.40, h, fill=lbl_bg, border=None)
    ctxt(slide, str(n), lx, ty+h/2-0.24, 0.40, 0.48, size=16, bold=True, color=WHITE)
    # tier label
    txt(slide, tier_lbl, lx+0.48, ty+0.06, w-0.56, 0.26, size=8, italic=True, color=bd)
    txt(slide, tech_lbl, lx+0.48, ty+0.30, w-0.56, 0.30, size=11, bold=True, color=DGRAY)
    if region:
        txt(slide, region, lx+0.48, ty+0.58, w-0.56, 0.22, size=7.5, italic=True, color=bd)
    # component pills
    iw = (w - 0.56 - 0.06*(len(items)-1)) / len(items)
    for j, (tag, name) in enumerate(items):
        ix = lx + 0.48 + j*(iw+0.06)
        iy = ty + (0.84 if region else 0.72)
        ih = h - iy + ty - 0.10
        box(slide, ix, iy, iw, ih, fill=WHITE, border=bd, bw=P(0.7), rounded=True)
        label_pill(slide, tag, ix+0.06, iy+0.06, iw-0.12, 0.22,
                   bg=lbl_bg, fg=WHITE, size=7)
        ctxt(slide, name, ix+0.04, iy+0.32, iw-0.08, ih-0.36, size=8, color=NAVY)

tier(s, 1, "ДАВХАРГА 1  —  Presentation (Клиент)",
     "Flutter Mobile Application", "Android  /  iOS  /  Dart 3.x",
     [("MOB", "Mobile\nClient"), ("UI", "Material\nDesign 3"),
      ("STATE", "Provider\nState Mgmt"), ("HTTP", "Dio HTTP\nClient")],
     TX, T1Y, TW, TH, BLUE, BLUE_L, BLUE)

vert_arrow(s, TX+TW/2, T1Y+TH+0.02, T2Y-0.04, BLUE, "HTTPS / REST API")

tier(s, 2, "ДАВХАРГА 2  —  Application (Серверийн логик)",
     "Node.js 18  /  Express.js",
     "AWS EC2  ap-northeast-1 (Токио)  —  Docker  node:18-alpine",
     [("API", "REST\nEndpoints"), ("SMTP", "Nodemailer\nEmail OTP"),
      ("CRON", "node-cron\nScheduler"), ("ADM", "Firebase\nAdmin SDK")],
     TX, T2Y, TW, TH, AMBER, AMBER_L, AMBER)

vert_arrow(s, TX+TW/2, T2Y+TH+0.02, T3Y-0.04, GREEN, "Firebase Admin SDK")

tier(s, 3, "ДАВХАРГА 3  —  Data (Өгөгдлийн давхарга)",
     "Google Cloud Firestore  (NoSQL)", "",
     [("COL1", "bookings\ncollection"), ("COL2", "users\ncollection"),
      ("COL3", "fixed_\nbookings"), ("COL4", "email_\nverifications")],
     TX, T3Y, TW, TH, GREEN, GREEN_L, GREEN)

# ── Right panel: Auth + AI + CI/CD ──────────────────────────────────
# Auth card
box(s, RX, T1Y, RW, 1.16, fill=PURP_L, border=PURPLE, bw=P(1.2), rounded=True)
box(s, RX, T1Y, RW, 0.34, fill=PURPLE, border=None, rounded=True)
ctxt(s, "БАТАЛГААЖУУЛАЛТ", RX, T1Y+0.05, RW, 0.26, size=8, bold=True, color=WHITE)
txt(s, "Firebase Authentication", RX+0.12, T1Y+0.40, RW-0.20, 0.26, size=9.5, bold=True, color=PURPLE)
txt(s, "Email + 6 оронтой OTP  \u2192  ID Token  \u2192  Firestore Rules",
    RX+0.12, T1Y+0.68, RW-0.20, 0.40, size=8.5, color=DGRAY)

# AI card
box(s, RX, T1Y+1.22, RW, 1.38, fill=GREEN_L, border=GREEN, bw=P(1.2), rounded=True)
box(s, RX, T1Y+1.22, RW, 0.34, fill=GREEN, border=None, rounded=True)
ctxt(s, "AI ТУСЛАХ СИСТЕМ", RX, T1Y+1.27, RW, 0.26, size=8, bold=True, color=WHITE)
txt(s, "Anthropic Claude Haiku 4.5  (үндсэн)",
    RX+0.12, T1Y+1.62, RW-0.20, 0.24, size=9, bold=True, color=GREEN)
txt(s, "Groq LLaMA 3.1-8B  (нөөц горим)",
    RX+0.12, T1Y+1.88, RW-0.20, 0.24, size=8.5, color=DGRAY)
txt(s, "Fallback chain: timeout 5 s \u2192 switch",
    RX+0.12, T1Y+2.10, RW-0.20, 0.44, size=8, italic=True, color=GREEN)

# CI/CD pipeline card
ci_y = T2Y + 0.06
box(s, RX, ci_y, RW, 2.58, fill=PURP_L, border=PURPLE, bw=P(1.4), rounded=True)
box(s, RX, ci_y, RW, 0.34, fill=PURPLE, border=None, rounded=True)
ctxt(s, "CI/CD PIPELINE", RX, ci_y+0.05, RW, 0.26, size=8, bold=True, color=WHITE)

steps = [("1", "git push\n(main)"), ("2", "GitHub\nActions"), ("3", "Docker\nBuild"),
         ("4", "ECR\nPush"), ("5", "EC2\nDeploy")]
sw = (RW - 0.16) / len(steps) - 0.04
for si, (sn, slbl) in enumerate(steps):
    sx = RX + 0.08 + si*(sw+0.04)
    box(s, sx, ci_y+0.42, sw, 0.90, fill=WHITE, border=PURPLE, bw=P(0.7), rounded=True)
    ctxt(s, sn,   sx, ci_y+0.44, sw, 0.26, size=10, bold=True, color=PURPLE)
    ctxt(s, slbl, sx, ci_y+0.70, sw, 0.58, size=7.5, color=DGRAY)
    if si < len(steps)-1:
        ctxt(s, "\u2192", sx+sw, ci_y+0.76, 0.06, 0.24, size=9, color=PURPLE)

txt(s, "Push \u2192 Build \u2192 Registry \u2192 Deploy  |  ~3-4 мин  |  ap-northeast-1",
    RX+0.10, ci_y+1.42, RW-0.18, 0.26, size=7.5, italic=True, color=PURPLE)
txt(s, "Орчин: node:18-alpine  \u2022  Image: dz-zaal-backend  \u2022  Port: 3000",
    RX+0.10, ci_y+1.66, RW-0.18, 0.24, size=7.5, color=DGRAY)
txt(s, "Secrets: AWS keys, EC2 SSH, ECR registry (GitHub Secrets)",
    RX+0.10, ci_y+1.90, RW-0.18, 0.24, size=7.5, color=DGRAY)

# Performance
pf_y = T3Y + 0.10
box(s, RX, pf_y, RW, 1.06, fill=LGRAY, border=NAVY, bw=P(1.0), rounded=True)
ctxt(s, "ГҮЙЦЭТГЭЛИЙН ХЭМЖИЛТ", RX, pf_y+0.04, RW, 0.22, size=7.5, bold=True, color=NAVY)
perf = [("Апп ачаалах","2.3 с"), ("Firestore","0.6 с"),
        ("EC2 API","0.8 с"), ("Claude Haiku","1.4 с")]
pw2 = (RW - 0.12) / len(perf) - 0.04
for pi, (plbl, pval) in enumerate(perf):
    px2 = RX + 0.06 + pi*(pw2+0.04)
    box(s, px2, pf_y+0.30, pw2, 0.68, fill=WHITE, border=NAVY, bw=P(0.7), rounded=True)
    ctxt(s, pval, px2, pf_y+0.32, pw2, 0.32, size=13, bold=True, color=NAVY)
    ctxt(s, plbl, px2, pf_y+0.64, pw2, 0.26, size=7, color=DGRAY)

print("  Slide 6 done")

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Technology Stack: layered visual + rationale
# ════════════════════════════════════════════════════════════════════
s = slides[6]; clear_content(s)

LW = 8.00
layers = [
    ("ДАВХАРГА 1\nPresentation", BLUE, BLUE_L,
     [("Flutter 3.x", "Dart"), ("Material\nDesign 3","MD3"),
      ("Provider","State"), ("Dio","HTTP"), ("Google\nFonts","Font")]),
    ("ДАВХАРГА 2\nApplication", AMBER, AMBER_L,
     [("Node.js 18","Runtime"), ("Express.js","Framework"),
      ("Nodemailer","SMTP"), ("node-cron","Scheduler"), ("Firebase\nAdmin","SDK")]),
    ("ДАВХАРГА 3\nData & Infra", GREEN, GREEN_L,
     [("Firestore","NoSQL DB"), ("Firebase\nAuth","Auth"),
      ("AWS EC2","ap-ne-1"), ("Docker","Container"), ("GitHub\nActions","CI/CD")]),
]

lh = (CH - 0.08) / 3 - 0.10
for li, (lname, bd, bg, techs) in enumerate(layers):
    lx = CX; ly = CY + li*(lh+0.10)
    box(s, lx, ly, LW, lh, fill=bg, border=bd, bw=P(1.6))
    # label bar on left
    box(s, lx, ly, 1.14, lh, fill=bd, border=None)
    ctxt(s, lname, lx, ly+lh/2-0.30, 1.14, 0.60, size=8.5, bold=True, color=WHITE)
    # tech cards
    tw = (LW - 1.22 - 0.06*(len(techs)-1)) / len(techs)
    for ti, (tname, trole) in enumerate(techs):
        tx2 = lx + 1.22 + ti*(tw+0.06)
        box(s, tx2, ly+0.12, tw, lh-0.24, fill=WHITE, border=bd, bw=P(0.8), rounded=True)
        # role tag at top of card
        label_pill(s, trole, tx2+0.06, ly+0.18, tw-0.12, 0.22, bg=bd, fg=WHITE, size=6.5)
        ctxt(s, tname, tx2+0.04, ly+0.46, tw-0.08, lh-0.72, size=9, bold=True, color=NAVY)

# Right panel: selection rationale
rx2 = CX + LW + 0.20
rw2 = CW - LW - 0.16
box(s, rx2, CY, rw2, CH-0.05, fill=NAVY, border=None, rounded=True)
ctxt(s, "Технологийн\nСонголтын Үндэслэл",
     rx2+0.10, CY+0.12, rw2-0.20, 0.56,
     size=11, bold=True, color=YELLOW)

rationale = [
    ("Flutter", "Нэг кодоос iOS/Android — хөгжүүлэлтийн зардал 50% хэмнэгдэнэ."),
    ("Firebase", "Бодит цагийн синк ба Auth нэгтгэл — бэкэнд нарийн төвөгтэй байдал буурна."),
    ("Node.js", "JavaScript full-stack — ижил хэл хэрэглэснээр цикл хурдан."),
    ("AWS EC2", "Токиогийн бүс — Монголд хамгийн ойр, latency хамгийн бага."),
    ("Docker", "Орчны тогтвортой байдал — dev/staging/prod ялгаа арилна."),
    ("Actions", "Push-triggered deploy — гараар орхигдоно, хүний алдаа бууруулна."),
]
for ri, (rtech, rdesc) in enumerate(rationale):
    ry2 = CY + 0.82 + ri*0.86
    box(s, rx2+0.10, ry2, rw2-0.20, 0.80, fill=RGBColor(0x14,0x2A,0x6E), border=None, rounded=True)
    label_pill(s, rtech, rx2+0.16, ry2+0.10, 0.96, 0.22, bg=YELLOW, fg=NAVY, size=7.5)
    txt(s, rdesc, rx2+1.20, ry2+0.08, rw2-1.36, 0.62, size=8.5, color=LGRAY)

print("  Slide 7 done")

# ════════════════════════════════════════════════════════════════════
# Strip emojis from all other slides (8 → 29)
# ════════════════════════════════════════════════════════════════════
print("\n  Stripping emojis from slides 8-29 …")
for i in range(7, len(slides)):
    strip_emojis_slide(slides[i])
    print(f"    Slide {i+1} cleaned")

# ════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════
prs.save('Goviyn_Sport_Iltgel_v9.pptx')
print(f"\n  Saved: Goviyn_Sport_Iltgel_v9.pptx   ({len(slides)} slides)")
print("  Done.")
